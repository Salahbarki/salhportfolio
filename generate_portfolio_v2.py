# -*- coding: utf-8 -*-
"""
Portfolio Maintenance Salah Barki - V2
Design premium minimaliste + Contenu ultra-technique
"""
import os

OUTPUT_DIR = r"C:\Users\ba2rb\Downloads\salpre"

# ============================================================
# DONNÉES TECHNIQUES DÉTAILLÉES
# ============================================================

PROFILE = {
    "name": "Salah Eddine Barki",
    "title": "Technicien de Maintenance Industrielle",
    "contact": "📍 Amiens, France  |  📱 +33 6 88 69 07 04  |  ✉️ salahbarki.seb@gmail.com",
    "summary": """Technicien multitechnique senior (10+ ans) spécialisé en diagnostic avancé sur systèmes automatisés 
Siemens, maintenance préventive/corrective curative sur presses injection, lignes agroalimentaire PET/canette, 
et bancs de test automobile. Expertise couvrant électrotechnique BT, hydraulique (120–200 bar), 
pneumatique (5–8 bar), mécanique de précision et programmation Ladder (TIA Portal V16 / Step 7). 
Certifié en consignation LOTO, analyse RCA, et planification GMAO (SAP PM)."""
}

FICHES = [
    {
        "num": "01",
        "titre": "Panne Encartonneuse — Ligne PET Coca-Cola",
        "entreprise": "Coca-Cola Europacific Partners",
        "lieu": "Grigny (91)",
        "type": "CURATIF URGENT — AUTOMATISME",
        "equip": "Encartonneuse Krones Variopac Pro (module couche + compression) — Ligne PET 2L — Cadence nominale 60 000 bph",
        "contexte": "Ligne critique approvisionnant la grande distribution IDF. Production 22h/jour, 6j/7. Toute immobilisation >20 min génère risque de rupture de stock.",
        "symptome": """Défaut HMI récurrent n°402 « PRODUIT_COINCÉ_ENTRÉE » toutes les 8–12 min. Accumulation de 15 à 18 colis 
rejetés/heure en sortie fardeleuse. Opérateurs forcent le redémarrage manuel, ce qui génère des à-coups mécaniques 
sur le module compression.""",
        "diagnostic": """
• Mesures électriques : alimentation capteur photoélectrique Sick WT24-2B440 (E0.4) = 23,8 VDC stable. 
  Signal logique oscillant 0-1-0 en 18–22 ms au niveau de l’entrée digitale CPU S7-1500 (DI 16×24VDC).
• Analyse programme TIA Portal V16 : dans FC301 (gestion encartonneuse), Network 12, le bit mémoire 
  M120.4 (VALID_PROD) est conditionné par un front montant sur E0.4 sans temporisation d’antirebond. 
  L’oscillation physique du faisceau (brouillard/buée) crée des fronts parasites interprétés comme validés.
• Vérification mécanique : jeu butée amortisseur ACE MC150 à l’entrée produit = 3,2 mm (tolérance 0,5 mm). 
  Le servo-moteur Kollmorgen AKD (axe pousseur) affiche un offset position de +2,3 mm vs consigne initiale 
  (paramètre Pn507 = +120 incréments au lieu de +85). Dérive mécanique due à l’amortisseur usé.
• Analyse des rejets : colis mal formés = produit poussé en biais (offset mécanique + signal instable 
  = mauvaise synchro pousseur/convoyeur).
        """.strip(),
        "actions": """
1. CONSIGNATION : Isolation énergétique LOTO (NF EN ISO 14118) : coupeur-sectionneur Q3 (électrique 400V) 
   + Q7 (pneumatique 6 bar) + condamnateur cadenas. Signalisation verticale.

2. ÉLECTRIQUE : Remplacement capteur Sick WT24-2B440 par modèle Sick WT27L-2F430 (portée 700 mm, 
   antirebond intégré 50 ms, purge d’air intégrée IP69K, raccord M12 5 pôles blindé). 
   Temps de réponse <1 ms. Câble remplacé par LappKabel UNITRONIC® Li2YCY (TP) 5×0,25 mm² blindé.

3. MÉCANIQUE : Dépose butée amortisseur ACE MC150 usée. Pose amortisseur neuf réglé à force 1 
   (course 25 mm). Recalage butée d’entrée au jeu de 0,5 mm (calibre à lame 0,5 mm passant, 0,6 mm non passant). 
   Serrage boulonnerie M10 classe 8.8 au couple 45 Nm (clé dynamométrique Hazet 5121-3CT).

4. AUTOMATISME (TIA Portal V16) :
   • Modification FC301 Network 12 : ajout temporisation S_PEXT (TP) de 50 ms sur l’entrée E0.4 
     avant validation du bit M120.4.
   • Recalage offset servo Kollmorgen AKD : Pn507 ajusté de +120 à +85 incréments. 
     Validation via WorkBench AKD avec oscilloscope intégré (courbe position consigne vs réelle 
     sous charge, écart <0,1 mm).
   • Sauvegarde version programme v2.7 sur serveur SVN interne + export PDF du bloc modifié.

5. ESSAIS : Run continu 45 min à cadence 55 000 bph. Monitoring via HMI Comfort KTP900 : 
   compteur rejets passé de 17/15 min à 2/15 min. Aucun arrêt forcé.

6. GMAO : Création ordre type « CTRL_ENCART_CAPTEUR » périodicité 2 000 h (environ 3 mois). 
   Mise à jour plan de maintenance constructeur Krones dans base SAP PM.
        """.strip(),
        "resultat": """• MTTR : réduit de 42 min à 8 min moyenne (–81 %).
• Taux rejet colis : divisé par 4 (17 rejets/h → 4 rejets/h).
• Zero arrêt forcé opérateur sur les 3 semaines suivantes.
• Amortisseur et capteur intégrés en stock sécurité classe A."""
    },
    {
        "num": "02",
        "titre": "Mise en Service Presse Injection + Auxiliaires",
        "entreprise": "TE Connectivity",
        "lieu": "Tanger (Maroc)",
        "type": "INSTALLATION & MISE EN ROUTE",
        "equip": "Presse injection électrique KraussMaffei CX 120-380 (120 T) + sécheur Motan Luxor CA A 80 + détecteur métal Sesotec RAPID VARIO",
        "contexte": "Relocalisation complète d’une cellule d’injection depuis ancien hall vers nouveau bâtiment de production. Objectif : mise en production J+2 pour contrat client automobile.",
        "symptome": "Nouvelle installation. Aucun symptôme. Objectif : démarrage qualité premier jet avec cycle stable <20 s.",
        "diagnostic": "Pas de diagnostic défaut. Analyse pré-installation : plan de pose, alignement géométrique, compatibilité réseaux (hydraulique 120 bar / eau 80 °C / électrique 400V triphasé + T + terre).",
        "actions": """
1. DÉPOSE & TRANSPORT : Dépose outillage (moule 24 cavités connecteurs auto 2,54 mm). 
   Vidange partielle circuit thermorégulation. Transport chariot élévateur Jungheinrich EFG 3,5T 
   avec fourches 1 800 mm + sangle textile 4T. Contrôle masse totale 4 200 kg.

2. ANCRAGE : Pose 4 plots antivibratoires Vibrostop Type B (charge 1 500 kg/unité, réduction vibration 85 %). 
   Nivellement laser à ±0,2 mm/m. Alignement axe presse vs convoyeur amont à ±1 mm (cordeau + laser).

3. ÉLECTRIQUE : Raccordement arrivée 400 V AC tri + N + T via câble 5G16 mm² (courant nominal 32 A). 
   Mesure terre : Rterre = 0,4 Ω (seuil NFC 15-100 : <1 Ω). Armoire électrique : Siemens S7-1200 CPU 1214C DC/DC/DC 
   + module SM1234 AI/AO. Raccordement bus PROFINET vers HMI KTP700 + variateur servo Sinamics V90.

4. HYDRAULIQUE (centraux) : Raccordement flexible 3/4" BSP 120 bar. Serrage raccords à 45 Nm. 
   Essai pression statique 130 bar (marge 10 %) pendant 10 min. Zero fuite. 
   Purge circuit par pompe manuelle (évacuation air) avant mise en pression machine.

5. THERMORÉGULATION : Circuit eau moule : température 80 °C (côté fixe) et 75 °C (côté mobile). 
   Débit contrôlé 8 L/min circuit A et 12 L/min circuit B (débitmètre à palette). 
   Circuit refroidissement tour aéroréfrigérant : 20 °C entrée, ΔT = 5 °C.

6. SÉCHEUR : Mise en service Motan Luxor CA A 80. Objectif point de rosée ≤ –40 °C (granulés PA66 hygroscopique). 
   Mesure après 45 min : –42 °C (hygromètre à condensation). Raccordement air chauffé 80 °C vers trémie ST 120. 
   Contrôle absence fuite thermique (caméra IR Flir TG165 : Δ max 3 °C).

7. DÉTECTEUR MÉTAL : Installation Sesotec RAPID VARIO en amont trémie. Réglage sensibilité : 
   Fe Ø 0,8 mm / NFe Ø 1,0 mm / AIS Ø 1,2 mm. Validation par 10 tests pièges (morceaux acier/cuivre/alu). 
   Taux détection 10/10. Aucun faux positif sur 50 passages granulés vides.

8. PARAMÉTRAGE CYCLE : Vitesse injection 85 mm/s, pression maintien 650 bar, temps refroidissement 8 s, 
   ouverture moule 1,2 s. Premier cycle = réglage coussin 5 mm. Batch validation 500 pièces.

9. QUALITÉ : Contrôle dimensionnel premier jet au pied à coulisse digital Mitutoyo 500-197-30 
   (résolution 0,01 mm). Tolérance connecteur ±0,05 mm. Résultat : 0 pièce hors tolérance sur 500.

10. DOCUMENTATION : Rédaction PV de recette (Procès-Verbal) signé Production + Qualité + Maintenance. 
    Formation 3 opérateurs : démarrage/arrêt, changement référence (temps cible 8 min), lecture alarmes HMI.
        """.strip(),
        "resultat": """• Mise en production effective : J+2 (mardi matin, déménagement vendredi soir).
• Cycle stable : 18,5 s moyenne (objectif <20 s).
• Zéro défaut qualité sur batch validation 500 pièces (100 % dans tolérance).
• Procédure mise en service standardisée créée (checklist 42 points) et partagée à l’équipe maintenance 
  pour réplications futures (2 autres presses relocalisées le trimestre suivant)."""
    },
    {
        "num": "03",
        "titre": "Débogage & Optimisation Programme Siemens S7",
        "entreprise": "TE Connectivity",
        "lieu": "Tanger (Maroc)",
        "type": "AUTOMATISME AVANCÉ",
        "equip": "Cellule assemblage/connectique — Automate Siemens S7-1200 CPU 1214C DC/DC/DC + HMI KTP700 Basic + servo Sinamics V90",
        "contexte": "Cellule critique assemblage connecteurs automobile. Cadence nominale 220 pièces/heure. Sous-capacité visible sur OEE mensuel (objectif 85 %, réel 72 %).",
        "symptome": """Temps de cycle moyen mesuré 28,5 s vs nominal 16,5 s (+12 s, +73 %). 
Alarme HMI fréquente : E204 « Timeout bras robot zone pick ». Arrêts aléatoires 3–4 fois/heure. 
Opérateurs redémarrage manuel (perte temps 2 min/arrêt).""",
        "diagnostic": """
• Connexion TIA Portal V16 via Ethernet (CP 1243-1). Monitoring Online des blocs.
• Fonction Trace intégrée S7-1200 : temps d’exécution OB1 = 18 ms. Temps exécution FC12 (gestion robot pick-and-place) 
  = 8,2 s (anormal, devrait être <3,5 s max).
• Analyse chronologique FC12 (Networks 15–22) : temporisation T37 (TON) réglée à 8 000 ms (8 s). 
  Bras robot KUKA KR 6 R900 sixx mesuré au chronomètre + suivi signaux E/S : temps cycle pick-and-place 
  réel = 5 500 ms (5,5 s). Marge sécurité excessive (2,5 s) bloque la suite du cycle.
• Bit oscillant : M45.2 (SIGNAL_PINCE_OK) alimenté par capteur proximité inductif Balluff BES M08MG1-PSC20B. 
  Détection instable : oscillation 15–25 ms due à décalage mécanique capteur (jeu support 2,1 mm, spécification 0,3 mm max). 
  Programme sans temporisation antirebond : bit non stable => FC12 attend validation pince indéfiniment.
• Programme non optimisé : mouvements bras strictement séquentiels : 
  (1) attente pince OK → (2) descente → (3) attente prise → (4) montée → (5) attente libération → (6) retour. 
  Mouvements compatibles non parallélisés (ex: montée bras + translation peuvent être simultanés avec interlocks).
        """.strip(),
        "actions": """
1. MÉCANIQUE : Réglage support capteur Balluff BES M08MG1 : réduction jeu à 0,3 mm via taraudage M8 
   + contre-écrou nylstop. Serrage couple 6 Nm (clé dynamométrique micro 1–10 Nm).

2. PROGRAMMATION TIA PORTAL V16 (modification blocs FC12, FC13) :
   • Network 18 : Remplacement temporisation T37 (TON) 8 000 ms → T37_SÉCURITÉ 6 500 ms 
     (mesure bras max 5 800 ms + marge 700 ms).
   • Network 16 : Ajout temporisation antirebond TON sur E0.5 (SIGNAL_PINCE) : 35 ms stable 
     avant validation bit M45.2. Supprime l’effet oscillation électromagnétique.
   • Réécriture séquence FC13 (nouveau bloc mouvements parallèles) : 
     - Phase 1 : Descente + fermeture pince (si produit présent).
     - Phase 2 : Montée bras + translation latérale simultanées (interlock : translation autorisée 
       uniquement si hauteur Z > seuil sécurité 150 mm).
     - Phase 3 : Ouverture pince + dépose + retour origine.
     - Division cycle en 3 phases parallèles via méthode SET/RSET d’étapes (programmation séquentielle structurée).
   • Optimisation code : suppression de 4 temporisations redondantes, gain temps de cycle théorique 2,1 s.
   • Ajout compteur CTU (C1) : comptage cycles machine. Alarme automatique si dérive temps cycle >10 % 
     vs moyenne glissante 50 cycles (détection précoce dérive mécanique).

3. HMI : Création écran diagnostic « TEMPS_CYCLE » affichant temps réel, moyenne glissante, 
   graphique trend 4 h, et compteur alarmes E204.

4. SÉCURITÉ & VERSIONNING : Sauvegarde programme version 2.4 sur serveur SVN + copie locale .ap14. 
   Export commentaires détaillés dans chaque Network (nom du concepteur, date, justification modification).

5. VALIDATION : Run 4 h à cadence max. Monitoring temps cycle via HMI et chronomètre externe. 
   Résultats : 16,2 s (min) — 16,8 s (moy) — 17,1 s (max). Aucun arrêt E204.
        """.strip(),
        "resultat": """• Temps de cycle : 28,5 s → 16,8 s moyenne (–41 %). Capacité théorique passée de 126 à 214 pièces/heure.
• OEE cellule : 72 % → 84 % (+12 points de pourcentage).
• Alarme E204 : fréquence 3,4 arrêts/h → 0 arrêt sur 4 h test, puis 0,2 arrêt/h en production continue (–94 %).
• Documentation : 14 Networks commentés + nouveau bloc FC13 documenté. Maintenance autonome sur modifications futures."""
    },
    {
        "num": "04",
        "titre": "RCA Banc de Test Électrique — Ligne Faisceaux",
        "entreprise": "Kromberg & Schubert",
        "lieu": "Kénitra (Maroc)",
        "type": "RCA & QUALITÉ",
        "equip": "Banc test électrique Komax Alpha 440 — Tests : continuité 4 fils Kelvin, isolation 500 VDC, résistance contact <10 mΩ",
        "contexte": "Ligne faisceaux moteur 32 voies pour client automobile (norme VW 60330 / LV 112-1). Production 2 000 pièces/jour. Just-in-time client : aucun stock tampon.",
        "symptome": """Taux rejet test continuité anormal : 12,3 % sur référence stable (batch >5 000 pièces sans changement outillage). 
Seuil qualité client <2 %. Conséquence : 246 pièces rebutées/jour x coût rebut estimé 18 € = ~4 430 €/jour de perte. 
Risque blocage ligne just-in-time client.""",
        "diagnostic": """
• Analyse données SPC (Statistical Process Control) 15 jours : dérive lente résistance contact testée. 
  Tolérance : <10 mΩ. Valeurs initiales 3–5 mΩ → dérive progressive 8–12 mΩ → >10 mΩ (rebut).
• Analyse Pareto défauts sur 342 rebut : 68 % liés aux pins de contact du banc (connecteur adaptateur 32 voies), 
  22 % câbles piquage usés, 10 % autres (mauvaise insertion opérateur).
• Inspection visuelle pins : oxydation noire (CuO) sur 14 des 32 pins. Cause racine : absence de nettoyage 
  pins dans la gamme préventive (dernier nettoyage non documenté depuis 8 mois). 
  Fumées de soudure étain-flux (colophane) aspirées partiellement : hotte extraction débit mesuré 80 m³/h 
  (norme 120 m³/h min). Filtre charbon saturé.
• Mesure micro-ohmètre Fluke 1507 : pins oxydés = 15–45 mΩ. Pins neufs de rechange = <2 mΩ.
• Analyse huile/contact : présence résidus flux isolant sur 6 pins (contamination électrique).
        """.strip(),
        "actions": """
1. ARRÊT & CONSIGNATION : Arrêt immédiat ligne test (LOTO électrique banc 24 VDC + 500 VDC isolation). 
   Signalisation. Isolation batch produit depuis dernière validation qualité OK (rebut contrôlé 1 200 pièces, 
   retri manuel 94 % bonnes, 6 % rebut confirmé).

2. REMPLACEMENT CONNECTEUR ADAPTATEUR : Dépose jeu pins contact usés 32 voies. Pose connecteur neuf 
   référence constructeur Komax (p/n 100.456.32). Serrage couple bague verrouillage 2,5 Nm.

3. RÉGLAGE MÉCANIQUE : Vérin de mise en position banc : réglage pression air 3,2 bar → 4,0 bar. 
   Contrôle à la jauge analogique (4,0 ±0,1 bar). Vérification non-déformation embout mâle-femelle 
   (pas de marquage plastique après 50 insertions tests).

4. MAINTENANCE EXTRACTION : Remplacement filtre charbon actif hotte (dimensions 500×300×50 mm, 
   classe F7). Révision ventilateur centrifuge : nettoyage roue, graissage paliers. 
   Mesure débit post-intervention : 130 m³/h (>norme 120). Ajout hotte aspirante mobile Plymovent en renfort 
   ponctuel sur poste soudure amont.

5. MODIFICATION GMAO / PROCÉDURES :
   • Nouvelle gamme préventive : nettoyage pins test à la brosse fibre de verre + isopropanol 99 % 
     — fréquence bi-hebdomadaire (toutes les 2 semaines) au lieu de trimestriel.
   • Ajout contrôle qualité mensuel : mesure micro-ohmètre sur 5 pins aléatoires. Seuil <5 mΩ.
   • Création fiche réflexe : si dérive SPC >6 mΩ sur 3 mesures consécutives → arrêt immédiat + contrôle pins.
   • Formation 2 opérateurs + 1 technicien maintenance à la nouvelle procédure.

6. CONTRÔLE ATTRIBUTAIRE : Mise en place test gabarit 32/32 pins tous les matins avant démarrage 
   (détection précoce oxydation ou déformation).
        """.strip(),
        "resultat": """• Taux rejet test continuité : 12,3 % → 1,1 % moyenne sur 30 jours. Stabilisé sous seuil client 2 %.
• Économie : ~220 pièces/jour sauvées x 18 € x 220 jours/an = ~870 000 €/an de rebut évité 
  (réduction drastique rebut + retri inutile).
• Procédure RCA formalisée (fiche 5 pourquoi + Ishikawa) partagée aux 4 lignes test du site 
  (Kromberg + 3 sous-traitants internes).
• Temps contrôle qualité départ : réduit de 25 min/batch à 10 min/batch (test plus fiable, moins de retest)."""
    },
    {
        "num": "05",
        "titre": "Dépannage Hydraulique Presse Découpe / Pliage",
        "entreprise": "Sovireso",
        "lieu": "Saint-Laurent-sur-Sèvre (85)",
        "type": "CURATIF MÉCANIQUE",
        "equip": "Presse hydraulique découpe/pliage AMADA HFE 80-25 — Vérin double effet Ø100/Ø70 mm — Centrale 80 bar / 25 L/min — Distributeur 4/3 voies NG10 proportionnel",
        "contexte": "Machine unique sur site pour découpe tôlerie fine (acier DC01 0,8–2 mm, inox 1,4301 1–1,5 mm). Pas de remplaçant. Production programme client 480 pièces/jour.",
        "symptome": """Perte brutale de pression lors de la descente rapide vers position pliage. 
Pression manomètre principal : chute de 80 bar à 32 bar. Bruit violent type « coup de bélier » à chaque inversion 
descente/montée. Arrêt machine. Production interrompue depuis 2 h.""",
        "diagnostic": """
• Mesures manomètres multi-points (WIKA 0–160 bar, classe 1,0) :
  - P pompe (amont filtre) : 80 bar stable.
  - P ligne descente (amont vérin tige) : 32 bar (anormal, devrait être 78–80 bar).
  - P ligne montée (amont vérin fond) : 78 bar (normal).
  → Panne localisée sur circuit descente (clapet anti-retour ou fuite interne vérin).

• Auscultation + stéthoscope mécanique : coup de bélier localisé accumulateur hydropneumatique 
  HAB 1-330 (bladder) côté montée. Fréquence bruit synchrone avec inversion cycle.

• Démontage clapet anti-retour hydraulique NG25 (ligne descente) : présence éclat joint torique NBR 90SH 
  Ø12×2 mm coincé dans le siège clapet (corps étranger). Joint dégradé par huile surchauffée : 
  température huile mesurée 68 °C (normale exploitation <55 °C). Cause surchauffe : filtre retour saturé 
  => pompe travaille en dépression/cavitation.

• Contrôle accumulateur HAB 1-330 : précharge azote mesurée 12 bar (normale constructeur 65 bar). 
  Membrane bladder fissurée visible lors démontage contrôle (huile dans précharge gaz).

• Analyse huile : viscosité 32 cSt → 28 cSt (dégradation polymère). Indice contamination NAS 10 
  (seuil acceptable NAS 8). Filtre retour 25 µm saturé (contamination ferreuse usure pompe à palettes).
        """.strip(),
        "actions": """
1. CONSIGNATION ÉNERGÉTIQUE : Coupure moteur pompe (disjoncteur magnétothermique 32 A). 
   Condamnateur cadenas + panneau « Ne pas démarrer — Maintenance ». Vidange huile 120 L 
   dans récipient homologué UN 1H2. Neutralisation absorbent si résidus sol.

2. CLAPET ANTI-RETOUR NG25 : Dépose complète. Nettoyage siège en carbure de tungstène 
   par rodage à la pâte à roder (grain 600). Contrôle étanchéité au kérosène (0 goutte/30 s). 
   Remplacement joint torique NBR 90SH Ø12×2 mm (parker) + joint plat cuivre anneau NG25. 
   Remontage : serrage croisé vis M10 classe 12.9 au couple 35 Nm (clé dynamométrique 3/8").

3. ACCUMULATEUR HAB 1-330 : Remplacement complet (vase + précharge). Précharge azote vérifiée 
   au manomètre de contrôle à 65 bar (6,5 MPa). Contrôle re-étanchéité 24 h après (perte 0 bar).

4. FILTRATION : Remplacement filtre retour 25 µm (MP Filtri MR2503A10A). Remplacement filtre aspiration 
   100 µm. Remplissage huile neuve Shell Tellus S2 MX 46 (ISO VG 46) 120 L. Purge par cycles lents 
   sans charge (10 cycles montée/descente vitesse 10 %).

5. CONTRÔLE POMPE : Démontage inspection palettes (pompe à palettes Denison T6C). 
  Jeu latéral mesuré 0,12 mm (tolérance constructeur 0,15 mm max) => acceptable mais programmée 
  révision T4 dans 6 mois (préventif opportuniste).

6. MISE EN PRESSION : Montée progressive par paliers : 0 → 20 bar (5 min) → 40 bar (5 min) → 80 bar (5 min). 
   Contrôle fuites aux raccords : 0 fuite (papier absorbant 10 min). Contrôle température huile après 30 min : 48 °C (OK).

7. ESSAIS SOUS CHARGE : Découpe tôle acier DC01 2 mm, 50 cycles continus. 
   Pression mesurée stable 80 ±2 bar sur l’ensemble du cycle (descente + effort + montée). 
   Disparition complète des coups de bélier (auscultation + capteur vibration portable SKF CMAS 100-SL : <2 mm/s).
        """.strip(),
        "resultat": """• Temps d’intervention : 3 h 20 (prévision constructeur 1 journée = 7 h). Production relancée même journée.
• Pression rétablie : 80 bar stable ±2 bar sur cycle complet descente/montée/pliage.
• Disparition totale à-coups hydrauliques => réduction vibration structure, allongement durée de vie vérins + tuyauterie.
• Maintenance préventive avancée : remplacement huile + filtres programmés avec marge de 2 mois (préventif opportuniste 
  intégré à la gamme annuelle). Révision pompe T4 programmée en planning SAP PM."""
    },
    {
        "num": "06",
        "titre": "Planification Maintenance Annuelle & Pilotage GMAO",
        "entreprise": "TE Connectivity",
        "lieu": "Tanger (Maroc)",
        "type": "PLANIFICATION & FIABILITÉ",
        "equip": "Site industriel complet : 142 équipements référencés (injection, assemblage, test, utilités, bâtiment)",
        "contexte": "Site de 450 personnes, production 24/5. Objectif groupe TE : taux réalisation maintenance préventive >90 % et réduction arrêts imprévus de 25 % sur l’exercice.",
        "symptome": """Taux réalisation maintenance préventive année N : 68 % (233/342 OT planifiés). 
Arrêts imprévus représentent 23 % du temps disponible production. Coût maintenance en croissance +12 % vs budget. 
MTBF moyen site 420 h (benchmark industrie connectique 650 h).""",
        "diagnostic": """
• Analyse data GMAO (SAP PM) + supervision (WinCC V7.5) sur 12 mois :
  - 342 OT préventifs planifiés, 233 réalisés (68 %), 109 reportés.
  - Pareto causes reports : 45 % conflit production (pas de fenêtre machine disponible), 
    30 % pièces détachées non disponibles (délai approvisionnement 5–7 j), 
    15 % sous-traitant spécialisé indisponible, 10 % autres.

• Absence classification criticité équipements : tous les équipements traités avec même priorité 
  (pas de matrice ABC). Conséquence : un compresseur d’air (critique) et un éclairage bureau (non critique) 
  avaient la même fréquence de préventif.

• Analyse MTBF/MTTR par famille : presses injection MTBF 380 h (faible), bancs test MTBF 620 h (moyen), 
  utilités MTBF 850 h (bon). Pas de plan d’action spécifique sur presses.
        """.strip(),
        "actions": """
1. COLLECTE & ANALYSE DATA : Export GMAO (OT, OT réalisés, reports, coûts). Export supervision WinCC 
   (arrêts, codes défauts, durées). Croisement Excel Power Query + tableau de bord Power BI simplifié.

2. MATRICE CRITICITÉ ABC (méthode RCM simplifiée) :
   • Classe A (20 équipements, 14 %) : critique production + sécurité + environnement. 
     → Préventif mensuel/bimestriel. Stock pièces critiques local (capteurs, vannes, joints). 
     • Exemples : presses injection, compresseur air 8 bar, tour aéroréfrigérant principal.
   • Classe B (45 équipements, 32 %) : important mais remplaçable. 
     → Préventif trimestriel. Approvisionnement pièces sous 48 h.
   • Classe C (77 équipements, 54 %) : standard / non critique. 
     → Préventif semestriel/annuel. Maintenance curative à la demande.

3. PLANNING ANNUEL INTÉGRÉ : Élaboration planning 12 mois dans SAP PM avec fenêtres fixes 
   négociées avec production : samedi matin (4 h), arrêts planifiés trimestriels (8 h), 
   vacances annuelles (grandes révisions). Réduction conflits production de 45 % à 8 %.

4. APPROVISIONNEMENT PIÈCES : Négociation fournisseurs locaux (Siemens, Festo, Parker). 
   Création stock sécurité classe A (consommables : capteurs inductifs, photoelectriques, vannes 5/2, 
   joints toriques standards). Délai moyen pièces critiques : 5 j → 2 j.

5. SOUS-TRAITANCE : Regroupement interventions par lot (thermographie, étalonnage instruments, 
   révision pompes spécialisées). Appel d’offres 3 fournisseurs. Réduction coûts sous-traitance : –18 %.

6. PILOTAGE KPI : Tableau de bord mensuel affiché maintenance + production :
   - Taux réalisation préventif (%).
   - MTBF (h) et MTTR (min) par classe.
   - Coût maintenance / tonne produit (€/t).
   - Top 3 pannes du mois + actions en cours.

7. PROJETS D’AMÉLIORATION : Élaboration et soumission 5 projets CAPEX/OPEX : 
   remplacement 2 variateurs obsolètes, mise en place vibration monitoring sur 3 groupes critiques, 
   formation interne automatisme (3 techniciens). 3/5 validés et réalisés en interne.
        """.strip(),
        "resultat": """• Taux réalisation préventif : 68 % → 94 % (année N+1). Objectif groupe 90 % dépassé.
• Arrêts imprévus : –28 % vs année N (production gagne 340 h/an disponibles).
• MTBF site : 420 h → 610 h (+45 %), rapproché benchmark industrie.
• Coûts maintenance : –15 % sous-traitants (achats groupés) + –8 % pièces (stock optimisé + négociation).
• 3 projets amélioration réalisés en interne avec ROI <18 mois."""
    },
]

# ============================================================
# 1) WORD
# ============================================================
def create_word():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.enum.table import WD_TABLE_ALIGNMENT

    doc = Document()

    # Marges étroites pour design magazine
    sections = doc.sections[0]
    sections.top_margin = Cm(2)
    sections.bottom_margin = Cm(2)
    sections.left_margin = Cm(2)
    sections.right_margin = Cm(2)

    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(10)
    style.element.rPr.rFonts.set(qn('w:eastAsia'), 'Calibri')

    # --- PAGE TITRE ---
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("PORTFOLIO TECHNIQUE")
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    r.font.name = 'Calibri Light'

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("d'Interventions Maintenance Industrielle")
    r.font.size = Pt(18)
    r.font.color.rgb = RGBColor(0xF9, 0x73, 0x16)
    r.font.name = 'Calibri'

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(PROFILE["name"])
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(PROFILE["contact"])
    r.font.size = Pt(10)
    r.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    p.space_after = Pt(20)

    # Résumé
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    r = p.add_run(PROFILE["summary"])
    r.font.size = Pt(10)
    r.italic = True
    r.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
    doc.add_paragraph()

    # --- CHIFFRES CLÉS ---
    doc.add_heading("Chiffres Clés", level=1)
    table = doc.add_table(rows=1, cols=4)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cells = table.rows[0].cells
    headers = [("10+", "Années d'expérience"), ("142", "Équipements gérés"), ("4", "Pays d'intervention"), ("∞", "Pannes résolues")]
    for i, (num, lab) in enumerate(headers):
        p = cells[i].paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(num)
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xF9, 0x73, 0x16)
        p.add_run("\n" + lab).font.size = Pt(9)
    doc.add_paragraph()

    # --- COMPÉTENCES ---
    doc.add_heading("Compétences Techniques", level=1)
    comp_data = [
        ("⚡ Électrotechnique & Automatisme", "Siemens S7-1200/1500, TIA Portal V16, Step 7, Ladder/GRAPH, PROFINET, Sinamics V90, câblage BT, schémas électriques IEC, variateurs, servo-motors, capteurs (inductifs, photo, ultrasons, encodeurs)"),
        ("🔧 Mécanique & Fluidique", "Presses injection (KraussMaffei), lignes PET (Krones), hydraulique 80–200 bar (pompes, vérins, distributeurs NG6/NG10, accumulateurs), pneumatique 5–8 bar (distributeurs 5/2, vérins, préparation air), transmissions (engrenages, courroies, chaînes), guidages, roulements"),
        ("📊 Méthodes & Gestion", "GMAO SAP PM, planification préventive/corrective, analyse RCA (5 pourquoi, Ishikawa, Pareto), collecte data SPC, tableaux de bord KPI (MTBF, MTTR, OEE), consignation LOTO (ISO 14118), normes HSE"),
    ]
    for title, detail in comp_data:
        p = doc.add_paragraph()
        r = p.add_run(title + " : ")
        r.bold = True
        r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        p.add_run(detail).font.color.rgb = RGBColor(0x37, 0x41, 0x51)
    doc.add_paragraph()

    # --- FICHES ---
    doc.add_heading("Fiches d'Intervention Détaillées", level=1)
    doc.add_paragraph("Format : Contexte → Diagnostic technique mesuré → Actions détaillées → Résultats chiffrés")

    for f in FICHES:
        doc.add_heading(f"{f['num']}. {f['titre']}", level=2)

        # Tableau meta
        table = doc.add_table(rows=2, cols=3)
        table.style = 'Light Grid Accent 1'
        cells = table.rows[0].cells
        cells[0].text = f"Entreprise : {f['entreprise']}"
        cells[1].text = f"Site : {f['lieu']}"
        cells[2].text = f"Type : {f['type']}"
        cells = table.rows[1].cells
        cells[0].text = f"Équipement : {f['equip']}"
        cells[0].merge(cells[1]).merge(cells[2])
        doc.add_paragraph()

        p = doc.add_paragraph()
        r = p.add_run("Contexte : ")
        r.bold = True
        r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        p.add_run(f['contexte'])

        p = doc.add_paragraph()
        r = p.add_run("Symptôme : ")
        r.bold = True
        r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        p.add_run(f['symptome'])

        p = doc.add_paragraph()
        r = p.add_run("Diagnostic technique : ")
        r.bold = True
        r.font.color.rgb = RGBColor(0xF9, 0x73, 0x16)
        p.add_run(f['diagnostic'])

        p = doc.add_paragraph()
        r = p.add_run("Actions réalisées : ")
        r.bold = True
        r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        for line in f['actions'].split('\n'):
            if line.strip():
                p = doc.add_paragraph(line.strip(), style='List Number')
                p.paragraph_format.space_after = Pt(4)

        p = doc.add_paragraph()
        r = p.add_run("Résultat : ")
        r.bold = True
        r.font.color.rgb = RGBColor(0x05, 0x96, 0x69)
        run = p.add_run(f['resultat'])
        run.font.color.rgb = RGBColor(0x05, 0x96, 0x69)
        run.bold = True
        p.paragraph_format.space_after = Pt(12)
        doc.add_paragraph()

    # --- MÉTHODOLOGIE ---
    doc.add_heading("Méthodologie d'Intervention", level=1)
    meth = [
        ("1. SÉCURISATION", "Consignation LOTO (ISO 14118). Protection individuelle (EPI). Signalisation de zone."),
        ("2. DIAGNOSTIC MESURÉ", "Données quantitatives : pression (bar), tension (V), temps (ms), température (°C), dimension (mm). Aucune hypothèse sans mesure."),
        ("3. ANALYSE CAUSE RACINE", "Arbre de défaillance, Pareto, 5 pourquoi. Identification cause physique + cause systémique."),
        ("4. INTERVENTION CIBLÉE", "Procédure pas à pas avec outillages spécifiés, couples de serrage, références pièces."),
        ("5. VALIDATION & REPORTING", "Essais sous charge, mesures avant/après. Saisie GMAO. Mise à jour documentation. Formation transfert."),
    ]
    for title, detail in meth:
        p = doc.add_paragraph()
        r = p.add_run(title + " — ")
        r.bold = True
        r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        p.add_run(detail)

    # Footer
    doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(f"{PROFILE['name']} — {PROFILE['title']}\n{PROFILE['contact']}")
    r.italic = True
    r.font.size = Pt(9)
    r.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    path = os.path.join(OUTPUT_DIR, "Portfolio_Salah_Barki_v2.docx")
    doc.save(path)
    print(f"[OK] Word créé : {path}")


# ============================================================
# 2) POWERPOINT (.pptx) — Design minimaliste premium
# ============================================================
def create_ppt():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Couleurs
    C_BG = RGBColor(0xFF, 0xFF, 0xFF)
    C_DARK = RGBColor(0x0F, 0x17, 0x2A)
    C_ACCENT = RGBColor(0xF9, 0x73, 0x16)
    C_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
    C_TEXT_GRAY = RGBColor(0x37, 0x41, 0x51)
    C_GREEN = RGBColor(0x05, 0x96, 0x69)
    C_LIGHT_TEXT = RGBColor(0x94, 0xA3, 0xB8)

    def add_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = C_BG

    def add_rect(slide, left, top, width, height, color):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def add_text(slide, left, top, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        return box

    def add_bullet_text(slide, left, top, width, height, lines, font_size, color, bullet='•'):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"{bullet} {line}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = 'Calibri'
            p.space_after = Pt(6)
        return box

    # ===== SLIDE 1 : TITRE =====
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    # Barre accent haut
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), C_ACCENT)
    # Titre
    add_text(slide, Inches(1), Inches(2.2), Inches(11.333), Inches(1.2), "PORTFOLIO TECHNIQUE", 48, C_DARK, bold=True, align=PP_ALIGN.CENTER, font_name='Calibri Light')
    add_text(slide, Inches(1), Inches(3.4), Inches(11.333), Inches(0.8), "d'Interventions Maintenance Industrielle", 22, C_ACCENT, align=PP_ALIGN.CENTER)
    # Ligne séparatrice
    add_rect(slide, Inches(5.5), Inches(4.2), Inches(2.333), Inches(0.01), C_LIGHT_TEXT)
    add_text(slide, Inches(1), Inches(4.5), Inches(11.333), Inches(0.6), PROFILE["name"], 26, C_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5.2), Inches(11.333), Inches(0.5), PROFILE["contact"], 12, C_LIGHT_TEXT, align=PP_ALIGN.CENTER)

    # ===== SLIDE 2 : PROFIL =====
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), C_DARK)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), "PROFIL & EXPERTISE", 24, C_BG, bold=True)
    # 2 colonnes
    add_rect(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.8), C_GRAY)
    add_text(slide, Inches(0.7), Inches(1.4), Inches(5.6), Inches(5.4), PROFILE["summary"], 14, C_TEXT_GRAY)
    # Chiffres à droite
    stats = [("10+", "Ans d'exp."), ("142", "Équip. gérés"), ("6", "Fiches phares"), ("4", "Pays")]
    y = 1.4
    for num, lab in stats:
        add_text(slide, Inches(7.2), Inches(y), Inches(2), Inches(0.5), num, 28, C_ACCENT, bold=True)
        add_text(slide, Inches(9), Inches(y+0.1), Inches(3), Inches(0.4), lab, 14, C_TEXT_GRAY)
        y += 1.3

    # ===== SLIDE 3 : COMPÉTENCES =====
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), C_DARK)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), "COMPÉTENCES & TECHNOLOGIES", 24, C_BG, bold=True)
    comp_blocks = [
        ("⚡ ÉLECTROTECHNIQUE", "Siemens S7-1200/1500\nTIA Portal V16 / Step 7\nPROFINET / Sinamics V90\nCâblage BT / Schémas IEC\nVariateurs / Servos"),
        ("🔧 MÉCANIQUE", "Presses injection\nLignes PET / Canettes\nTransmissions / Guidages\nPneumatique 5–8 bar\nHydraulique 80–200 bar"),
        ("📊 MÉTHODES", "GMAO SAP PM\nRCA / 5 pourquoi\nSPC / Pareto\nPlanification préventive\nKPI (MTBF, MTTR, OEE)"),
    ]
    x = 0.5
    for title, body in comp_blocks:
        add_rect(slide, Inches(x), Inches(1.3), Inches(4), Inches(5.5), C_GRAY)
        add_text(slide, Inches(x+0.2), Inches(1.5), Inches(3.6), Inches(0.5), title, 14, C_ACCENT, bold=True)
        add_text(slide, Inches(x+0.2), Inches(2.1), Inches(3.6), Inches(4.6), body, 13, C_TEXT_GRAY)
        x += 4.4

    # ===== FICHES =====
    for f in FICHES:
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        # Header
        add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), C_DARK)
        add_text(slide, Inches(0.4), Inches(0.15), Inches(12), Inches(0.5), f"{f['num']}  {f['titre']}", 22, C_BG, bold=True)
        add_text(slide, Inches(0.4), Inches(0.6), Inches(12), Inches(0.4), f"{f['entreprise']} — {f['lieu']}  |  {f['type']}", 11, C_LIGHT_TEXT)

        # Colonne gauche : DIAGNOSTIC (gris)
        add_rect(slide, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.8), C_GRAY)
        add_text(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.4), "DIAGNOSTIC TECHNIQUE", 13, C_ACCENT, bold=True)
        diag_text = f"Équipement : {f['equip']}\n\nContexte : {f['contexte']}\n\nSymptôme : {f['symptome']}\n\n{f['diagnostic']}"
        add_text(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5.2), diag_text, 11, C_TEXT_GRAY)

        # Colonne droite : ACTIONS (blanc avec ligne accent)
        add_rect(slide, Inches(6.7), Inches(1.3), Inches(0.05), Inches(5.8), C_ACCENT)
        add_text(slide, Inches(6.9), Inches(1.4), Inches(6), Inches(0.4), "ACTIONS DÉTAILLÉES", 13, C_DARK, bold=True)
        # Actions en bullet
        lines = [l.strip() for l in f['actions'].split('\n') if l.strip()]
        add_bullet_text(slide, Inches(6.9), Inches(1.8), Inches(6), Inches(4.5), lines, 11, C_TEXT_GRAY)

        # Footer résultat (bandeau vert)
        add_rect(slide, Inches(0.3), Inches(6.4), Inches(12.7), Inches(0.9), C_GREEN)
        add_text(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), f"RÉSULTAT : {f['resultat']}", 12, C_BG, bold=True)

    # ===== SLIDE FINALE =====
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), C_ACCENT)
    add_text(slide, Inches(1), Inches(2.5), Inches(11.333), Inches(1), "MÉTHODOLOGIE D'INTERVENTION", 32, C_DARK, bold=True, align=PP_ALIGN.CENTER, font_name='Calibri Light')
    meth_lines = [
        "1. SÉCURISATION — Consignation LOTO (ISO 14118), EPI, signalisation",
        "2. DIAGNOSTIC MESURÉ — Bar, V, °C, mm, ms. Aucune hypothèse sans mesure",
        "3. RCA — 5 pourquoi, Ishikawa, Pareto. Cause physique + cause systémique",
        "4. INTERVENTION CIBLÉE — Procédure pas à pas, outillages, couples, références",
        "5. VALIDATION & REPORTING — Essais charge, mesures avant/après, GMAO, formation"
    ]
    y = 3.5
    for line in meth_lines:
        add_text(slide, Inches(2), Inches(y), Inches(9.333), Inches(0.5), line, 16, C_TEXT_GRAY, align=PP_ALIGN.CENTER)
        y += 0.6
    add_text(slide, Inches(1), Inches(6.3), Inches(11.333), Inches(0.5), f"{PROFILE['name']} — {PROFILE['contact']}", 12, C_LIGHT_TEXT, align=PP_ALIGN.CENTER)

    path = os.path.join(OUTPUT_DIR, "Portfolio_Salah_Barki_v2.pptx")
    prs.save(path)
    print(f"[OK] PowerPoint créé : {path}")


# ============================================================
# 3) PDF — Design magazine technique
# ============================================================
def create_pdf():
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
    from reportlab.pdfgen.canvas import Canvas

    path = os.path.join(OUTPUT_DIR, "Portfolio_Salah_Barki_v2.pdf")
    doc = SimpleDocTemplate(path, pagesize=A4,
                          rightMargin=1.8*cm, leftMargin=1.8*cm,
                          topMargin=1.8*cm, bottomMargin=1.8*cm)

    styles = getSampleStyleSheet()

    # Styles custom
    s_title = ParagraphStyle(name='s_title_v2', parent=styles['Title'], fontSize=26, textColor=colors.HexColor('#0f172a'), alignment=TA_CENTER, spaceAfter=8, fontName='Helvetica-Bold')
    s_sub = ParagraphStyle(name='s_sub_v2', fontSize=14, textColor=colors.HexColor('#f97316'), alignment=TA_CENTER, spaceAfter=12, fontName='Helvetica-Bold')
    s_name = ParagraphStyle(name='s_name_v2', fontSize=16, textColor=colors.HexColor('#0f172a'), alignment=TA_CENTER, spaceAfter=6, fontName='Helvetica-Bold')
    s_contact = ParagraphStyle(name='s_contact_v2', fontSize=9, textColor=colors.HexColor('#94a3b8'), alignment=TA_CENTER, spaceAfter=14)
    s_h1 = ParagraphStyle(name='s_h1_v2', parent=styles['Heading1'], fontSize=16, textColor=colors.HexColor('#0f172a'), spaceAfter=10, spaceBefore=14, fontName='Helvetica-Bold', leftIndent=0, borderWidth=0, borderColor=colors.HexColor('#f97316'), borderPadding=5)
    s_h2 = ParagraphStyle(name='s_h2_v2', parent=styles['Heading2'], fontSize=13, textColor=colors.HexColor('#0f172a'), spaceAfter=8, spaceBefore=12, fontName='Helvetica-Bold')
    s_h3 = ParagraphStyle(name='s_h3_v2', parent=styles['Heading3'], fontSize=11, textColor=colors.HexColor('#f97316'), spaceAfter=6, spaceBefore=8, fontName='Helvetica-Bold')
    s_body = ParagraphStyle(name='s_body_v2', parent=styles['Normal'], fontSize=9.5, leading=13, alignment=TA_JUSTIFY, spaceAfter=6, textColor=colors.HexColor('#374151'))
    s_bullet = ParagraphStyle(name='s_bullet_v2', parent=s_body, leftIndent=14, spaceAfter=4, bulletIndent=6, bulletFontName='Helvetica-Bold', bulletFontSize=9, bulletColor=colors.HexColor('#f97316'))
    s_result = ParagraphStyle(name='s_result_v2', parent=s_body, textColor=colors.HexColor('#059669'), fontSize=9.5, leading=13, spaceAfter=8, leftIndent=6, backColor=colors.HexColor('#f0fdf4'), borderColor=colors.HexColor('#059669'), borderWidth=2, borderPadding=6, leftPadding=8, rightPadding=8, topPadding=6, bottomPadding=6)
    s_diag = ParagraphStyle(name='s_diag_v2', parent=s_body, backColor=colors.HexColor('#f8fafc'), leftIndent=6, borderColor=colors.HexColor('#e2e8f0'), borderWidth=1, borderPadding=6, leftPadding=8, rightPadding=8, topPadding=6, bottomPadding=6)
    s_meta = ParagraphStyle(name='s_meta_v2', fontSize=8.5, textColor=colors.HexColor('#64748b'), alignment=TA_LEFT, spaceAfter=4, fontName='Helvetica')

    story = []

    # --- TITRE ---
    story.append(Paragraph("PORTFOLIO TECHNIQUE", s_title))
    story.append(Paragraph("d'Interventions Maintenance Industrielle", s_sub))
    story.append(Paragraph(PROFILE["name"], s_name))
    story.append(Paragraph(PROFILE["contact"].replace('📍','').replace('📱','').replace('✉️',''), s_contact))
    story.append(Spacer(1, 0.2*cm))

    # Résumé encadré
    story.append(Paragraph(PROFILE["summary"], ParagraphStyle(name='s_sum_v2', parent=s_body, backColor=colors.HexColor('#f8fafc'), borderColor=colors.HexColor('#e2e8f0'), borderWidth=1, borderPadding=8, leftPadding=10, rightPadding=10, topPadding=8, bottomPadding=8)))
    story.append(Spacer(1, 0.4*cm))

    # Chiffres clés
    data = [
        [Paragraph("<b>10+</b>", ParagraphStyle(name='sk1', alignment=TA_CENTER, fontSize=18, textColor=colors.HexColor('#f97316'))),
         Paragraph("<b>142</b>", ParagraphStyle(name='sk2', alignment=TA_CENTER, fontSize=18, textColor=colors.HexColor('#f97316'))),
         Paragraph("<b>4</b>", ParagraphStyle(name='sk3', alignment=TA_CENTER, fontSize=18, textColor=colors.HexColor('#f97316'))),
         Paragraph("<b>6</b>", ParagraphStyle(name='sk4', alignment=TA_CENTER, fontSize=18, textColor=colors.HexColor('#f97316')))],
        [Paragraph("Années d'expérience", ParagraphStyle(name='sk5', alignment=TA_CENTER, fontSize=9, textColor=colors.HexColor('#64748b'))),
         Paragraph("Équipements gérés", ParagraphStyle(name='sk6', alignment=TA_CENTER, fontSize=9, textColor=colors.HexColor('#64748b'))),
         Paragraph("Pays", ParagraphStyle(name='sk7', alignment=TA_CENTER, fontSize=9, textColor=colors.HexColor('#64748b'))),
         Paragraph("Fiches phares", ParagraphStyle(name='sk8', alignment=TA_CENTER, fontSize=9, textColor=colors.HexColor('#64748b')))],
    ]
    t = Table(data, colWidths=[4*cm, 4*cm, 4*cm, 4*cm])
    t.setStyle(TableStyle([
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('LEFTPADDING', (0,0), (-1,-1), 6),
        ('RIGHTPADDING', (0,0), (-1,-1), 6),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#f8fafc')),
    ]))
    story.append(t)
    story.append(Spacer(1, 0.4*cm))

    # Compétences
    story.append(Paragraph("COMPÉTENCES & TECHNOLOGIES", s_h1))
    comps = [
        ("ÉLECTROTECHNIQUE & AUTOMATISME", "Siemens S7-1200/1500, TIA Portal V16, Step 7, Ladder/GRAPH, PROFINET, Sinamics V90, câblage BT, schémas IEC, variateurs, servo-moteurs, capteurs inductifs/photo/encodeurs."),
        ("MÉCANIQUE & FLUIDIQUE", "Presses injection KraussMaffei, lignes PET Krones, hydraulique 80–200 bar (pompes, vérins, distributeurs NG6/NG10, accumulateurs), pneumatique 5–8 bar, transmissions, guidages, roulements."),
        ("MÉTHODES & MANAGEMENT", "GMAO SAP PM, planification préventive/corrective, analyse RCA (5 pourquoi, Ishikawa, Pareto), SPC, KPI (MTBF, MTTR, OEE), consignation LOTO ISO 14118, normes HSE."),
    ]
    for cat, detail in comps:
        story.append(Paragraph(f"<b>{cat}</b> — {detail}", s_body))
    story.append(PageBreak())

    # FICHES
    story.append(Paragraph("FICHES D'INTERVENTION DÉTAILLÉES", s_h1))
    story.append(Paragraph("Format structuré : Contexte → Diagnostic technique mesuré → Actions détaillées (outillages, couples, références) → Résultats chiffrés.", s_body))
    story.append(Spacer(1, 0.3*cm))

    for f in FICHES:
        story.append(Paragraph(f"{f['num']}. {f['titre']}", s_h2))
        story.append(Paragraph(f"<b>{f['entreprise']}</b> — {f['lieu']}  |  {f['type']}", s_meta))
        story.append(Spacer(1, 0.1*cm))

        story.append(Paragraph(f"<b>Équipement :</b> {f['equip']}", s_body))
        story.append(Paragraph(f"<b>Contexte :</b> {f['contexte']}", s_body))
        story.append(Paragraph(f"<b>Symptôme :</b> {f['symptome']}", s_body))
        story.append(Paragraph(f"<b>Diagnostic technique :</b><br/>{f['diagnostic']}", s_diag))

        story.append(Paragraph("<b>Actions réalisées :</b>", s_h3))
        for line in f['actions'].split('\n'):
            if line.strip():
                story.append(Paragraph(line.strip(), s_bullet))

        story.append(Paragraph(f"<b>RÉSULTAT :</b> {f['resultat']}", s_result))
        story.append(Spacer(1, 0.3*cm))

    # Méthodologie
    story.append(PageBreak())
    story.append(Paragraph("MÉTHODOLOGIE D'INTERVENTION", s_h1))
    meths = [
        ("1. SÉCURISATION", "Consignation LOTO (NF EN ISO 14118). EPI complète. Signalisation de zone et condamnateur cadenas."),
        ("2. DIAGNOSTIC MESURÉ", "Données quantitatives : pression (bar), tension (V), température (°C), dimension (mm), temps (ms). Aucune hypothèse sans mesure instrumentée."),
        ("3. RCA", "Analyse cause racine : arbre de défaillance, Pareto, 5 pourquoi, Ishikawa. Identification cause physique + cause systémique (procédure, formation, maintenance)."),
        ("4. INTERVENTION CIBLÉE", "Procédure pas à pas : outillages spécifiés, couples de serrage, références constructeur, normes appliquées."),
        ("5. VALIDATION & REPORTING", "Essais sous charge, mesures avant/après, saisie GMAO, mise à jour documentation, formation transfert de compétences."),
    ]
    for t, d in meths:
        story.append(Paragraph(f"<b>{t}</b> — {d}", s_body))
    story.append(Spacer(1, 0.5*cm))

    story.append(Paragraph(f"<b>{PROFILE['name']}</b> — {PROFILE['title']}<br/>{PROFILE['contact']}", ParagraphStyle(name='s_footer_v2', alignment=TA_CENTER, fontSize=9, textColor=colors.HexColor('#94a3b8'))))

    doc.build(story)
    print(f"[OK] PDF créé : {path}")


# ============================================================
# 4) HTML — Design moderne premium
# ============================================================
def create_html():
    html_content = """<!DOCTYPE html>
<html lang="fr">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Portfolio Technique — Salah Eddine Barki</title>
<style>
:root {
  --bg: #ffffff;
  --dark: #0f172a;
  --accent: #f97316;
  --gray: #f1f5f9;
  --text: #374151;
  --muted: #94a3b8;
  --green: #059669;
  --green-bg: #f0fdf4;
  --border: #e2e8f0;
}
* { margin:0; padding:0; box-sizing:border-box; }
body {
  font-family: 'Segoe UI', Roboto, Helvetica, Arial, sans-serif;
  background: var(--bg);
  color: var(--text);
  line-height: 1.55;
  font-size: 15px;
}
.container { max-width: 980px; margin: 0 auto; padding: 0 24px; }

/* HEADER */
header {
  background: var(--bg);
  border-bottom: 3px solid var(--accent);
  padding: 48px 0 32px;
  text-align: center;
}
header h1 { font-size: 2.8rem; font-weight: 300; color: var(--dark); letter-spacing: -1px; }
header .subtitle { font-size: 1.15rem; color: var(--accent); font-weight: 600; margin: 8px 0; text-transform: uppercase; letter-spacing: 1px; }
header .name { font-size: 1.6rem; color: var(--dark); font-weight: 700; margin-top: 12px; }
header .contact { font-size: 0.9rem; color: var(--muted); margin-top: 10px; }
header .summary {
  max-width: 780px; margin: 20px auto 0;
  background: var(--gray);
  padding: 18px 22px;
  border-radius: 10px;
  font-size: 0.95rem;
  color: var(--text);
  text-align: justify;
  border: 1px solid var(--border);
}

/* STATS */
.stats { display: grid; grid-template-columns: repeat(auto-fit, minmax(140px, 1fr)); gap: 16px; margin: 28px 0; }
.stat-card {
  background: var(--gray);
  border: 1px solid var(--border);
  border-radius: 10px;
  padding: 20px 10px;
  text-align: center;
}
.stat-card .num { font-size: 1.9rem; font-weight: 700; color: var(--accent); }
.stat-card .lab { font-size: 0.8rem; color: var(--muted); text-transform: uppercase; letter-spacing: 0.5px; margin-top: 4px; }

/* SECTIONS */
.section { padding: 36px 0; }
.section-title {
  font-size: 1.4rem; font-weight: 700; color: var(--dark);
  margin-bottom: 18px;
  display: flex; align-items: center; gap: 10px;
}
.section-title::after {
  content: "";
  flex: 1; height: 2px; background: var(--border);
}

/* COMPETENCES */
.comp-grid { display: grid; grid-template-columns: repeat(auto-fit, minmax(280px, 1fr)); gap: 16px; }
.comp-card {
  background: var(--gray);
  border-radius: 10px; padding: 18px;
  border-left: 4px solid var(--accent);
}
.comp-card h3 { font-size: 0.95rem; color: var(--accent); margin-bottom: 8px; text-transform: uppercase; letter-spacing: 0.5px; }
.comp-card p { font-size: 0.9rem; color: var(--text); }

/* FICHE */
.fiche { margin-bottom: 32px; }
.fiche-header {
  background: var(--dark);
  color: #fff;
  padding: 14px 18px;
  border-radius: 10px 10px 0 0;
  display: flex; justify-content: space-between; align-items: center; flex-wrap: wrap; gap: 8px;
}
.fiche-header h3 { font-size: 1.1rem; font-weight: 600; }
.fiche-header .meta { font-size: 0.8rem; color: var(--muted); }
.fiche-header .badge {
  background: var(--accent);
  color: #fff;
  padding: 4px 10px;
  border-radius: 20px;
  font-size: 0.7rem;
  font-weight: 700;
  text-transform: uppercase;
}
.fiche-body {
  background: var(--bg);
  border: 1px solid var(--border);
  border-top: none;
  border-radius: 0 0 10px 10px;
  padding: 18px;
}
.fiche-grid {
  display: grid;
  grid-template-columns: 1fr 1fr;
  gap: 18px;
}
@media (max-width: 768px) { .fiche-grid { grid-template-columns: 1fr; } }

.fiche-col {
  background: var(--gray);
  border-radius: 8px;
  padding: 14px;
}
.fiche-col h4 {
  font-size: 0.8rem; text-transform: uppercase; letter-spacing: 1px;
  color: var(--accent); margin-bottom: 10px;
}
.fiche-col p, .fiche-col li {
  font-size: 0.88rem; color: var(--text); margin-bottom: 6px;
}
.fiche-col ul { list-style: none; padding-left: 0; }
.fiche-col ul li { position: relative; padding-left: 14px; }
.fiche-col ul li::before {
  content: "›"; position: absolute; left: 0; color: var(--accent); font-weight: bold;
}

.result-box {
  background: var(--green-bg);
  border-left: 4px solid var(--green);
  padding: 12px 14px;
  border-radius: 0 8px 8px 0;
  margin-top: 12px;
  font-size: 0.9rem;
  color: var(--green);
  font-weight: 600;
}

/* METHODOLOGY */
.method-grid { display: grid; grid-template-columns: repeat(auto-fit, minmax(220px, 1fr)); gap: 14px; }
.method-card {
  background: var(--gray);
  border-radius: 10px; padding: 16px;
  text-align: center;
  border-top: 3px solid var(--accent);
}
.method-card .step { font-size: 1.6rem; color: var(--accent); font-weight: 700; margin-bottom: 6px; }
.method-card h4 { font-size: 0.9rem; color: var(--dark); margin-bottom: 6px; }
.method-card p { font-size: 0.85rem; color: var(--text); }

/* FOOTER */
footer {
  text-align: center; padding: 30px 0; margin-top: 20px;
  border-top: 1px solid var(--border);
  font-size: 0.85rem; color: var(--muted);
}
</style>
</head>
<body>

<header>
  <div class="container">
    <h1>Portfolio Technique</h1>
    <div class="subtitle">d'Interventions Maintenance Industrielle</div>
    <div class="name">Salah Eddine Barki</div>
    <div class="contact">📍 Amiens, France &nbsp;|&nbsp; 📱 +33 6 88 69 07 04 &nbsp;|&nbsp; ✉️ salahbarki.seb@gmail.com</div>
    <div class="summary">
      Technicien multitechnique senior (10+ ans) spécialisé en diagnostic avancé sur systèmes automatisés Siemens,
      maintenance préventive / corrective / curative sur presses injection, lignes agroalimentaire PET / canette,
      et bancs de test automobile. Expertise couvrant électrotechnique BT, hydraulique (120–200 bar),
      pneumatique (5–8 bar), mécanique de précision et programmation Ladder (TIA Portal V16 / Step 7).
      Certifié en consignation LOTO, analyse RCA, et planification GMAO (SAP PM).
    </div>
    <div class="stats">
      <div class="stat-card"><div class="num">10+</div><div class="lab">Années</div></div>
      <div class="stat-card"><div class="num">142</div><div class="lab">Équipements</div></div>
      <div class="stat-card"><div class="num">6</div><div class="lab">Fiches phares</div></div>
      <div class="stat-card"><div class="num">4</div><div class="lab">Pays</div></div>
    </div>
  </div>
</header>

<section class="section">
  <div class="container">
    <div class="section-title">Compétences & Technologies</div>
    <div class="comp-grid">
      <div class="comp-card">
        <h3>⚡ Électrotechnique & Automatisme</h3>
        <p>Siemens S7-1200/1500, TIA Portal V16, Step 7, Ladder/GRAPH, PROFINET, Sinamics V90, câblage BT, schémas IEC, variateurs, servo-moteurs, capteurs inductifs / photoélectriques / encodeurs.</p>
      </div>
      <div class="comp-card">
        <h3>🔧 Mécanique & Fluidique</h3>
        <p>Presses injection KraussMaffei, lignes PET Krones, hydraulique 80–200 bar (pompes, vérins, distributeurs NG6/NG10, accumulateurs), pneumatique 5–8 bar, transmissions, guidages, roulements.</p>
      </div>
      <div class="comp-card">
        <h3>📊 Méthodes & Management</h3>
        <p>GMAO SAP PM, planification préventive/corrective, analyse RCA (5 pourquoi, Ishikawa, Pareto), SPC, KPI (MTBF, MTTR, OEE), consignation LOTO ISO 14118, normes HSE.</p>
      </div>
    </div>
  </div>
</section>

<section class="section">
  <div class="container">
    <div class="section-title">Fiches d'Intervention Détaillées</div>
    <p style="margin-bottom:22px; color:var(--muted); font-size:0.9rem;">
      Format structuré : Contexte → Diagnostic technique mesuré → Actions détaillées (outillages, couples, références) → Résultats chiffrés.
    </p>
"""

    for f in FICHES:
        actions_html = "\n".join([f"<li>{l.strip()}</li>" for l in f['actions'].split('\n') if l.strip()])
        html_content += f"""
    <div class="fiche">
      <div class="fiche-header">
        <div>
          <h3>{f['num']}. {f['titre']}</h3>
          <div class="meta">{f['entreprise']} — {f['lieu']}</div>
        </div>
        <span class="badge">{f['type']}</span>
      </div>
      <div class="fiche-body">
        <p style="font-size:0.85rem; color:var(--muted); margin-bottom:12px;"><b>Équipement :</b> {f['equip']}</p>
        <div class="fiche-grid">
          <div class="fiche-col">
            <h4>Diagnostic Technique</h4>
            <p><b>Contexte :</b> {f['contexte']}</p>
            <p><b>Symptôme :</b> {f['symptome']}</p>
            <p style="white-space:pre-line;">{f['diagnostic']}</p>
          </div>
          <div class="fiche-col">
            <h4>Actions Détaillées</h4>
            <ul>
              {actions_html}
            </ul>
          </div>
        </div>
        <div class="result-box">RÉSULTAT : {f['resultat']}</div>
      </div>
    </div>
"""

    html_content += """
  </div>
</section>

<section class="section">
  <div class="container">
    <div class="section-title">Méthodologie d'Intervention</div>
    <div class="method-grid">
      <div class="method-card">
        <div class="step">1</div>
        <h4>SÉCURISATION</h4>
        <p>Consignation LOTO (ISO 14118). EPI. Signalisation de zone.</p>
      </div>
      <div class="method-card">
        <div class="step">2</div>
        <h4>DIAGNOSTIC MESURÉ</h4>
        <p>Bar, V, °C, mm, ms. Aucune hypothèse sans mesure.</p>
      </div>
      <div class="method-card">
        <div class="step">3</div>
        <h4>RCA</h4>
        <p>5 pourquoi, Ishikawa, Pareto. Cause physique + systémique.</p>
      </div>
      <div class="method-card">
        <div class="step">4</div>
        <h4>INTERVENTION CIBLÉE</h4>
        <p>Procédure pas à pas, outillages, couples, références.</p>
      </div>
      <div class="method-card">
        <div class="step">5</div>
        <h4>VALIDATION</h4>
        <p>Essais charge, mesures A/A, GMAO, formation.</p>
      </div>
    </div>
  </div>
</section>

<footer>
  <div class="container">
    <strong>Salah Eddine Barki</strong> — Technicien de Maintenance Industrielle<br>
    📱 +33 6 88 69 07 04 &nbsp;|&nbsp; ✉️ salahbarki.seb@gmail.com &nbsp;|&nbsp; 📍 Amiens, France
  </div>
</footer>

</body>
</html>
"""

    path = os.path.join(OUTPUT_DIR, "Portfolio_Salah_Barki_v2.html")
    with open(path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    print(f"[OK] HTML créé : {path}")


if __name__ == "__main__":
    create_word()
    create_ppt()
    create_pdf()
    create_html()
    print("\n=== PORTFOLIO V2 COMPLET ===")
