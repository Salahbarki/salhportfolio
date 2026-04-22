# -*- coding: utf-8 -*-
"""
Portfolio PPT V5 — Style Industrial Precision
Bleu industriel #003857 + Orange sécurité #fc7127
"""
import os

OUTPUT_DIR = r"C:\Users\ba2rb\Downloads\salpre"

PROFILE = {
    "name": "Salah Eddine Barki",
    "title": "Technicien de Maintenance Industrielle",
    "contact": "📍 Amiens, France  |  📱 +33 6 88 69 07 04  |  ✉️ salahbarki.seb@gmail.com",
    "summary": """Technicien de maintenance industrielle senior avec plus de 10 ans d'expérience sur des équipements 
automatisés et mécaniques complexes : presses injection, lignes de production agroalimentaire (PET/canettes), 
bancs de test automobile et utilités industrielles. Compétences transverses : électrotechnique, automatisme 
(Siemens S7, TIA Portal), mécanique, hydraulique et pneumatique."""
}

FICHES = [
    {
        "num": "01",
        "titre": "Panne encartonneuse — Ligne PET",
        "entreprise": "Coca-Cola",
        "lieu": "Grigny (91)",
        "type": "CURATIF URGENT",
        "typeColor": "secondary",
        "equip": "Encartonneuse / casseuse de colis — Ligne PET 2L — 60 000 b/h",
        "contexte": "Ligne critique. Production 22h/jour. Arrêt >20 min = risque rupture stock.",
        "symptome": "Défaut HMI #402 toutes les 8–12 min. 15–18 colis rejetés/heure. Opérateurs redémarrage manuel forcé.",
        "diagnostic": "Capteur photo signal oscillant (condensation). Butée jeu 3 mm (spec 0.5 mm). Amortisseur usé. Servo décalé 2–3 mm.",
        "actions": [
            "Capteur anti-buée + câble blindé",
            "Remplacement amortisseur + recalage butée 0.5 mm",
            "Temporisation 50 ms programme TIA Portal",
            "Recalage servo-pousseur",
            "Run 45 min validation : 2 rejets vs 12–15"
        ],
        "resultat": "MTTR 40 min → 8 min (–80%). Rejets ÷5 (17/h → 3/h).",
        "metriques": [("MTTR", "40→8 min"), ("Rejets", "–80%"), ("Dispo.", "+12%")]
    },
    {
        "num": "02",
        "titre": "Mise en service presse injection",
        "entreprise": "TE Connectivity",
        "lieu": "Tanger",
        "type": "INSTALLATION",
        "typeColor": "primary",
        "equip": "Presse 120T + sécheur + trémie + détecteur métal",
        "contexte": "Relocalisation cellule injection. Objectif production J+2.",
        "symptome": "Nouvelle installation. Objectif cycle <20 s, 0 défaut qualité.",
        "diagnostic": "Compatibilité arrivées 400V/120 bar/eau 20°C. Plan constructeur. Niveau huile OK.",
        "actions": [
            "Plots anti-vib + nivellement laser ±0.2 mm",
            "Raccordements élec/hyd/eau + tests étanchéité",
            "Sécheur : point rosée –42°C (spec –40°C)",
            "Détecteur métal : 10/10 détection validée",
            "500 pièces : 0 hors tolérance"
        ],
        "resultat": "Production J+2. Cycle 18.5 s. 0 défaut/500. Checklist 42 pts réutilisée.",
        "metriques": [("Délai", "J+2"), ("Cycle", "18.5 s"), ("Qualité", "100%")]
    },
    {
        "num": "03",
        "titre": "Débogage automate Siemens S7",
        "entreprise": "TE Connectivity",
        "lieu": "Tanger",
        "type": "AUTOMATISME",
        "typeColor": "primary",
        "equip": "Cellule assemblage — S7-1200 + HMI + servo",
        "contexte": "Cellule critique auto. OEE 72% vs objectif 85%.",
        "symptome": "Cycle 28–30 s vs 16–17 s nominal. Alarme 'Timeout bras' 3–4×/h.",
        "diagnostic": "Tempo sécurité 8 s (bras réel 5.5 s). Capteur pince instable (support desserré 2 mm). Code séquentiel bloquant.",
        "actions": [
            "Serrage capteur + écrou nylstop",
            "Tempo 8 s → 6.5 s + antirebond 35 ms",
            "Parallélisation mouvements compatibles",
            "Compteur cycle + alarme dérive auto",
            "HMI diagnostic temps cycle temps réel"
        ],
        "resultat": "Cycle 28.5→16.8 s (–41%). OEE 72%→84%. Alarmes 3.4→<1/h.",
        "metriques": [("Cycle", "–41%"), ("OEE", "+12 pts"), ("Alarmes", "–94%")]
    },
    {
        "num": "04",
        "titre": "RCA banc de test électrique",
        "entreprise": "Kromberg & Schubert",
        "lieu": "Kénitra",
        "type": "RCA & QUALITÉ",
        "typeColor": "secondary",
        "equip": "Banc test élec — Ligne faisceaux 32 voies",
        "contexte": "Client automobile JIT. Tolérance <2% rebut.",
        "symptome": "Rejet 12% vs 1–2% normal. 240+ pièces rebut/jour.",
        "diagnostic": "Pareto : 2/3 pins oxydés (14/32). Hotte filtre saturé 80 vs 120 m³/h. Pas de nettoyage pins >6 mois.",
        "actions": [
            "Retri batch : 94% bonnes / 6% rebut",
            "Remplacement connecteur 32 voies",
            "Réglage pression vérin 3.2→4.0 bar",
            "Remplacement filtre hotte + hotte mobile",
            "Nettoyage pins 2 semaines + QC mensuel"
        ],
        "resultat": "Rejet 12.3%→1.1%. Économie ~220 pièces/jour. QC –60% temps.",
        "metriques": [("Rejet", "12.3→1.1%"), ("Économie", "+220 j"), ("QC", "–60%")]
    },
    {
        "num": "05",
        "titre": "Dépannage hydraulique presse",
        "entreprise": "Sovireso",
        "lieu": "St-Laurent-sur-Sèvre",
        "type": "CURATIF MÉCANIQUE",
        "typeColor": "secondary",
        "equip": "Presse hydraulique — Vérin Ø100/70 — Centrale 80 bar",
        "contexte": "Machine unique site. Pas remplaçante. 480 pièces/jour.",
        "symptome": "Pression 80→30 bar. Coups de bélier. Arrêt depuis 2 h.",
        "diagnostic": "Clapet anti-retour bloqué (morceau joint). Accumulateur précharge 12 vs 65 bar. Filtre retour encrassé. Huile 68°C vs 55°C.",
        "actions": [
            "LOTO + vidange 120 L huile",
            "Dépose/nettoyage clapet + joint neuf",
            "Remplacement accumulateur + précharge 65 bar",
            "Filtres retour/aspiration neufs + huile ISO VG 46",
            "Essais 50 cycles : 80 bar stable ±2"
        ],
        "resultat": "3h20 vs prévision 7h. Production relancée jour même. Pression stable.",
        "metriques": [("Durée", "3h20"), ("Prévision", "7h00"), ("Pression", "80 bar")]
    },
    {
        "num": "06",
        "titre": "Planification GMAO annuelle",
        "entreprise": "TE Connectivity",
        "lieu": "Tanger",
        "type": "PLANIFICATION",
        "typeColor": "primary",
        "equip": "Site 142 équipements (injection, assemblage, test, utilités)",
        "contexte": "Site 450 pers. Objectif : préventif >90%, arrêts –25%.",
        "symptome": "Préventif 68% (233/342). Arrêts imprévus 23%. Coûts +12%. MTBF 420 h.",
        "diagnostic": "45% conflit production. 30% pièces indisponibles. Pas classification ABC.",
        "actions": [
            "Classification ABC (A=14%, B=32%, C=54%)",
            "Planning 12 mois négocié production",
            "Stock sécurité A + délais 5j→2j",
            "Achats groupés sous-traitants –18%",
            "KPI mensuel : MTBF, MTTR, coût/tonne"
        ],
        "resultat": "Préventif 68→94%. Arrêts –28%. MTBF 420→610 h (+45%). Coûts –15%.",
        "metriques": [("Préventif", "68→94%"), ("Arrêts", "–28%"), ("MTBF", "+45%")]
    }
]


def create_ppt():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors Industrial Precision
    C_PRIMARY = RGBColor(0x00, 0x38, 0x57)      # #003857
    C_SECONDARY = RGBColor(0xA2, 0x3F, 0x00)    # #a23f00
    C_SECONDARY_L = RGBColor(0xFC, 0x71, 0x27)  # #fc7127
    C_BG = RGBColor(0xF8, 0xFA, 0xFA)           # #f8fafa
    C_SURFACE = RGBColor(0xF2, 0xF4, 0xF4)     # #f2f4f4
    C_TEXT = RGBColor(0x19, 0x1C, 0x1D)         # #191c1d
    C_MUTED = RGBColor(0x72, 0x78, 0x7F)        # #72787f
    C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    C_ERROR = RGBColor(0xBA, 0x1A, 0x1A)        # #ba1a1a

    def add_bg(slide, color=C_BG):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_rect(slide, left, top, width, height, color):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def add_text(slide, left, top, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT, font_name='Inter'):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        return box

    blank = prs.slide_layouts[6]

    # ===== SLIDE 1: HERO =====
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    # Top primary bar
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(3.8), C_PRIMARY)
    # Secondary accent bar
    add_rect(slide, Inches(0), Inches(3.78), Inches(13.333), Inches(0.08), C_SECONDARY_L)
    
    add_text(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.8), "TECHMAINTENANCE PRO", 16, C_SECONDARY_L, bold=True, align=PP_ALIGN.LEFT)
    add_text(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.2), "Portfolio", 60, C_WHITE, bold=True, align=PP_ALIGN.LEFT)
    add_text(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(0.6), "d'Interventions Maintenance Industrielle", 22, C_WHITE, align=PP_ALIGN.LEFT)
    
    add_text(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.6), PROFILE["name"], 28, C_PRIMARY, bold=True, align=PP_ALIGN.LEFT)
    add_text(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.4), PROFILE["contact"].replace('📍','').replace('📱','').replace('✉️',''), 11, C_MUTED, align=PP_ALIGN.LEFT)
    
    # Stats pills
    stats = [("10+", "Années"), ("7", "Entreprises"), ("4", "Pays"), ("6", "Fiches")]
    x = 0.8
    for num, lab in stats:
        add_rect(slide, Inches(x), Inches(5.6), Inches(2.5), Inches(1.2), C_SURFACE)
        add_rect(slide, Inches(x), Inches(5.6), Inches(2.5), Inches(0.06), C_SECONDARY_L)
        add_text(slide, Inches(x+0.1), Inches(5.75), Inches(2.3), Inches(0.6), num, 28, C_SECONDARY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x+0.1), Inches(6.35), Inches(2.3), Inches(0.4), lab, 11, C_TEXT, align=PP_ALIGN.CENTER)
        x += 2.8

    # ===== SLIDE 2: COMPETENCES =====
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), C_PRIMARY)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), "COMPÉTENCES & TECHNOLOGIES", 22, C_WHITE, bold=True)
    
    comps = [
        ("ÉLECTROTECHNIQUE", "Siemens S7\nTIA Portal / Step 7\nCâblage BT, schémas IEC\nVariateurs, capteurs", C_SECONDARY),
        ("MÉCANIQUE", "Presses injection\nLignes PET / Canettes\nHydraulique / Pneumatique\nTransmissions, guidages", C_PRIMARY),
        ("MÉTHODES", "GMAO SAP PM\nRCA, 5 pourquoi\nPareto, SPC\nKPI MTBF, MTTR, OEE", C_PRIMARY),
    ]
    x = 0.5
    for title, body, color in comps:
        add_rect(slide, Inches(x), Inches(1.3), Inches(4.0), Inches(5.5), C_SURFACE)
        add_rect(slide, Inches(x), Inches(1.3), Inches(4.0), Inches(0.06), color)
        add_text(slide, Inches(x+0.15), Inches(1.5), Inches(3.7), Inches(0.5), title, 14, color, bold=True)
        add_text(slide, Inches(x+0.15), Inches(2.1), Inches(3.7), Inches(4.5), body, 13, C_TEXT)
        x += 4.3

    # ===== FICHES (1 slide each) =====
    for f in FICHES:
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        # Header
        add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), C_PRIMARY)
        top_color = C_SECONDARY_L if f["typeColor"] == "secondary" else C_PRIMARY
        add_rect(slide, Inches(0), Inches(0.96), Inches(13.333), Inches(0.06), C_SECONDARY_L if f["typeColor"] == "secondary" else C_SECONDARY_L)
        
        add_text(slide, Inches(0.4), Inches(0.12), Inches(12), Inches(0.5), f"{f['num']}. {f['titre']}", 20, C_WHITE, bold=True)
        add_text(slide, Inches(0.4), Inches(0.55), Inches(12), Inches(0.4), f"{f['entreprise']} — {f['lieu']}  |  {f['type']}", 10, C_MUTED)
        
        # Left: Diagnostic + Context
        add_rect(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(5.8), C_SURFACE)
        add_rect(slide, Inches(0.3), Inches(1.2), Inches(0.06), Inches(5.8), C_SECONDARY_L if f["typeColor"] == "secondary" else C_SECONDARY_L)
        
        y_off = 1.35
        add_text(slide, Inches(0.5), Inches(y_off), Inches(5.9), Inches(0.35), "CONTEXTE", 11, C_SECONDARY if f["typeColor"] == "secondary" else C_PRIMARY, bold=True)
        y_off += 0.3
        add_text(slide, Inches(0.5), Inches(y_off), Inches(5.9), Inches(0.6), f['contexte'], 10, C_TEXT)
        y_off += 0.55
        add_text(slide, Inches(0.5), Inches(y_off), Inches(5.9), Inches(0.35), "SYMPTÔME", 11, C_SECONDARY if f["typeColor"] == "secondary" else C_PRIMARY, bold=True)
        y_off += 0.3
        add_text(slide, Inches(0.5), Inches(y_off), Inches(5.9), Inches(0.8), f['symptome'], 10, C_TEXT)
        y_off += 0.75
        add_text(slide, Inches(0.5), Inches(y_off), Inches(5.9), Inches(0.35), "DIAGNOSTIC", 11, C_SECONDARY if f["typeColor"] == "secondary" else C_PRIMARY, bold=True)
        y_off += 0.3
        add_text(slide, Inches(0.5), Inches(y_off), Inches(5.9), Inches(2.0), f['diagnostic'], 10, C_TEXT)
        
        # Right: Actions
        add_rect(slide, Inches(6.7), Inches(1.2), Inches(6.3), Inches(4.5), C_WHITE)
        add_text(slide, Inches(6.9), Inches(1.35), Inches(6), Inches(0.4), "ACTIONS", 11, C_PRIMARY, bold=True)
        y_off = 1.75
        for i, action in enumerate(f['actions']):
            add_text(slide, Inches(6.9), Inches(y_off), Inches(6), Inches(0.45), f"• {action}", 10, C_TEXT)
            y_off += 0.42
        
        # Result bar
        add_rect(slide, Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.9), C_SECONDARY_L if f["typeColor"] == "secondary" else C_PRIMARY)
        add_text(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7), f"RÉSULTAT : {f['resultat']}", 11, C_WHITE, bold=True)
    
    # ===== FINAL: METHODOLOGY =====
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), C_PRIMARY)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), "MÉTHODOLOGIE D'INTERVENTION", 22, C_WHITE, bold=True)
    
    meths = [
        ("1", "SÉCURISATION", "Consignation LOTO, EPI, signalisation"),
        ("2", "DIAGNOSTIC", "Mesures, alarmes HMI, schémas, GMAO"),
        ("3", "RCA", "5 pourquoi, Ishikawa, Pareto"),
        ("4", "INTERVENTION", "Réparation, réglage, programmation"),
        ("5", "VALIDATION", "Essais charge, mesures A/A, GMAO, formation"),
    ]
    x = 0.4
    for num, title, desc in meths:
        add_rect(slide, Inches(x), Inches(1.4), Inches(2.4), Inches(4.5), C_SURFACE)
        add_rect(slide, Inches(x), Inches(1.4), Inches(2.4), Inches(0.06), C_SECONDARY_L)
        add_text(slide, Inches(x+0.1), Inches(1.55), Inches(2.2), Inches(0.5), num, 32, C_SECONDARY_L, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x+0.1), Inches(2.2), Inches(2.2), Inches(0.4), title, 12, C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x+0.1), Inches(2.6), Inches(2.2), Inches(2.5), desc, 11, C_TEXT, align=PP_ALIGN.CENTER)
        x += 2.55
    
    add_text(slide, Inches(1), Inches(6.2), Inches(11.333), Inches(0.5), f"{PROFILE['name']} — {PROFILE['contact']}", 11, C_MUTED, align=PP_ALIGN.CENTER)

    path = os.path.join(OUTPUT_DIR, "Portfolio_Salah_Barki_v5.pptx")
    prs.save(path)
    print(f"[OK] PPT v5 Industrial Precision : {path}")


if __name__ == "__main__":
    create_ppt()
