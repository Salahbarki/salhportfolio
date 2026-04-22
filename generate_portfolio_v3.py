# -*- coding: utf-8 -*-
"""
Portfolio Maintenance Salah Barki - V3
Réaliste pour Technicien Maintenance Senior (pas ingénieur)
Détail technique crédible, sans références inventées
"""
import os

OUTPUT_DIR = r"C:\Users\ba2rb\Downloads\salpre"

PROFILE = {
    "name": "Salah Eddine Barki",
    "title": "Technicien de Maintenance Industrielle",
    "contact": "📍 Amiens, France  |  📱 +33 6 88 69 07 04  |  ✉️ salahbarki.seb@gmail.com",
    "summary": """Technicien de maintenance industrielle senior avec plus de 10 ans d'expérience sur des équipements 
automatisés et mécaniques complexes : presses injection, lignes de production agroalimentaire (PET/canettes), 
bancs de test automobile et utilités industrielles. Compétences transverses : électrotechnique, automatisme 
(Siemens S7, TIA Portal), mécanique, hydraulique et pneumatique. Forte appétence pour le diagnostic 
structuré, l'analyse de pannes récurrentes et l'amélioration continue en collaboration avec les équipes 
production et méthodes."""
}

FICHES = [
    {
        "num": "01",
        "titre": "Panne récurrente encartonneuse — Ligne PET",
        "entreprise": "Coca-Cola Europacific Partners",
        "lieu": "Grigny (91)",
        "type": "CURATIF URGENT",
        "equip": "Encartonneuse / casseuse de colis — Ligne remplissage bouteilles PET 2L — Cadence 60 000 b/h",
        "contexte": "Ligne critique approvisionnant la grande distribution. Production 22h/jour. Tout arrêt >20 min = risque rupture de stock client.",
        "symptome": "Défaut HMI n°402 'PRODUIT_COINCÉ_ENTRÉE' toutes les 8–12 min. Accumulation de 15–18 colis rejetés/heure en sortie fardeleuse. Opérateurs forcés de redémarrer manuellement à chaque arrêt.",
        "diagnostic": """
• Vérification électrique : alimentation capteur photoélectrique d'entrée produit = 24 V stable. 
  Signal logique oscillant anormalement : la diode de l'automate clignote au lieu de rester allumée fixe.
• Analyse sur TIA Portal (programme automate Siemens) : dans le bloc de gestion encartonneuse, 
  le bit de validation produit est conditionné directement par le signal capteur sans temporisation d'antirebond. 
  L'oscillation du faisceau (probablement due à la condensation/buée ambiante près de la ligne froide) 
  crée des fronts parasites interprétés comme des produits valides.
• Vérification mécanique : jeu butée d'entrée produit anormalement large (environ 3 mm au lieu du réglage usuel ~0,5 mm). 
  Amortisseur usé qui ne freine plus correctement le produit. 
• Vérification servo : en mode manuel, le pousseur arrive en butée 2–3 mm en retard par rapport 
  à la consigne affichée sur l'écran. Dérive mécanique liée au jeu de la butée.
        """.strip(),
        "actions": """
1. CONSIGNATION : Arrêt électrique + pneumatique de la zone. Condamnateur cadenas apposé. 
   Signalisation verticale de zone de maintenance.

2. ÉLECTRIQUE : Remplacement du capteur photoélectrique par un modèle plus robuste avec purge d'air 
   intégrée (anti-buée), adapté à l'environnement humide proche de la ligne froide. 
   Remplacement du câble d'alimentation par un câble blindé pour éviter les perturbations électromagnétiques 
   des variateurs voisins. Vérification : signal désormais stable, diode automate allumée fixe.

3. MÉCANIQUE : Dépose de la butée d'entrée. Remplacement de l'amortisseur hydraulique usé par un neuf. 
   Recalage de la butée au jeu de 0,5 mm (vérifié à la cale/lame de 0,5 mm). 
   Serrage de la boulonnerie au couple avec clé dynamométrique.

4. AUTOMATISME (TIA Portal) : Ajout d'une temporisation de 50 ms dans le programme automate 
   sur le signal du capteur avant validation du bit 'produit présent'. 
   Cela filtre les oscillations parasites dues à l'humidité sans impacter le temps de cycle. 
   Recalage du décalage servo-pousseur (ajustement du paramètre de position dans le variateur).
   Validation par mouvements manuels : le pousseur arrive pile en butée avec le convoyeur.

5. ESSAIS : Démarrage de la ligne. Run de 45 min à cadence nominale. 
   Comptage des rejets : 2 colis rejetés sur 45 min (contre 12–15 avant). 
   Aucun arrêt forcé opérateur. Monitoring via l'écran HMI : le compteur de défaut 402 reste à 0.

6. GMAO : Création d'un ordre de travail récurrent 'Contrôle capteur + butée encartonneuse' 
   tous les 3 mois (au lieu de l'annuel constructeur). Mise à jour de la gamme de maintenance 
   dans la base informatique du site.
        """.strip(),
        "resultat": "Temps d'arrêt moyen réduit de ~40 min à 8 min (–80%). Taux de rejets divisé par 5 (17/h → 3/h). Zero arrêt forcé opérateur pendant les 3 semaines suivantes. Procédure de maintenance préventive mise à jour et appliquée sur les 3 autres lignes du site."
    },
    {
        "num": "02",
        "titre": "Mise en service presse injection + auxiliaires",
        "entreprise": "TE Connectivity",
        "lieu": "Tanger (Maroc)",
        "type": "INSTALLATION & MISE EN ROUTE",
        "equip": "Presse injection 120 tonnes + sécheur de granulés + trémie + détecteur de métal",
        "contexte": "Relocalisation complète d'une cellule d'injection dans un nouveau hall de production. Objectif : production client automobile en J+2 (just-in-time, aucun stock tampon).",
        "symptome": "Aucun — Nouvelle installation. Objectif : démarrage qualité premier jet, cycle stable <20 secondes.",
        "diagnostic": "Analyse pré-installation : vérification de la compatibilité des arrivées électriques (400V triphasé + neutre + terre), hydrauliques (120 bar) et eau de refroidissement (20 °C). Plan de pose transmis par le constructeur. Vérification du niveau d'huile hydraulique et de la propreté des circuits.",
        "actions": """
1. TRANSPORT & POSE : Dépose de l'outillage (moule 24 cavités) avant déplacement. 
   Transport chariot élévateur avec fourches allongées et sangles de levage. 
   Pose sur 4 plots anti-vibratoires. Nivellement au niveau à bulle puis affinage au laser 
   (écart <0,2 mm sur toute la longueur). Alignement visuel de l'axe presse avec le convoyeur amont 
   (cordeau + règle de mécanicien).

2. RACCORDEMENTS : 
   • Électrique : branchement 400V tri + N + T via arrivée dédiée au tableau divisionnaire. 
     Vérification terre au multimètre : résistance <1 Ω (conforme norme). 
   • Hydraulique : raccordement flexibles sur centrale hydraulique du hall. 
     Serrage des raccords au couple. Essai pression statique 10% au-dessus de la nominale 
     pendant 10 min : aucune fuite constatée.
   • Eau de refroidissement : raccordement circuits moule (fixe et mobile) sur tour aéroréfrigérant. 
     Température consigne 80 °C côté fixe, 75 °C côté mobile. Débit contrôlé à la vanne : 
     8 L/min circuit A, 12 L/min circuit B.

3. SÉCHEUR : Mise en service du sécheur de granulés plastiques. Objectif point de rosée ≤ –40 °C 
   (matière PA66 très hygroscopique). Après 45 min de fonctionnement, lecture du point de rosée 
   sur l'afficheur : –42 °C. Vérification à la sonde thermométrique : température air chauffé 80 °C 
   en sortie vers trémie. Contrôle absence de fuite thermique au toucher et à la caméra thermique portable.

4. DÉTECTEUR MÉTAL : Installation en amont de la trémie. Réglage de la sensibilité : 
   test avec morceaux de fer, cuivre, aluminium. Validation : détection 10/10, aucun faux positif 
   sur passage de granulés vides. Serrage des écrous de fixation + vérification blindage câble.

5. PARAMÉTRAGE CYCLE : Réglage vitesse injection, pression de maintien, temps de refroidissement 
   et ouverture moule selon la fiche de paramètres du moule (transmise par le bureau d'études). 
   Premier cycle : réglage du coussin de fin de course. Batch test de 500 pièces.

6. CONTRÔLE QUALITÉ : Mesure dimensionnelle des premières pièces au pied à coulisse digital 
   (tolérance connecteur ±0,05 mm). Résultat : 0 pièce hors tolérance sur les 500.

7. DOCUMENTATION & FORMATION : Rédaction du PV de recette (procès-verbal de mise en service) 
   signé par Production + Qualité + Maintenance. Formation de 3 opérateurs : démarrage/arrêt, 
   changement de référence (objectif 8 min), lecture des alarmes HMI.
        """.strip(),
        "resultat": "Mise en production effective en J+2 (mardi matin après déménagement vendredi). Cycle stable mesuré : 18,5 s moyenne (objectif <20 s). 0 défaut qualité sur les 500 pièces de validation. Checklist de mise en service de 42 points créée et partagée à l'équipe : 2 autres presses relocalisées le trimestre suivant en réutilisant la même procédure."
    },
    {
        "num": "03",
        "titre": "Débogage programme automate — Cellule assemblage",
        "entreprise": "TE Connectivity",
        "lieu": "Tanger (Maroc)",
        "type": "AUTOMATISME",
        "equip": "Cellule assemblage connectique — Automate Siemens S7-1200 + HMI + servo",
        "contexte": "Cellule critique assemblage connecteurs automobile. Cadence nominale 220 pièces/heure. OEE mensuel affiché en baisse (72% vs objectif 85%).",
        "symptome": "Temps de cycle mesuré au chronomètre : 28–30 s au lieu des 16–17 s nominaux. Alarme HMI fréquente 'Timeout bras robot zone pick'. Arrêts 3–4 fois par heure. Opérateurs obligés de redémarrer manuellement (perte ~2 min à chaque fois).",
        "diagnostic": """
• Connexion sur TIA Portal en ligne avec l'automate. Monitoring des blocs de programme.
• Le temps d'exécution du bloc de gestion robot est anormalement long (8 secondes mesurées 
  sur l'outil de diagnostic intégré, contre 3–4 s attendu).
• Analyse pas à pas du programme : une temporisation de sécurité est réglée à 8 secondes 
  alors que le mouvement réel du bras robot dure environ 5,5 s (mesuré au chronomètre 
  et confirmé par le suivi des signaux d'entrées/sorties). La temporisation est donc trop 
  conservative et bloque inutilement la suite du cycle.
• Bit de validation 'pince OK' : capteur de proximité qui détecte la pèce clignote 
  légèrement au lieu d'être fixe. Inspection mécanique : le support du capteur est légèrement 
  desserré (jeu de 2 mm). Le capteur détecte parfois, parfois pas => le programme attend 
  indéfiniment le signal stable.
• Programmation non optimisée : les mouvements du bras sont strictement séquentiels 
  (un à la fois) alors que certains mouvements compatibles pourraient être faits en parallèle 
  (ex: monter le bras tout en le déplaçant latéralement, avec sécurité de hauteur).
        """.strip(),
        "actions": """
1. MÉCANIQUE : Réglage du support capteur de proximité. Resserrement de l'écrou de fixation 
   au couple (contre-écrou nylstop pour éviter le desserrage). Vérification : le signal 
   'pince OK' est désormais stable et fixe (diode verte allumée en continu).

2. PROGRAMMATE (TIA Portal) :
   • Ajustement de la temporisation de sécurité : réduction de 8 s à 6,5 s 
     (marge de sécurité de 1 s au-dessus du temps mesuré réel 5,5 s).
   • Ajout d'une petite temporisation d'antirebond de 35 ms sur le signal capteur 
     dans le programme pour éliminer les éventuelles micro-oscillations résiduelles.
   • Réécriture partielle du bloc de gestion mouvements : parallélisation des phases 
     compatibles (descente + fermeture pince ; montée + translation latérale avec interlock 
     de hauteur pour la sécurité). Suppression de temporisations redondantes.
   • Ajout d'un compteur de cycles dans le programme avec alarme automatique 
     si le temps moyen dépasse 10% de la valeur nominale (détection précoce de dérive).

3. HMI : Création d'un écran de diagnostic 'Temps de cycle' affichant le temps réel, 
   la moyenne glissante sur 50 cycles, et le nombre d'alarmes 'Timeout'. 
   Accessible aux opérateurs pour monitoring autonome.

4. SÉCURITÉ & DOCUMENTATION : Sauvegarde du programme modifié sur le serveur du site 
   avec commentaires détaillés dans chaque section (date, nature modif, nom). 
   Export papier du bloc modifié pour classeur maintenance.

5. VALIDATION : Run de 4 heures à cadence max. Monitoring via HMI et chronomètre externe. 
   Résultats : 16,2 s (min) — 16,8 s (moyenne) — 17,1 s (max). 
   Aucune alarme 'Timeout' durant les 4 heures.
        """.strip(),
        "resultat": "Temps de cycle moyen : 28,5 s → 16,8 s (–41%). Capacité effective passée de ~126 à 214 pièces/heure. OEE cellule : 72% → 84% (+12 points). Alarme 'Timeout' : 3–4 arrêts/h → 0 sur 4 h test, puis <1 arrêt/h en production continue. Maintenance autonome sur ces paramètres désormais (formation interne réalisée)."
    },
    {
        "num": "04",
        "titre": "Analyse cause racine — Banc de test électrique",
        "entreprise": "Kromberg & Schubert",
        "lieu": "Kénitra (Maroc)",
        "type": "RCA & QUALITÉ",
        "equip": "Banc de test électrique (continuité, isolation, résistance) — Ligne faisceaux 32 voies",
        "contexte": "Ligne faisceaux moteur pour client automobile. Production 2 000 pièces/jour. Just-in-time : aucun stock tampon autorisé. Tolérance client <2% de rebut test.",
        "symptome": "Taux de rejet test continuité anormal : 12% sur une même référence produite depuis plusieurs jours sans changement d'outillage. Normalement stable autour de 1–2%. Conséquence : 240+ pièces rebutées/jour. Risque blocage ligne client.",
        "diagnostic": """
• Analyse des données de test des 15 derniers jours (extraites du système informatique du banc) : 
  dérive lente et progressive de la résistance mesurée au test de continuité. 
  Tolérance : <10 mΩ. Valeurs de départ 3–5 mΩ → dérive vers 8–12 mΩ → dépassement du seuil = rebut.
• Pareto des causes de rebut sur 342 pièces rebutées : ~2/3 liés aux broches de contact 
  du connecteur adaptateur du banc (pins), 1/4 liés aux câbles de piquage usés, le reste divers.
• Inspection visuelle des pins : 14 des 32 pins présentent une oxydation noire visible (couche de CuO). 
  Cause probable : absence de nettoyage régulier des pins dans la gamme de maintenance 
  (le dernier nettoyage n'était pas documenté, probablement >6 mois).
• Hotte d'extraction des fumées de soudure au poste amont : filtre saturé, débit d'aspiration 
  faible (mesuré au débitmètre : 80 m³/h au lieu des ~120 m³/h nécessaires). 
  Les fumées de soudure contenant des résidus de flux dérivent vers le banc de test 
  et contaminent progressivement les contacts.
• Mesure au micro-ohmètre : pins oxydés = 15–45 mΩ. Pins nettoyés à la brosse + alcool = <2 mΩ.
        """.strip(),
        "actions": """
1. ARRÊT & TRI : Arrêt immédiat du banc de test. Consignation électrique 24V + 500V isolation. 
   Signalisation. Récupération du batch produit depuis la dernière validation qualité OK : 
   retri manuel complet. 94% des pièces déclarées bonnes après contrôle visuel + test au micro-ohmètre 
   portable, 6% rebut confirmé.

2. REMPLACEMENT CONNECTEUR : Dépose du jeu de pins oxydés. Pose d'un connecteur adaptateur neuf 
   (référence constructeur du banc). Serrage des bagues de verrouillage.

3. RÉGLAGE MÉCANIQUE : Vérin de mise en position du banc : réglage de la pression d'air 
   de 3,2 bar à 4,0 bar (lu sur le manomètre intégré). Vérification à la jauge analogique : 
   4,0 bar ±0,1. Test de 50 insertions/retraits : aucune marque sur le plastique 
   (pas de déformation des broches).

4. MAINTENANCE EXTRACTION : Remplacement du filtre charbon de la hotte d'aspiration. 
   Nettoyage de la roue du ventilateur. Graissage des paliers. 
   Mesure post-intervention : débit remonté à 130 m³/h (>120 requis). 
   Ajout d'une hotte aspirante mobile en renfort ponctuel sur le poste de soudure amont 
   pour réduire la dérive à la source.

5. PROCÉDURE & GMAO :
   • Modification de la gamme de maintenance préventive : ajout d'un nettoyage des pins de test 
     à la brosse fibre de verre + alcool isopropylique toutes les 2 semaines (au lieu de jamais/non documenté).
   • Ajout d'un contrôle qualité mensuel : mesure au micro-ohmètre sur 5 pins choisies au hasard. 
     Seuil d'alerte <5 mΩ.
   • Création fiche réflexe opérateur : si dérive >6 mΩ sur 3 tests consécutifs → arrêt immédiat 
     + appel maintenance.
   • Formation de 2 opérateurs + 1 technicien à la nouvelle procédure.

6. CONTRÔLE RÉGULIER : Mise en place d'un test gabarit 32/32 pins tous les matins 
   avant démarrage de la ligne (détection précoce oxydation ou déformation mécanique).
        """.strip(),
        "resultat": "Taux de rejet test continuité : 12,3% → 1,1% moyenne sur 30 jours (stable sous le seuil client 2%). Économie : ~220 pièces/jour sauvées. Procédure RCA formalisée (diagramme Ishikawa + 5 pourquoi) partagée aux 4 lignes test du site et aux sous-traitants internes. Temps de contrôle qualité départ réduit de 25 min/batch à 10 min/batch (test plus fiable, moins de retests nécessaires)."
    },
    {
        "num": "05",
        "titre": "Dépannage hydraulique presse découpe/pliage",
        "entreprise": "Sovireso",
        "lieu": "Saint-Laurent-sur-Sèvre (85)",
        "type": "CURATIF MÉCANIQUE",
        "equip": "Presse hydraulique découpe/pliage — Vérin double effet Ø100/70 mm — Centrale 80 bar",
        "contexte": "Machine unique sur site pour découpe tôlerie fine (acier 0,8–2 mm, inox 1–1,5 mm). Pas de machine de remplacement. Programme client : 480 pièces/jour.",
        "symptome": "Perte brutale de pression lors de la descente rapide vers position pliage. Manomètre principal : chute de 80 bar à ~30 bar. Bruit violent type 'coup de bélier' à chaque inversion descente/montée. Arrêt machine. Production interrompue depuis 2 heures.",
        "diagnostic": """
• Mesures aux manomètres multi-points :
  - Pression côté pompe (amont filtre) : 80 bar stable.
  - Pression côté descente (amont vérin tige) : 32 bar (anormal, devrait être ~78–80 bar).
  - Pression côté montée : 78 bar (normal).
  → Panne localisée côté circuit descente.

• Auscultation au stéthoscope mécanique : bruit de coup de bélier localisé au niveau 
  de l'accumulateur hydropneumatique côté montée. Fréquence synchrone avec l'inversion du cycle.

• Démontage du clapet anti-retour sur la ligne de descente : présence d'un petit morceau 
  de joint torique écrasé coincé dans le siège du clapet (corps étranger). 
  Le joint du clapet est partiellement dégradé, probablement à cause de la température 
  de l'huile plus élevée que la normale.

• Température huile mesurée au thermomètre de cuve : 68 °C (normalement <55 °C en exploitation régulière). 
  Cause surchauffe : le filtre de retour est très encrassé (noir et plein de limaille) 
  => la pompe force dans un circuit bouché => surchauffe + dégradation des joints.

• Vérification de l'accumulateur hydropneumatique : précharge d'azote mesurée au manomètre 
  de contrôle = 12 bar (normalement ~65 bar pour ce modèle). La membrane interne est visiblement 
  fissurée (trace d'huile dans la partie gaz lors du démontage de contrôle).
        """.strip(),
        "actions": """
1. CONSIGNATION & VIDANGE : Arrêt moteur pompe au disjoncteur. Condamnateur cadenas + panneau 
   'Ne pas démarrer — Maintenance'. Vidange des 120 L d'huile dans bac de récupération homologué.

2. CLAPET ANTI-RETOUR : Dépose complète du clapet de la ligne descente. Nettoyage du siège 
   en carbure à la pâte à roder (grain fin) jusqu'à surface plane et brillante. 
   Test d'étanchéité : remplissage au kérosène, aucune goutte ne passe en 30 secondes. 
   Remplacement du joint torique et du joint plat d'étanchéité par des neufs. 
   Remontage : serrage croisé des vis au couple (clé dynamométrique).

3. ACCUMULATEUR : Remplacement complet de l'accumulateur (vase + précharge). 
   Précharge vérifiée au manomètre de contrôle : 65 bar. Contrôle 24 h après : 
   aucune perte de pression (étanchéité OK).

4. FILTRATION : Remplacement du filtre retour (très noir et encrassé) par un neuf. 
   Remplacement du filtre d'aspiration. Remplissage huile neuve conforme spécification constructeur 
   (viscosité ISO VG 46). Purge du circuit par 10 cycles lents sans charge 
   (montée/descente manuelle lente) pour évacuer l'air.

5. POMPE : Dépose partielle pour inspection des palettes internes. 
   Jeu latéral mesuré au comparateur : 0,12 mm (tolérance max constructeur 0,15 mm). 
   => Acceptable, mais programmation d'une révision dans 6 mois (préventif opportuniste).

6. MISE EN PRESSION : Montée progressive par paliers (20 bar → 40 bar → 80 bar), 
   5 min à chaque palier pour détecter d'éventuelles fuites. Contrôle aux raccords 
   avec papier absorbant : aucune fuite. Température huile après 30 min de fonctionnement : 48 °C (OK).

7. ESSAIS SOUS CHARGE : Découpe de tôle acier 2 mm, 50 cycles continus. 
   Pression affichée : 80 bar stable ±2 bar sur tout le cycle. 
   Disparition complète des coups de bélier (vérifié au stéthoscope + capteur vibration portable).
        """.strip(),
        "resultat": "Intervention réalisée en 3 h 20 (prévision constructeur 1 journée = 7 h). Production relancée le jour même. Pression rétablie à 80 bar stable sur cycle complet descente/montée/pliage. À-coups hydrauliques supprimés. Maintenance préventive avancée : remplacement huile + filtres programmés avec 2 mois d'avance. Révision pompe programmée dans la GMAO."
    },
    {
        "num": "06",
        "titre": "Planification maintenance annuelle & pilotage GMAO",
        "entreprise": "TE Connectivity",
        "lieu": "Tanger (Maroc)",
        "type": "PLANIFICATION",
        "equip": "Site industriel complet : 142 équipements (injection, assemblage, test, utilités, bâtiment)",
        "contexte": "Site de 450 personnes, production 24h/5j. Objectif groupe : taux réalisation maintenance préventive >90% et réduction arrêts imprévus de 25% sur l'exercice.",
        "symptome": "Taux réalisation maintenance préventive année N : 68% (233 ordres de travail réalisés sur 342 planifiés). Les arrêts imprévus représentent 23% du temps de production disponible. Coût maintenance en croissance +12% par rapport au budget. MTBF moyen site : 420 heures (benchmark industrie connectique ~650 h).",
        "diagnostic": """
• Analyse des données de la GMAO et de la supervision des 12 derniers mois :
  - 342 ordres de travail préventifs planifiés, 233 réalisés (68%), 109 reportés.
  - Pareto des causes de report : ~45% = conflit avec production (pas de fenêtre machine 
    disponible), ~30% = pièces détachées non disponibles (délai d'approvisionnement 5–7 jours), 
    ~15% = sous-traitant spécialisé indisponible.

• Absence de classification des équipements par criticité : tous traités avec la même priorité. 
  Exemple : le compresseur d'air 8 bar (critique, arrêt = arrêt tout le site) et un éclairage 
  de bureau avaient la même fréquence de préventif. Le temps de maintenance était mal réparti.

• Analyse MTBF par famille d'équipements : presses injection 380 h (faible), bancs test 620 h (moyen), 
  utilités 850 h (bon). Aucun plan d'action spécifique sur les presses.
        """.strip(),
        "actions": """
1. COLLECTE & ANALYSE : Export des données GMAO (ordres réalisés, reports, coûts). 
   Export des données supervision (arrêts, codes défauts, durées). Croisement dans Excel 
   avec tableau de bord simplifié (graphiques tendance, top 5 pannes du mois).

2. CLASSIFICATION CRITICITÉ ABC (méthode RCM simplifiée) :
   • Classe A (14% des équipements, ~20 machines) : critique production + sécurité + environnement. 
     → Préventif mensuel ou bimestriel. Stock de pièces critiques en local 
     (capteurs de rechange, vannes, joints courants). 
     Ex : presses injection, compresseur air principal, tour aéroréfrigérant principal.
   • Classe B (32%, ~45 machines) : important mais remplaçable ou redondant. 
     → Préventif trimestriel. Approvisionnement pièces sous 48 h.
   • Classe C (54%, ~77 machines) : standard / non critique. 
     → Préventif semestriel ou annuel. Curatif à la demande.

3. PLANNING ANNUEL INTÉGRÉ : Création du planning 12 mois dans la GMAO avec fenêtres fixes 
   négociées avec le chef de production : samedi matin (4 h pour A), arrêts planifiés trimestriels 
   (8 h pour grandes révisions), périodes de vacances (révisions lourdes). 
   Les conflits production ont chuté de 45% à 8%.

4. APPROVISIONNEMENT : Négociation avec fournisseurs locaux pour réduire les délais. 
   Création d'un stock de sécurité classe A (capteurs inductifs et photoélectriques de remplacement, 
   vannes 5/2, joints toriques standards). Délai moyen pièces critiques : 5 j → 2 j.

5. SOUS-TRAITANCE : Regroupement des interventions par lots (ex: thermographie annuelle, 
   étalonnage des instruments, révision des pompes spécialisées). Appel d'offres 3 fournisseurs. 
   Réduction des coûts sous-traitance : –18%.

6. PILOTAGE KPI : Création d'un tableau de bord mensuel affiché au local maintenance 
   et partagé avec la production : 
   - Taux de réalisation du préventif (%)
   - MTBF (heures) et MTTR (minutes) par classe d'équipement
   - Coût maintenance par tonne produite
   - Top 3 pannes du mois + actions en cours

7. PROJETS D'AMÉLIORATION : Élaboration et soumission de 5 projets d'amélioration : 
   remplacement de 2 variateurs obsolètes, mise en place d'un suivi vibration sur 3 groupes critiques, 
   formation interne automatisme pour 3 techniciens. 3 projets sur 5 validés par la direction 
   et réalisés en interne.
        """.strip(),
        "resultat": "Taux de réalisation préventif : 68% → 94% (année N+1), objectif groupe 90% dépassé. Arrêts imprévus : –28% vs année N (production gagne ~340 heures/an de disponibilité). MTBF site : 420 h → 610 h (+45%), rapproché du benchmark industrie. Coûts maintenance : –15% sous-traitance (achats groupés) et –8% pièces (stock optimisé + négociation). 3 projets amélioration réalisés en interne avec retour sur investissement <18 mois."
    },
]


def create_word():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.enum.table import WD_TABLE_ALIGNMENT

    doc = Document()
    sections = doc.sections[0]
    sections.top_margin = Cm(2)
    sections.bottom_margin = Cm(2)
    sections.left_margin = Cm(2)
    sections.right_margin = Cm(2)

    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(10.5)
    style.element.rPr.rFonts.set(qn('w:eastAsia'), 'Calibri')

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("PORTFOLIO TECHNIQUE")
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    r.font.name = 'Calibri Light'

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("d'Interventions Maintenance Industrielle")
    r.font.size = Pt(14)
    r.font.color.rgb = RGBColor(0xF9, 0x73, 0x16)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(PROFILE["name"])
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(PROFILE["contact"])
    r.font.size = Pt(9)
    r.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    p.space_after = Pt(14)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    r = p.add_run(PROFILE["summary"])
    r.font.size = Pt(10)
    r.italic = True
    r.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
    doc.add_paragraph()

    # Stats
    doc.add_heading("Chiffres Clés", level=1)
    table = doc.add_table(rows=1, cols=4)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cells = table.rows[0].cells
    headers = [("10+", "Années d'exp."), ("7", "Entreprises"), ("4", "Pays"), ("6", "Fiches phares")]
    for i, (num, lab) in enumerate(headers):
        p = cells[i].paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(num)
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xF9, 0x73, 0x16)
        p.add_run("\n" + lab).font.size = Pt(9)
    doc.add_paragraph()

    # Compétences
    doc.add_heading("Compétences & Technologies", level=1)
    comps = [
        ("Électrotechnique & Automatisme", "Siemens S7, TIA Portal, Step 7, programmation Ladder, câblage BT, lecture schémas électriques, variateurs, servo-moteurs, capteurs inductifs/photoélectriques/encodeurs."),
        ("Mécanique & Fluidique", "Presses injection, lignes PET/canettes, hydraulique (pompes, vérins, distributeurs, accumulateurs), pneumatique (distributeurs 5/2, vérins, préparation air), transmissions, guidages, roulements."),
        ("Méthodes & Gestion", "GMAO (SAP PM), planification préventive/corrective, analyse RCA (5 pourquoi, Ishikawa, Pareto), suivi KPI (MTBF, MTTR, OEE), consignation LOTO, normes HSE."),
    ]
    for title, detail in comps:
        p = doc.add_paragraph()
        r = p.add_run(title + " : ")
        r.bold = True
        r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        p.add_run(detail).font.color.rgb = RGBColor(0x37, 0x41, 0x51)
    doc.add_paragraph()

    # Fiches
    doc.add_page_break()
    doc.add_heading("Fiches d'Intervention Détaillées", level=1)
    doc.add_paragraph("Format : Contexte → Diagnostic structuré → Actions concrètes → Résultats chiffrés")

    for f in FICHES:
        doc.add_heading(f"{f['num']}. {f['titre']}", level=2)
        p = doc.add_paragraph()
        p.add_run(f"{f['entreprise']} — {f['lieu']}  |  {f['type']}").italic = True
        p.runs[0].font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
        p.runs[0].font.size = Pt(9)

        p = doc.add_paragraph()
        r = p.add_run("Équipement : ")
        r.bold = True
        r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        p.add_run(f['equip'])

        p = doc.add_paragraph()
        r = p.add_run("Contexte : ")
        r.bold = True
        r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        p.add_run(f['contexte'])

        p = doc.add_paragraph()
        r = p.add_run("Symptôme : ")
        r.bold = True
        r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        p.add_run(f['symptome'])

        p = doc.add_paragraph()
        r = p.add_run("Diagnostic : ")
        r.bold = True
        r.font.color.rgb = RGBColor(0xF9, 0x73, 0x16)
        p.add_run(f['diagnostic'])

        p = doc.add_paragraph()
        r = p.add_run("Actions réalisées : ")
        r.bold = True
        r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        for line in f['actions'].split('\n'):
            if line.strip():
                p = doc.add_paragraph(line.strip(), style='List Bullet')
                p.paragraph_format.space_after = Pt(3)

        p = doc.add_paragraph()
        r = p.add_run("Résultat : ")
        r.bold = True
        r.font.color.rgb = RGBColor(0x05, 0x96, 0x69)
        run = p.add_run(f['resultat'])
        run.font.color.rgb = RGBColor(0x05, 0x96, 0x69)
        run.bold = True
        p.paragraph_format.space_after = Pt(10)
        doc.add_paragraph()

    # Méthodo
    doc.add_heading("Méthodologie d'Intervention", level=1)
    meth = [
        ("1. SÉCURISATION", "Consignation LOTO (arrêt électrique/pneumatique/hydraulique). EPI. Signalisation de zone."),
        ("2. DIAGNOSTIC STRUCTURÉ", "Observation, mesures (multimètre, manomètre, chronomètre), lecture alarmes HMI, consultation schémas et GMAO. Aucune hypothèse sans vérification."),
        ("3. RCA", "Analyse cause racine : Pareto, 5 pourquoi, Ishikawa. Identifier cause physique + cause systémique."),
        ("4. INTERVENTION", "Réparation/remplacement, réglage, programmation. Respect des couples, procédures, normes."),
        ("5. VALIDATION & REPORTING", "Essais sous charge, mesures avant/après. Saisie GMAO. Mise à jour documentation. Formation transfert."),
    ]
    for t, d in meth:
        p = doc.add_paragraph()
        r = p.add_run(t + " — ")
        r.bold = True
        r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        p.add_run(d)

    doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(f"{PROFILE['name']} — {PROFILE['title']}\n{PROFILE['contact']}")
    r.italic = True
    r.font.size = Pt(9)
    r.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    path = os.path.join(OUTPUT_DIR, "Portfolio_Salah_Barki_v3.docx")
    doc.save(path)
    print(f"[OK] Word v3 : {path}")


def create_ppt():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    C_BG = RGBColor(0xFF, 0xFF, 0xFF)
    C_DARK = RGBColor(0x0F, 0x17, 0x2A)
    C_ACCENT = RGBColor(0xF9, 0x73, 0x16)
    C_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
    C_TEXT = RGBColor(0x37, 0x41, 0x51)
    C_MUTED = RGBColor(0x94, 0xA3, 0xB8)
    C_GREEN = RGBColor(0x05, 0x96, 0x69)

    def add_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = C_BG

    def add_text(slide, left, top, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Calibri'
        p.alignment = align
        return box

    def add_bullet_text(slide, left, top, width, height, lines, font_size, color):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = 'Calibri'
            p.space_after = Pt(5)
        return box

    blank = prs.slide_layouts[6]

    # Slide 1 Title
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    s.fill.solid(); s.fill.fore_color.rgb = C_ACCENT; s.line.fill.background()
    add_text(slide, Inches(1), Inches(2.3), Inches(11.333), Inches(1.2), "PORTFOLIO TECHNIQUE", 44, C_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(3.4), Inches(11.333), Inches(0.8), "d'Interventions Maintenance Industrielle", 20, C_ACCENT, align=PP_ALIGN.CENTER)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.1), Inches(2.333), Inches(0.01))
    s.fill.solid(); s.fill.fore_color.rgb = C_MUTED; s.line.fill.background()
    add_text(slide, Inches(1), Inches(4.4), Inches(11.333), Inches(0.6), PROFILE["name"], 24, C_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5), PROFILE["contact"], 11, C_MUTED, align=PP_ALIGN.CENTER)

    # Slide 2 Profil
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
    s.fill.solid(); s.fill.fore_color.rgb = C_DARK; s.line.fill.background()
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.5), "PROFIL & EXPERTISE", 22, C_BG, bold=True)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.1), Inches(6.2), Inches(5.8))
    s.fill.solid(); s.fill.fore_color.rgb = C_GRAY; s.line.fill.background()
    add_text(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.4), PROFILE["summary"], 13, C_TEXT)
    stats = [("10+", "Ans exp."), ("7", "Entreprises"), ("4", "Pays"), ("6", "Fiches")]
    y = 1.3
    for num, lab in stats:
        add_text(slide, Inches(7.2), Inches(y), Inches(2), Inches(0.5), num, 26, C_ACCENT, bold=True)
        add_text(slide, Inches(9.0), Inches(y+0.1), Inches(3), Inches(0.4), lab, 13, C_TEXT)
        y += 1.2

    # Slide 3 Competences
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
    s.fill.solid(); s.fill.fore_color.rgb = C_DARK; s.line.fill.background()
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.5), "COMPÉTENCES & TECHNOLOGIES", 22, C_BG, bold=True)
    comp_blocks = [
        ("⚡ ÉLECTROTECHNIQUE", "Siemens S7, TIA Portal\nStep 7, Ladder\nCâblage BT, schémas IEC\nVariateurs, servos, capteurs"),
        ("🔧 MÉCANIQUE", "Presses injection\nLignes PET / Canettes\nHydraulique / Pneumatique\nTransmissions, guidages"),
        ("📊 MÉTHODES", "GMAO SAP PM\nRCA, 5 pourquoi\nPareto, SPC\nKPI MTBF, MTTR, OEE"),
    ]
    x = 0.4
    for title, body in comp_blocks:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.2), Inches(4.1), Inches(5.6))
        s.fill.solid(); s.fill.fore_color.rgb = C_GRAY; s.line.fill.background()
        add_text(slide, Inches(x+0.15), Inches(1.35), Inches(3.8), Inches(0.5), title, 14, C_ACCENT, bold=True)
        add_text(slide, Inches(x+0.15), Inches(1.9), Inches(3.8), Inches(4.8), body, 13, C_TEXT)
        x += 4.3

    # Fiches
    for f in FICHES:
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
        s.fill.solid(); s.fill.fore_color.rgb = C_DARK; s.line.fill.background()
        add_text(slide, Inches(0.4), Inches(0.12), Inches(12), Inches(0.5), f"{f['num']}  {f['titre']}", 20, C_BG, bold=True)
        add_text(slide, Inches(0.4), Inches(0.55), Inches(12), Inches(0.4), f"{f['entreprise']} — {f['lieu']}  |  {f['type']}", 10, C_MUTED)

        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.2), Inches(6.2), Inches(5.9))
        s.fill.solid(); s.fill.fore_color.rgb = C_GRAY; s.line.fill.background()
        add_text(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.4), "DIAGNOSTIC", 12, C_ACCENT, bold=True)
        diag_text = f"Équipement : {f['equip']}\n\nContexte : {f['contexte']}\n\nSymptôme : {f['symptome']}\n\n{f['diagnostic']}"
        add_text(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(5.3), diag_text, 10, C_TEXT)

        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(1.2), Inches(0.05), Inches(5.9))
        s.fill.solid(); s.fill.fore_color.rgb = C_ACCENT; s.line.fill.background()
        add_text(slide, Inches(6.9), Inches(1.3), Inches(6), Inches(0.4), "ACTIONS", 12, C_DARK, bold=True)
        lines = [l.strip() for l in f['actions'].split('\n') if l.strip()]
        add_bullet_text(slide, Inches(6.9), Inches(1.7), Inches(6), Inches(4.5), lines, 10, C_TEXT)

        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(6.4), Inches(12.7), Inches(0.9))
        s.fill.solid(); s.fill.fore_color.rgb = C_GREEN; s.line.fill.background()
        add_text(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), f"RÉSULTAT : {f['resultat']}", 11, C_BG, bold=True)

    # Final
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    s.fill.solid(); s.fill.fore_color.rgb = C_ACCENT; s.line.fill.background()
    add_text(slide, Inches(1), Inches(2.4), Inches(11.333), Inches(1), "MÉTHODOLOGIE", 32, C_DARK, bold=True, align=PP_ALIGN.CENTER)
    meth_lines = [
        "1. SÉCURISATION — Consignation LOTO, EPI, signalisation",
        "2. DIAGNOSTIC — Mesures, alarmes HMI, schémas, GMAO. Pas d'hypothèse sans vérification",
        "3. RCA — 5 pourquoi, Ishikawa, Pareto. Cause physique + systémique",
        "4. INTERVENTION — Réparation, réglage, programmation. Couples, procédures, normes",
        "5. VALIDATION — Essais charge, mesures A/A, GMAO, formation"
    ]
    y = 3.4
    for line in meth_lines:
        add_text(slide, Inches(2), Inches(y), Inches(9.333), Inches(0.5), line, 15, C_TEXT, align=PP_ALIGN.CENTER)
        y += 0.55
    add_text(slide, Inches(1), Inches(6.2), Inches(11.333), Inches(0.5), f"{PROFILE['name']} — {PROFILE['contact']}", 11, C_MUTED, align=PP_ALIGN.CENTER)

    path = os.path.join(OUTPUT_DIR, "Portfolio_Salah_Barki_v3.pptx")
    prs.save(path)
    print(f"[OK] PPT v3 : {path}")


def create_pdf():
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY

    path = os.path.join(OUTPUT_DIR, "Portfolio_Salah_Barki_v3.pdf")
    doc = SimpleDocTemplate(path, pagesize=A4, rightMargin=1.8*cm, leftMargin=1.8*cm, topMargin=1.8*cm, bottomMargin=1.8*cm)
    styles = getSampleStyleSheet()

    s_title = ParagraphStyle(name='s_title_v3', parent=styles['Title'], fontSize=24, textColor=colors.HexColor('#0f172a'), alignment=TA_CENTER, spaceAfter=6, fontName='Helvetica-Bold')
    s_sub = ParagraphStyle(name='s_sub_v3', fontSize=13, textColor=colors.HexColor('#f97316'), alignment=TA_CENTER, spaceAfter=10, fontName='Helvetica-Bold')
    s_name = ParagraphStyle(name='s_name_v3', fontSize=15, textColor=colors.HexColor('#0f172a'), alignment=TA_CENTER, spaceAfter=4, fontName='Helvetica-Bold')
    s_contact = ParagraphStyle(name='s_contact_v3', fontSize=9, textColor=colors.HexColor('#94a3b8'), alignment=TA_CENTER, spaceAfter=12)
    s_h1 = ParagraphStyle(name='s_h1_v3', parent=styles['Heading1'], fontSize=15, textColor=colors.HexColor('#0f172a'), spaceAfter=8, spaceBefore=12, fontName='Helvetica-Bold')
    s_h2 = ParagraphStyle(name='s_h2_v3', parent=styles['Heading2'], fontSize=12, textColor=colors.HexColor('#0f172a'), spaceAfter=6, spaceBefore=10, fontName='Helvetica-Bold')
    s_body = ParagraphStyle(name='s_body_v3', parent=styles['Normal'], fontSize=9.5, leading=13, alignment=TA_JUSTIFY, spaceAfter=5, textColor=colors.HexColor('#374151'))
    s_bullet = ParagraphStyle(name='s_bullet_v3', parent=s_body, leftIndent=12, spaceAfter=3, bulletIndent=6, bulletFontName='Helvetica-Bold', bulletFontSize=9, bulletColor=colors.HexColor('#f97316'))
    s_result = ParagraphStyle(name='s_result_v3', parent=s_body, textColor=colors.HexColor('#059669'), fontSize=9.5, leading=13, spaceAfter=8, leftIndent=6, backColor=colors.HexColor('#f0fdf4'), borderColor=colors.HexColor('#059669'), borderWidth=2, borderPadding=6, leftPadding=8, rightPadding=8, topPadding=6, bottomPadding=6)
    s_diag = ParagraphStyle(name='s_diag_v3', parent=s_body, backColor=colors.HexColor('#f8fafc'), leftIndent=6, borderColor=colors.HexColor('#e2e8f0'), borderWidth=1, borderPadding=6, leftPadding=8, rightPadding=8, topPadding=6, bottomPadding=6)

    story = []
    story.append(Paragraph("PORTFOLIO TECHNIQUE", s_title))
    story.append(Paragraph("d'Interventions Maintenance Industrielle", s_sub))
    story.append(Paragraph(PROFILE["name"], s_name))
    story.append(Paragraph(PROFILE["contact"].replace('📍','').replace('📱','').replace('✉️',''), s_contact))
    story.append(Paragraph(PROFILE["summary"], ParagraphStyle(name='s_sum_v3', parent=s_body, backColor=colors.HexColor('#f8fafc'), borderColor=colors.HexColor('#e2e8f0'), borderWidth=1, borderPadding=8, leftPadding=10, rightPadding=10, topPadding=8, bottomPadding=8)))
    story.append(Spacer(1, 0.3*cm))

    data = [
        [Paragraph("<b>10+</b>", ParagraphStyle(name='sk1', alignment=TA_CENTER, fontSize=16, textColor=colors.HexColor('#f97316'))),
         Paragraph("<b>7</b>", ParagraphStyle(name='sk2', alignment=TA_CENTER, fontSize=16, textColor=colors.HexColor('#f97316'))),
         Paragraph("<b>4</b>", ParagraphStyle(name='sk3', alignment=TA_CENTER, fontSize=16, textColor=colors.HexColor('#f97316'))),
         Paragraph("<b>6</b>", ParagraphStyle(name='sk4', alignment=TA_CENTER, fontSize=16, textColor=colors.HexColor('#f97316')))],
        [Paragraph("Années d'exp.", ParagraphStyle(name='sk5', alignment=TA_CENTER, fontSize=9, textColor=colors.HexColor('#64748b'))),
         Paragraph("Entreprises", ParagraphStyle(name='sk6', alignment=TA_CENTER, fontSize=9, textColor=colors.HexColor('#64748b'))),
         Paragraph("Pays", ParagraphStyle(name='sk7', alignment=TA_CENTER, fontSize=9, textColor=colors.HexColor('#64748b'))),
         Paragraph("Fiches phares", ParagraphStyle(name='sk8', alignment=TA_CENTER, fontSize=9, textColor=colors.HexColor('#64748b')))],
    ]
    t = Table(data, colWidths=[4*cm, 4*cm, 4*cm, 4*cm])
    t.setStyle(TableStyle([
        ('ALIGN', (0,0), (-1,-1), 'CENTER'), ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('LEFTPADDING', (0,0), (-1,-1), 6), ('RIGHTPADDING', (0,0), (-1,-1), 6), ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#f8fafc')),
    ]))
    story.append(t)
    story.append(Spacer(1, 0.3*cm))

    story.append(Paragraph("COMPÉTENCES & TECHNOLOGIES", s_h1))
    comps = [
        ("ÉLECTROTECHNIQUE & AUTOMATISME", "Siemens S7, TIA Portal, Step 7, programmation Ladder, câblage BT, lecture schémas électriques, variateurs, servo-moteurs, capteurs inductifs/photoélectriques/encodeurs."),
        ("MÉCANIQUE & FLUIDIQUE", "Presses injection, lignes PET/canettes, hydraulique (pompes, vérins, distributeurs, accumulateurs), pneumatique (distributeurs 5/2, vérins, préparation air), transmissions, guidages, roulements."),
        ("MÉTHODES & GESTION", "GMAO (SAP PM), planification préventive/corrective, analyse RCA (5 pourquoi, Ishikawa, Pareto), suivi KPI (MTBF, MTTR, OEE), consignation LOTO, normes HSE."),
    ]
    for cat, detail in comps:
        story.append(Paragraph(f"<b>{cat}</b> — {detail}", s_body))
    story.append(PageBreak())

    story.append(Paragraph("FICHES D'INTERVENTION", s_h1))
    story.append(Paragraph("Format : Contexte → Diagnostic structuré → Actions concrètes → Résultats chiffrés", s_body))
    story.append(Spacer(1, 0.2*cm))

    for f in FICHES:
        story.append(Paragraph(f"{f['num']}. {f['titre']}", s_h2))
        story.append(Paragraph(f"<i>{f['entreprise']} — {f['lieu']}  |  {f['type']}</i>", ParagraphStyle(name='s_meta_v3', parent=s_body, textColor=colors.HexColor('#64748b'), fontSize=9)))
        story.append(Paragraph(f"<b>Équipement :</b> {f['equip']}", s_body))
        story.append(Paragraph(f"<b>Contexte :</b> {f['contexte']}", s_body))
        story.append(Paragraph(f"<b>Symptôme :</b> {f['symptome']}", s_body))
        story.append(Paragraph(f"<b>Diagnostic :</b><br/>{f['diagnostic']}", s_diag))
        story.append(Paragraph("<b>Actions réalisées :</b>", s_body))
        for line in f['actions'].split('\n'):
            if line.strip():
                story.append(Paragraph(line.strip(), s_bullet))
        story.append(Paragraph(f"<b>RÉSULTAT :</b> {f['resultat']}", s_result))
        story.append(Spacer(1, 0.2*cm))

    story.append(PageBreak())
    story.append(Paragraph("MÉTHODOLOGIE D'INTERVENTION", s_h1))
    meths = [
        ("1. SÉCURISATION", "Consignation LOTO (arrêt électrique/pneumatique/hydraulique). EPI. Signalisation de zone."),
        ("2. DIAGNOSTIC STRUCTURÉ", "Observation, mesures (multimètre, manomètre, chronomètre), lecture alarmes HMI, consultation schémas et GMAO. Aucune hypothèse sans vérification."),
        ("3. RCA", "Analyse cause racine : Pareto, 5 pourquoi, Ishikawa. Identifier cause physique + cause systémique."),
        ("4. INTERVENTION", "Réparation/remplacement, réglage, programmation. Respect des couples, procédures, normes."),
        ("5. VALIDATION & REPORTING", "Essais sous charge, mesures avant/après. Saisie GMAO. Mise à jour documentation. Formation transfert."),
    ]
    for t, d in meths:
        story.append(Paragraph(f"<b>{t}</b> — {d}", s_body))
    story.append(Spacer(1, 0.4*cm))
    story.append(Paragraph(f"<b>{PROFILE['name']}</b> — {PROFILE['title']}<br/>{PROFILE['contact']}", ParagraphStyle(name='s_footer_v3', alignment=TA_CENTER, fontSize=9, textColor=colors.HexColor('#94a3b8'))))

    doc.build(story)
    print(f"[OK] PDF v3 : {path}")


def create_html():
    html_content = """<!DOCTYPE html>
<html lang="fr">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Portfolio — Salah Eddine Barki</title>
<style>
:root { --bg:#ffffff; --dark:#0f172a; --accent:#f97316; --gray:#f1f5f9; --text:#374151; --muted:#94a3b8; --green:#059669; --green-bg:#f0fdf4; --border:#e2e8f0; }
* { margin:0; padding:0; box-sizing:border-box; }
body { font-family:'Segoe UI',Roboto,Helvetica,Arial,sans-serif; background:var(--bg); color:var(--text); line-height:1.55; font-size:15px; }
.container { max-width:980px; margin:0 auto; padding:0 24px; }
header { background:var(--bg); border-bottom:3px solid var(--accent); padding:48px 0 32px; text-align:center; }
header h1 { font-size:2.8rem; font-weight:300; color:var(--dark); letter-spacing:-1px; }
header .subtitle { font-size:1.15rem; color:var(--accent); font-weight:600; margin:8px 0; text-transform:uppercase; letter-spacing:1px; }
header .name { font-size:1.6rem; color:var(--dark); font-weight:700; margin-top:12px; }
header .contact { font-size:0.9rem; color:var(--muted); margin-top:10px; }
header .summary { max-width:780px; margin:20px auto 0; background:var(--gray); padding:18px 22px; border-radius:10px; font-size:0.95rem; color:var(--text); text-align:justify; border:1px solid var(--border); }
.stats { display:grid; grid-template-columns:repeat(auto-fit,minmax(140px,1fr)); gap:16px; margin:28px 0; }
.stat-card { background:var(--gray); border:1px solid var(--border); border-radius:10px; padding:20px 10px; text-align:center; }
.stat-card .num { font-size:1.9rem; font-weight:700; color:var(--accent); }
.stat-card .lab { font-size:0.8rem; color:var(--muted); text-transform:uppercase; letter-spacing:0.5px; margin-top:4px; }
.section { padding:36px 0; }
.section-title { font-size:1.4rem; font-weight:700; color:var(--dark); margin-bottom:18px; display:flex; align-items:center; gap:10px; }
.section-title::after { content:""; flex:1; height:2px; background:var(--border); }
.comp-grid { display:grid; grid-template-columns:repeat(auto-fit,minmax(280px,1fr)); gap:16px; }
.comp-card { background:var(--gray); border-radius:10px; padding:18px; border-left:4px solid var(--accent); }
.comp-card h3 { font-size:0.95rem; color:var(--accent); margin-bottom:8px; text-transform:uppercase; letter-spacing:0.5px; }
.comp-card p { font-size:0.9rem; color:var(--text); }
.fiche { margin-bottom:32px; }
.fiche-header { background:var(--dark); color:#fff; padding:14px 18px; border-radius:10px 10px 0 0; display:flex; justify-content:space-between; align-items:center; flex-wrap:wrap; gap:8px; }
.fiche-header h3 { font-size:1.1rem; font-weight:600; }
.fiche-header .meta { font-size:0.8rem; color:var(--muted); }
.fiche-header .badge { background:var(--accent); color:#fff; padding:4px 10px; border-radius:20px; font-size:0.7rem; font-weight:700; text-transform:uppercase; }
.fiche-body { background:var(--bg); border:1px solid var(--border); border-top:none; border-radius:0 0 10px 10px; padding:18px; }
.fiche-grid { display:grid; grid-template-columns:1fr 1fr; gap:18px; }
@media(max-width:768px){ .fiche-grid { grid-template-columns:1fr; } }
.fiche-col { background:var(--gray); border-radius:8px; padding:14px; }
.fiche-col h4 { font-size:0.8rem; text-transform:uppercase; letter-spacing:1px; color:var(--accent); margin-bottom:10px; }
.fiche-col p,.fiche-col li { font-size:0.88rem; color:var(--text); margin-bottom:6px; }
.fiche-col ul { list-style:none; padding-left:0; }
.fiche-col ul li { position:relative; padding-left:14px; }
.fiche-col ul li::before { content:"›"; position:absolute; left:0; color:var(--accent); font-weight:bold; }
.result-box { background:var(--green-bg); border-left:4px solid var(--green); padding:12px 14px; border-radius:0 8px 8px 0; margin-top:12px; font-size:0.9rem; color:var(--green); font-weight:600; }
.method-grid { display:grid; grid-template-columns:repeat(auto-fit,minmax(220px,1fr)); gap:14px; }
.method-card { background:var(--gray); border-radius:10px; padding:16px; text-align:center; border-top:3px solid var(--accent); }
.method-card .step { font-size:1.6rem; color:var(--accent); font-weight:700; margin-bottom:6px; }
.method-card h4 { font-size:0.9rem; color:var(--dark); margin-bottom:6px; }
.method-card p { font-size:0.85rem; color:var(--text); }
footer { text-align:center; padding:30px 0; margin-top:20px; border-top:1px solid var(--border); font-size:0.85rem; color:var(--muted); }
</style>
</head>
<body>
<header>
  <div class="container">
    <h1>Portfolio Technique</h1>
    <div class="subtitle">d'Interventions Maintenance Industrielle</div>
    <div class="name">Salah Eddine Barki</div>
    <div class="contact">📍 Amiens, France &nbsp;|&nbsp; 📱 +33 6 88 69 07 04 &nbsp;|&nbsp; ✉️ salahbarki.seb@gmail.com</div>
    <div class="summary">""" + PROFILE["summary"] + """</div>
    <div class="stats">
      <div class="stat-card"><div class="num">10+</div><div class="lab">Années</div></div>
      <div class="stat-card"><div class="num">7</div><div class="lab">Entreprises</div></div>
      <div class="stat-card"><div class="num">4</div><div class="lab">Pays</div></div>
      <div class="stat-card"><div class="num">6</div><div class="lab">Fiches phares</div></div>
    </div>
  </div>
</header>
<section class="section">
  <div class="container">
    <div class="section-title">Compétences & Technologies</div>
    <div class="comp-grid">
      <div class="comp-card"><h3>⚡ Électrotechnique & Automatisme</h3><p>Siemens S7, TIA Portal, Step 7, programmation Ladder, câblage BT, lecture schémas électriques, variateurs, servo-moteurs, capteurs inductifs/photoélectriques/encodeurs.</p></div>
      <div class="comp-card"><h3>🔧 Mécanique & Fluidique</h3><p>Presses injection, lignes PET/canettes, hydraulique (pompes, vérins, distributeurs, accumulateurs), pneumatique (distributeurs 5/2, vérins, préparation air), transmissions, guidages, roulements.</p></div>
      <div class="comp-card"><h3>📊 Méthodes & Gestion</h3><p>GMAO (SAP PM), planification préventive/corrective, analyse RCA (5 pourquoi, Ishikawa, Pareto), suivi KPI (MTBF, MTTR, OEE), consignation LOTO, normes HSE.</p></div>
    </div>
  </div>
</section>
<section class="section">
  <div class="container">
    <div class="section-title">Fiches d'Intervention Détaillées</div>
    <p style="margin-bottom:22px; color:var(--muted); font-size:0.9rem;">Format : Contexte → Diagnostic structuré → Actions concrètes → Résultats chiffrés</p>
"""
    for f in FICHES:
        actions_html = "\n".join([f"<li>{l.strip()}</li>" for l in f['actions'].split('\n') if l.strip()])
        html_content += f"""
    <div class="fiche">
      <div class="fiche-header">
        <div><h3>{f['num']}. {f['titre']}</h3><div class="meta">{f['entreprise']} — {f['lieu']}</div></div>
        <span class="badge">{f['type']}</span>
      </div>
      <div class="fiche-body">
        <p style="font-size:0.85rem; color:var(--muted); margin-bottom:12px;"><b>Équipement :</b> {f['equip']}</p>
        <div class="fiche-grid">
          <div class="fiche-col">
            <h4>Diagnostic</h4>
            <p><b>Contexte :</b> {f['contexte']}</p>
            <p><b>Symptôme :</b> {f['symptome']}</p>
            <p style="white-space:pre-line;">{f['diagnostic']}</p>
          </div>
          <div class="fiche-col">
            <h4>Actions</h4>
            <ul>{actions_html}</ul>
          </div>
        </div>
        <div class="result-box">RÉSULTAT : {f['resultat']}</div>
      </div>
    </div>
"""
    html_content += """
  </div>
</section>
<section class="section">
  <div class="container">
    <div class="section-title">Méthodologie d'Intervention</div>
    <div class="method-grid">
      <div class="method-card"><div class="step">1</div><h4>SÉCURISATION</h4><p>Consignation LOTO, EPI, signalisation de zone.</p></div>
      <div class="method-card"><div class="step">2</div><h4>DIAGNOSTIC</h4><p>Mesures, alarmes HMI, schémas, GMAO. Pas d'hypothèse sans vérification.</p></div>
      <div class="method-card"><div class="step">3</div><h4>RCA</h4><p>5 pourquoi, Ishikawa, Pareto. Cause physique + systémique.</p></div>
      <div class="method-card"><div class="step">4</div><h4>INTERVENTION</h4><p>Réparation, réglage, programmation. Respect procédures.</p></div>
      <div class="method-card"><div class="step">5</div><h4>VALIDATION</h4><p>Essais charge, mesures A/A, GMAO, formation.</p></div>
    </div>
  </div>
</section>
<footer>
  <div class="container">
    <strong>Salah Eddine Barki</strong> — Technicien de Maintenance Industrielle<br>
    📱 +33 6 88 69 07 04 &nbsp;|&nbsp; ✉️ salahbarki.seb@gmail.com &nbsp;|&nbsp; 📍 Amiens, France
  </div>
</footer>
</body>
</html>
"""
    path = os.path.join(OUTPUT_DIR, "Portfolio_Salah_Barki_v3.html")
    with open(path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    print(f"[OK] HTML v3 : {path}")


if __name__ == "__main__":
    create_word()
    create_ppt()
    create_pdf()
    create_html()
    print("\n=== PORTFOLIO V3 COMPLET ===")
