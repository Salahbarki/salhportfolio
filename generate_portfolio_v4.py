# -*- coding: utf-8 -*-
"""
Portfolio Maintenance Salah Barki - V4
Design futuriste tech : light mode + navy profond + vert emeraude
Inspire par Linear / Notion / modern SaaS dashboard
"""
import os

OUTPUT_DIR = r"C:\Users\ba2rb\Downloads\salpre"

PROFILE = {
    "name": "Salah Eddine Barki",
    "title": "Technicien de Maintenance Industrielle",
    "contact": "📍 Amiens, France  |  📱 +33 6 88 69 07 04  |  ✉️ salahbarki.seb@gmail.com",
    "summary": """Technicien de maintenance industrielle senior avec plus de 10 ans d'expérience sur des équipements 
automatisés et mécaniques complexes : presses injection, lignes de production agroalimentaire (PET/canettes), 
bancs de test automobile et utilités industrielles. Compétences transverses : électrotechnique, automatisme 
(Siemens S7, TIA Portal), mécanique, hydraulique et pneumatique. Forte appétence pour le diagnostic 
structuré, l'analyse de pannes récurrentes et l'amélioration continue en collaboration avec les équipes 
production et méthodes."""
}

FICHES = [
    {
        "num": "01",
        "titre": "Panne récurrente encartonneuse — Ligne PET",
        "entreprise": "Coca-Cola Europacific Partners",
        "lieu": "Grigny (91)",
        "type": "CURATIF URGENT",
        "equip": "Encartonneuse / casseuse de colis — Ligne remplissage bouteilles PET 2L — Cadence 60 000 b/h",
        "contexte": "Ligne critique approvisionnant la grande distribution. Production 22h/jour. Tout arrêt >20 min = risque rupture de stock client.",
        "symptome": "Défaut HMI n°402 'PRODUIT_COINCÉ_ENTRÉE' toutes les 8–12 min. Accumulation de 15–18 colis rejetés/heure en sortie fardeleuse. Opérateurs forcés de redémarrer manuellement à chaque arrêt.",
        "diagnostic": "Vérification électrique : alimentation capteur photoélectrique d'entrée produit = 24 V stable. Signal logique oscillant anormalement : la diode de l'automate clignote au lieu de rester allumée fixe.\n\nAnalyse sur TIA Portal (programme automate Siemens) : dans le bloc de gestion encartonneuse, le bit de validation produit est conditionné directement par le signal capteur sans temporisation d'antirebond. L'oscillation du faisceau (probablement due à la condensation/buée ambiante près de la ligne froide) crée des fronts parasites interprétés comme des produits valides.\n\nVérification mécanique : jeu butée d'entrée produit anormalement large (environ 3 mm au lieu du réglage usuel ~0,5 mm). Amortisseur usé qui ne freine plus correctement le produit.\n\nVérification servo : en mode manuel, le pousseur arrive en butée 2–3 mm en retard par rapport à la consigne affichée sur l'écran. Dérive mécanique liée au jeu de la butée.",
        "actions": [
            "CONSIGNATION : Arrêt électrique + pneumatique de la zone. Condamnateur cadenas apposé. Signalisation verticale de zone de maintenance.",
            "ÉLECTRIQUE : Remplacement du capteur photoélectrique par un modèle plus robuste avec purge d'air intégrée (anti-buée), adapté à l'environnement humide proche de la ligne froide. Remplacement du câble d'alimentation par un câble blindé pour éviter les perturbations électromagnétiques des variateurs voisins. Vérification : signal désormais stable, diode automate allumée fixe.",
            "MÉCANIQUE : Dépose de la butée d'entrée. Remplacement de l'amortisseur hydraulique usé par un neuf. Recalage de la butée au jeu de 0,5 mm (vérifié à la cale/lame de 0,5 mm). Serrage de la boulonnerie au couple avec clé dynamométrique.",
            "AUTOMATISME (TIA Portal) : Ajout d'une temporisation de 50 ms dans le programme automate sur le signal du capteur avant validation du bit 'produit présent'. Cela filtre les oscillations parasites dues à l'humidité sans impacter le temps de cycle. Recalage du décalage servo-pousseur (ajustement du paramètre de position dans le variateur). Validation par mouvements manuels : le pousseur arrive pile en butée avec le convoyeur.",
            "ESSAIS : Démarrage de la ligne. Run de 45 min à cadence nominale. Comptage des rejets : 2 colis rejetés sur 45 min (contre 12–15 avant). Aucun arrêt forcé opérateur. Monitoring via l'écran HMI : le compteur de défaut 402 reste à 0.",
            "GMAO : Création d'un ordre de travail récurrent 'Contrôle capteur + butée encartonneuse' tous les 3 mois (au lieu de l'annuel constructeur). Mise à jour de la gamme de maintenance dans la base informatique du site."
        ],
        "resultat": "Temps d'arrêt moyen réduit de ~40 min à 8 min (–80%). Taux de rejets divisé par 5 (17/h → 3/h). Zero arrêt forcé opérateur pendant les 3 semaines suivantes. Procédure de maintenance préventive mise à jour et appliquée sur les 3 autres lignes du site."
    },
    {
        "num": "02",
        "titre": "Mise en service presse injection + auxiliaires",
        "entreprise": "TE Connectivity",
        "lieu": "Tanger (Maroc)",
        "type": "INSTALLATION",
        "equip": "Presse injection 120 tonnes + sécheur de granulés + trémie + détecteur de métal",
        "contexte": "Relocalisation complète d'une cellule d'injection dans un nouveau hall de production. Objectif : production client automobile en J+2 (just-in-time, aucun stock tampon).",
        "symptome": "Aucun — Nouvelle installation. Objectif : démarrage qualité premier jet, cycle stable <20 secondes.",
        "diagnostic": "Analyse pré-installation : vérification de la compatibilité des arrivées électriques (400V triphasé + neutre + terre), hydrauliques (120 bar) et eau de refroidissement (20 °C). Plan de pose transmis par le constructeur. Vérification du niveau d'huile hydraulique et de la propreté des circuits.",
        "actions": [
            "TRANSPORT & POSE : Dépose de l'outillage (moule 24 cavités) avant déplacement. Transport chariot élévateur avec fourches allongées et sangles de levage. Pose sur 4 plots anti-vibratoires. Nivellement au niveau à bulle puis affinage au laser (écart <0,2 mm sur toute la longueur). Alignement visuel de l'axe presse avec le convoyeur amont (cordeau + règle de mécanicien).",
            "RACCORDEMENTS ÉLECTRIQUE : branchement 400V tri + N + T via arrivée dédiée au tableau divisionnaire. Vérification terre au multimètre : résistance <1 Ω (conforme norme).",
            "RACCORDEMENTS HYDRAULIQUE : raccordement flexibles sur centrale hydraulique du hall. Serrage des raccords au couple. Essai pression statique 10% au-dessus de la nominale pendant 10 min : aucune fuite constatée.",
            "EAU REFROIDISSEMENT : raccordement circuits moule (fixe et mobile) sur tour aéroréfrigérant. Température consigne 80 °C côté fixe, 75 °C côté mobile. Débit contrôlé à la vanne : 8 L/min circuit A, 12 L/min circuit B.",
            "SÉCHEUR : Mise en service du sécheur de granulés plastiques. Objectif point de rosée ≤ –40 °C (matière PA66 très hygroscopique). Après 45 min de fonctionnement, lecture du point de rosée sur l'afficheur : –42 °C. Vérification à la sonde thermométrique : température air chauffé 80 °C en sortie vers trémie.",
            "DÉTECTEUR MÉTAL : Installation en amont de la trémie. Réglage de la sensibilité : test avec morceaux de fer, cuivre, aluminium. Validation : détection 10/10, aucun faux positif sur passage de granulés vides. Serrage des écrous de fixation + vérification blindage câble.",
            "PARAMÉTRAGE CYCLE : Réglage vitesse injection, pression de maintien, temps de refroidissement et ouverture moule selon la fiche de paramètres du moule (transmise par le bureau d'études). Premier cycle : réglage du coussin de fin de course. Batch test de 500 pièces.",
            "CONTRÔLE QUALITÉ : Mesure dimensionnelle des premières pièces au pied à coulisse digital (tolérance connecteur ±0,05 mm). Résultat : 0 pièce hors tolérance sur les 500.",
            "DOCUMENTATION & FORMATION : Rédaction du PV de recette (procès-verbal de mise en service) signé par Production + Qualité + Maintenance. Formation de 3 opérateurs : démarrage/arrêt, changement de référence (objectif 8 min), lecture des alarmes HMI."
        ],
        "resultat": "Mise en production effective en J+2 (mardi matin après déménagement vendredi). Cycle stable mesuré : 18,5 s moyenne (objectif <20 s). 0 défaut qualité sur les 500 pièces de validation. Checklist de mise en service de 42 points créée et partagée à l'équipe : 2 autres presses relocalisées le trimestre suivant en réutilisant la même procédure."
    },
    {
        "num": "03",
        "titre": "Débogage & optimisation programme automate — Cellule assemblage",
        "entreprise": "TE Connectivity",
        "lieu": "Tanger (Maroc)",
        "type": "AUTOMATISME",
        "equip": "Cellule assemblage connectique — Automate Siemens S7-1200 + HMI + servo",
        "contexte": "Cellule critique assemblage connecteurs automobile. Cadence nominale 220 pièces/heure. OEE mensuel affiché en baisse (72% vs objectif 85%).",
        "symptome": "Temps de cycle mesuré au chronomètre : 28–30 s au lieu des 16–17 s nominaux. Alarme HMI fréquente 'Timeout bras robot zone pick'. Arrêts 3–4 fois par heure. Opérateurs obligés de redémarrer manuellement (perte ~2 min à chaque fois).",
        "diagnostic": "Connexion sur TIA Portal en ligne avec l'automate. Monitoring des blocs de programme.\n\nLe temps d'exécution du bloc de gestion robot est anormalement long (8 secondes mesurées sur l'outil de diagnostic intégré, contre 3–4 s attendu).\n\nAnalyse pas à pas du programme : une temporisation de sécurité est réglée à 8 secondes alors que le mouvement réel du bras robot dure environ 5,5 s (mesuré au chronomètre et confirmé par le suivi des signaux d'entrées/sorties). La temporisation est donc trop conservative et bloque inutilement la suite du cycle.\n\nBit de validation 'pince OK' : capteur de proximité qui détecte la pièce clignote légèrement au lieu d'être fixe. Inspection mécanique : le support du capteur est légèrement desserré (jeu de 2 mm). Le capteur détecte parfois, parfois pas => le programme attend indéfiniment le signal stable.\n\nProgrammation non optimisée : les mouvements du bras sont strictement séquentiels (un à la fois) alors que certains mouvements compatibles pourraient être faits en parallèle (ex: monter le bras tout en le déplaçant latéralement, avec sécurité de hauteur).",
        "actions": [
            "MÉCANIQUE : Réglage du support capteur de proximité. Resserrement de l'écrou de fixation au couple (contre-écrou nylstop pour éviter le desserrage). Vérification : le signal 'pince OK' est désormais stable et fixe (diode verte allumée en continu).",
            "PROGRAMMATION TIA PORTAL : Ajustement de la temporisation de sécurité : réduction de 8 s à 6,5 s (marge de sécurité de 1 s au-dessus du temps mesuré réel 5,5 s). Ajout d'une petite temporisation d'antirebond de 35 ms sur le signal capteur dans le programme pour éliminer les éventuelles micro-oscillations résiduelles.",
            "OPTIMISATION CODE : Réécriture partielle du bloc de gestion mouvements : parallélisation des phases compatibles (descente + fermeture pince ; montée + translation latérale avec interlock de hauteur pour la sécurité). Suppression de temporisations redondantes.",
            "SURVEILLANCE : Ajout d'un compteur de cycles dans le programme avec alarme automatique si le temps moyen dépasse 10% de la valeur nominale (détection précoce de dérive).",
            "HMI : Création d'un écran de diagnostic 'Temps de cycle' affichant le temps réel, la moyenne glissante sur 50 cycles, et le nombre d'alarmes 'Timeout'. Accessible aux opérateurs pour monitoring autonome.",
            "SÉCURITÉ & DOCUMENTATION : Sauvegarde du programme modifié sur le serveur du site avec commentaires détaillés dans chaque section (date, nature modif, nom). Export papier du bloc modifié pour classeur maintenance.",
            "VALIDATION : Run de 4 heures à cadence max. Monitoring temps cycle via HMI et chronomètre externe. Résultats : 16,2 s (min) — 16,8 s (moyenne) — 17,1 s (max). Aucune alarme 'Timeout' durant les 4 heures."
        ],
        "resultat": "Temps de cycle moyen : 28,5 s → 16,8 s (–41%). Capacité effective passée de ~126 à 214 pièces/heure. OEE cellule : 72% → 84% (+12 points). Alarme 'Timeout' : 3–4 arrêts/h → 0 sur 4 h test, puis <1 arrêt/h en production continue. Maintenance autonome sur ces paramètres désormais (formation interne réalisée)."
    },
    {
        "num": "04",
        "titre": "Analyse cause racine — Banc de test électrique",
        "entreprise": "Kromberg & Schubert",
        "lieu": "Kénitra (Maroc)",
        "type": "RCA & QUALITÉ",
        "equip": "Banc de test électrique (continuité, isolation, résistance) — Ligne faisceaux 32 voies",
        "contexte": "Ligne faisceaux moteur pour client automobile. Production 2 000 pièces/jour. Just-in-time : aucun stock tampon autorisé. Tolérance client <2% de rebut test.",
        "symptome": "Taux de rejet test continuité anormal : 12% sur une même référence produite depuis plusieurs jours sans changement d'outillage. Normalement stable autour de 1–2%. Conséquence : 240+ pièces rebutées/jour. Risque blocage ligne client.",
        "diagnostic": "Analyse des données de test des 15 derniers jours (extraites du système informatique du banc) : dérive lente et progressive de la résistance mesurée au test de continuité. Tolérance : <10 mΩ. Valeurs de départ 3–5 mΩ → dérive progressive 8–12 mΩ → dépassement du seuil = rebut.\n\nPareto des causes de rebut sur 342 pièces rebutées : ~2/3 liés aux broches de contact du connecteur adaptateur du banc (pins), 1/4 liés aux câbles de piquage usés, le reste divers.\n\nInspection visuelle des pins : 14 des 32 pins présentent une oxydation noire visible (couche de CuO). Cause probable : absence de nettoyage régulier des pins dans la gamme de maintenance (le dernier nettoyage n'était pas documenté, probablement >6 mois).\n\nHotte d'extraction des fumées de soudure au poste amont : filtre saturé, débit d'aspiration faible (mesuré au débitmètre : 80 m³/h au lieu des ~120 m³/h nécessaires). Les fumées de soudure contenant des résidus de flux dérivent vers le banc de test et contaminent progressivement les contacts.\n\nMesure au micro-ohmètre : pins oxydés = 15–45 mΩ. Pins nettoyés à la brosse + alcool = <2 mΩ.",
        "actions": [
            "ARRÊT & TRI : Arrêt immédiat du banc de test. Consignation électrique 24V + 500V isolation. Signalisation. Récupération du batch produit depuis la dernière validation qualité OK : retri manuel complet. 94% des pièces déclarées bonnes après contrôle visuel + test au micro-ohmètre portable, 6% rebut confirmé.",
            "REMPLACEMENT CONNECTEUR : Dépose du jeu de pins oxydés. Pose d'un connecteur adaptateur neuf (référence constructeur du banc). Serrage des bagues de verrouillage.",
            "RÉGLAGE MÉCANIQUE : Vérin de mise en position du banc : réglage de la pression d'air de 3,2 bar à 4,0 bar (lu sur le manomètre intégré). Vérification à la jauge analogique : 4,0 bar ±0,1. Test de 50 insertions/retraits : aucune marque sur le plastique (pas de déformation des broches).",
            "MAINTENANCE EXTRACTION : Remplacement du filtre charbon de la hotte d'aspiration. Nettoyage de la roue du ventilateur. Graissage des paliers. Mesure post-intervention : débit remonté à 130 m³/h (>120 requis). Ajout d'une hotte aspirante mobile en renfort ponctuel sur le poste de soudure amont pour réduire la dérive à la source.",
            "MODIFICATION PROCÉDURE & GMAO : Modification de la gamme de maintenance préventive : ajout d'un nettoyage des pins de test à la brosse fibre de verre + alcool isopropylique toutes les 2 semaines (au lieu de jamais/non documenté). Ajout d'un contrôle qualité mensuel : mesure au micro-ohmètre sur 5 pins choisies au hasard. Seuil d'alerte <5 mΩ.",
            "FORMATION : Création fiche réflexe opérateur : si dérive >6 mΩ sur 3 tests consécutifs → arrêt immédiat + appel maintenance. Formation de 2 opérateurs + 1 technicien à la nouvelle procédure.",
            "CONTRÔLE RÉGULIER : Mise en place d'un test gabarit 32/32 pins tous les matins avant démarrage de la ligne (détection précoce oxydation ou déformation mécanique)."
        ],
        "resultat": "Taux de rejet test continuité : 12,3% → 1,1% moyenne sur 30 jours (stable sous le seuil client 2%). Économie : ~220 pièces/jour sauvées. Procédure RCA formalisée (diagramme Ishikawa + 5 pourquoi) partagée aux 4 lignes test du site et aux sous-traitants internes. Temps de contrôle qualité départ réduit de 25 min/batch à 10 min/batch (test plus fiable, moins de retests nécessaires)."
    },
    {
        "num": "05",
        "titre": "Dépannage hydraulique presse découpe / pliage",
        "entreprise": "Sovireso",
        "lieu": "Saint-Laurent-sur-Sèvre (85)",
        "type": "CURATIF MÉCANIQUE",
        "equip": "Presse hydraulique découpe/pliage — Vérin double effet Ø100/70 mm — Centrale 80 bar",
        "contexte": "Machine unique sur site pour découpe tôlerie fine (acier 0,8–2 mm, inox 1–1,5 mm). Pas de machine de remplacement. Programme client : 480 pièces/jour.",
        "symptome": "Perte brutale de pression lors de la descente rapide vers position pliage. Manomètre principal : chute de 80 bar à ~30 bar. Bruit violent type 'coup de bélier' à chaque inversion descente/montée. Arrêt machine. Production interrompue depuis 2 heures.",
        "diagnostic": "Mesures aux manomètres multi-points :\n- Pression côté pompe (amont filtre) : 80 bar stable.\n- Pression côté descente (amont vérin tige) : 32 bar (anormal, devrait être ~78–80 bar).\n- Pression côté montée : 78 bar (normal).\n→ Panne localisée côté circuit descente.\n\nAuscultation au stéthoscope mécanique : bruit de coup de bélier localisé au niveau de l'accumulateur hydropneumatique côté montée. Fréquence synchrone avec l'inversion du cycle.\n\nDémontage du clapet anti-retour sur la ligne de descente : présence d'un petit morceau de joint torique écrasé coincé dans le siège du clapet (corps étranger). Le joint du clapet est partiellement dégradé, probablement à cause de la température de l'huile plus élevée que la normale.\n\nTempérature huile mesurée au thermomètre de cuve : 68 °C (normalement <55 °C en exploitation régulière). Cause surchauffe : le filtre de retour est très encrassé (noir et plein de limaille) => la pompe force dans un circuit bouché => surchauffe + dégradation des joints.\n\nVérification de l'accumulateur hydropneumatique : précharge d'azote mesurée au manomètre de contrôle = 12 bar (normalement ~65 bar pour ce modèle). La membrane interne est visiblement fissurée (trace d'huile dans la partie gaz lors du démontage de contrôle).",
        "actions": [
            "CONSIGNATION & VIDANGE : Arrêt moteur pompe au disjoncteur. Condamnateur cadenas + panneau 'Ne pas démarrer — Maintenance'. Vidange des 120 L d'huile dans bac de récupération homologué.",
            "CLAPET ANTI-RETOUR : Dépose complète du clapet de la ligne descente. Nettoyage du siège en carbure à la pâte à roder (grain fin) jusqu'à surface plane et brillante. Test d'étanchéité : remplissage au kérosène, aucune goutte ne passe en 30 secondes. Remplacement du joint torique et du joint plat d'étanchéité par des neufs. Remontage : serrage croisé des vis au couple (clé dynamométrique).",
            "ACCUMULATEUR : Remplacement complet de l'accumulateur (vase + précharge). Précharge vérifiée au manomètre de contrôle : 65 bar. Contrôle 24 h après : aucune perte de pression (étanchéité OK).",
            "FILTRATION : Remplacement du filtre retour (très noir et encrassé) par un neuf. Remplacement du filtre d'aspiration. Remplissage huile neuve conforme spécification constructeur (viscosité ISO VG 46). Purge du circuit par 10 cycles lents sans charge (montée/descente manuelle lente) pour évacuer l'air.",
            "POMPE : Dépose partielle pour inspection des palettes internes. Jeu latéral mesuré au comparateur : 0,12 mm (tolérance max constructeur 0,15 mm). => Acceptable, mais programmation d'une révision dans 6 mois (préventif opportuniste).",
            "MISE EN PRESSION : Montée progressive par paliers (20 bar → 40 bar → 80 bar), 5 min à chaque palier pour détecter d'éventuelles fuites. Contrôle aux raccords avec papier absorbant : aucune fuite. Température huile après 30 min de fonctionnement : 48 °C (OK).",
            "ESSAIS SOUS CHARGE : Découpe de tôle acier 2 mm, 50 cycles continus. Pression affichée : 80 bar stable ±2 bar sur tout le cycle. Disparition complète des coups de bélier (vérifié au stéthoscope + capteur vibration portable)."
        ],
        "resultat": "Intervention réalisée en 3 h 20 (prévision constructeur 1 journée = 7 h). Production relancée le jour même. Pression rétablie à 80 bar stable sur cycle complet descente/montée/pliage. À-coups hydrauliques supprimés. Maintenance préventive avancée : remplacement huile + filtres programmés avec 2 mois d'avance. Révision pompe programmée dans la GMAO."
    },
    {
        "num": "06",
        "titre": "Planification maintenance annuelle & pilotage GMAO",
        "entreprise": "TE Connectivity",
        "lieu": "Tanger (Maroc)",
        "type": "PLANIFICATION",
        "equip": "Site industriel complet : 142 équipements (injection, assemblage, test, utilités, bâtiment)",
        "contexte": "Site de 450 personnes, production 24h/5j. Objectif groupe : taux réalisation maintenance préventive >90% et réduction arrêts imprévus de 25% sur l'exercice.",
        "symptome": "Taux réalisation maintenance préventive année N : 68% (233 ordres de travail réalisés sur 342 planifiés). Les arrêts imprévus représentent 23% du temps de production disponible. Coût maintenance en croissance +12% par rapport au budget. MTBF moyen site : 420 heures (benchmark industrie connectique ~650 h).",
        "diagnostic": "Analyse des données de la GMAO et de la supervision des 12 derniers mois :\n- 342 ordres de travail préventifs planifiés, 233 réalisés (68%), 109 reportés.\n- Pareto des causes de report : ~45% = conflit avec production (pas de fenêtre machine disponible), ~30% = pièces détachées non disponibles (délai d'approvisionnement 5–7 jours), ~15% = sous-traitant spécialisé indisponible.\n\nAbsence de classification des équipements par criticité : tous traités avec la même priorité. Exemple : le compresseur d'air 8 bar (critique, arrêt = arrêt tout le site) et un éclairage de bureau avaient la même fréquence de préventif. Le temps de maintenance était mal réparti.\n\nAnalyse MTBF par famille d'équipements : presses injection 380 h (faible), bancs test 620 h (moyen), utilités 850 h (bon). Aucun plan d'action spécifique sur les presses.",
        "actions": [
            "COLLECTE & ANALYSE : Export des données GMAO (ordres réalisés, reports, coûts). Export des données supervision (arrêts, codes défauts, durées). Croisement dans Excel avec tableau de bord simplifié (graphiques tendance, top 5 pannes du mois).",
            "CLASSIFICATION CRITICITÉ ABC (méthode RCM simplifiée) : Classe A (14% des équipements, ~20 machines) = critique production + sécurité + environnement. → Préventif mensuel ou bimestriel. Stock de pièces critiques en local (capteurs de rechange, vannes, joints courants). Ex : presses injection, compresseur air principal, tour aéroréfrigérant principal. Classe B (32%, ~45 machines) = important mais remplaçable ou redondant. → Préventif trimestriel. Approvisionnement pièces sous 48 h. Classe C (54%, ~77 machines) = standard / non critique. → Préventif semestriel ou annuel. Curatif à la demande.",
            "PLANNING ANNUEL INTÉGRÉ : Création du planning 12 mois dans la GMAO avec fenêtres fixes négociées avec le chef de production : samedi matin (4 h pour A), arrêts planifiés trimestriels (8 h pour grandes révisions), périodes de vacances (révisions lourdes). Les conflits production ont chuté de 45% à 8%.",
            "APPROVISIONNEMENT : Négociation avec fournisseurs locaux pour réduire les délais. Création d'un stock de sécurité classe A (capteurs inductifs et photoélectriques de remplacement, vannes 5/2, joints toriques standards). Délai moyen pièces critiques : 5 j → 2 j.",
            "SOUS-TRAITANCE : Regroupement des interventions par lots (ex: thermographie annuelle, étalonnage des instruments, révision des pompes spécialisées). Appel d'offres 3 fournisseurs. Réduction des coûts sous-traitance : –18%.",
            "PILOTAGE KPI : Création d'un tableau de bord mensuel affiché au local maintenance et partagé avec la production : taux de réalisation du préventif (%), MTBF (heures) et MTTR (minutes) par classe d'équipement, coût maintenance par tonne produite, top 3 pannes du mois + actions en cours.",
            "PROJETS D'AMÉLIORATION : Élaboration et soumission de 5 projets d'amélioration : remplacement de 2 variateurs obsolètes, mise en place d'un suivi vibration sur 3 groupes critiques, formation interne automatisme pour 3 techniciens. 3 projets sur 5 validés par la direction et réalisés en interne."
        ],
        "resultat": "Taux de réalisation préventif : 68% → 94% (année N+1), objectif groupe 90% dépassé. Arrêts imprévus : –28% vs année N (production gagne ~340 heures/an de disponibilité). MTBF site : 420 h → 610 h (+45%), rapproché du benchmark industrie. Coûts maintenance : –15% sous-traitance (achats groupés) et –8% pièces (stock optimisé + négociation). 3 projets amélioration réalisés en interne avec retour sur investissement <18 mois."
    },
]


# ============================================================
# 1) WORD V4 — Futuristic Navy + Emerald
# ============================================================
def create_word():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.enum.table import WD_TABLE_ALIGNMENT

    doc = Document()
    sections = doc.sections[0]
    sections.top_margin = Cm(2)
    sections.bottom_margin = Cm(2)
    sections.left_margin = Cm(2)
    sections.right_margin = Cm(2)

    style = doc.styles['Normal']
    font = style.font
    font.name = 'Segoe UI'
    font.size = Pt(10.5)
    style.element.rPr.rFonts.set(qn('w:eastAsia'), 'Segoe UI')

    # Colors
    C_NAVY = RGBColor(0x0F, 0x17, 0x2A)
    C_EMERALD = RGBColor(0x10, 0xB9, 0x81)
    C_DARK = RGBColor(0x1E, 0x29, 0x3B)
    C_TEXT = RGBColor(0x37, 0x41, 0x51)
    C_MUTED = RGBColor(0x64, 0x71, 0x8B)
    C_GREEN = RGBColor(0x05, 0x96, 0x69)

    # Title page
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("PORTFOLIO")
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = C_NAVY
    r.font.name = 'Segoe UI Light'

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("d'Interventions")
    r.font.size = Pt(20)
    r.font.color.rgb = C_EMERALD
    r.font.name = 'Segoe UI'

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Maintenance Industrielle")
    r.font.size = Pt(20)
    r.font.color.rgb = C_EMERALD

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(PROFILE["name"])
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = C_NAVY

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(PROFILE["contact"])
    r.font.size = Pt(9)
    r.font.color.rgb = C_MUTED
    p.space_after = Pt(20)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    r = p.add_run(PROFILE["summary"])
    r.font.size = Pt(10)
    r.italic = True
    r.font.color.rgb = C_TEXT
    doc.add_paragraph()

    # Stats table
    doc.add_heading("Chiffres Clés", level=1)
    table = doc.add_table(rows=1, cols=4)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cells = table.rows[0].cells
    headers = [("10+", "Années d'exp."), ("7", "Entreprises"), ("4", "Pays"), ("6", "Fiches phares")]
    for i, (num, lab) in enumerate(headers):
        p = cells[i].paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(num)
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = C_EMERALD
        p.add_run("\n" + lab).font.size = Pt(9)
        p.add_run("\n" + lab).font.color.rgb = C_MUTED
    doc.add_paragraph()

    # Skills
    doc.add_heading("Compétences & Technologies", level=1)
    comps = [
        ("⚡ Électrotechnique & Automatisme", "Siemens S7, TIA Portal, Step 7, programmation Ladder, câblage BT, lecture schémas électriques, variateurs, servo-moteurs, capteurs inductifs/photoélectriques/encodeurs."),
        ("🔧 Mécanique & Fluidique", "Presses injection, lignes PET/canettes, hydraulique (pompes, vérins, distributeurs, accumulateurs), pneumatique (distributeurs 5/2, vérins, préparation air), transmissions, guidages, roulements."),
        ("📊 Méthodes & Gestion", "GMAO (SAP PM), planification préventive/corrective, analyse RCA (5 pourquoi, Ishikawa, Pareto), suivi KPI (MTBF, MTTR, OEE), consignation LOTO, normes HSE."),
    ]
    for title, detail in comps:
        p = doc.add_paragraph()
        r = p.add_run(title + " : ")
        r.bold = True
        r.font.color.rgb = C_NAVY
        p.add_run(detail).font.color.rgb = C_TEXT
    doc.add_paragraph()

    # Fiches
    doc.add_page_break()
    doc.add_heading("Fiches d'Intervention", level=1)
    doc.add_paragraph("Format : Contexte → Diagnostic structuré → Actions concrètes → Résultats chiffrés")

    for f in FICHES:
        doc.add_heading(f"{f['num']}. {f['titre']}", level=2)
        p = doc.add_paragraph()
        p.add_run(f"{f['entreprise']} — {f['lieu']}  |  {f['type']}").italic = True
        p.runs[0].font.color.rgb = C_MUTED
        p.runs[0].font.size = Pt(9)

        p = doc.add_paragraph()
        r = p.add_run("Équipement : ")
        r.bold = True
        r.font.color.rgb = C_NAVY
        p.add_run(f['equip'])

        p = doc.add_paragraph()
        r = p.add_run("Contexte : ")
        r.bold = True
        r.font.color.rgb = C_NAVY
        p.add_run(f['contexte'])

        p = doc.add_paragraph()
        r = p.add_run("Symptôme : ")
        r.bold = True
        r.font.color.rgb = C_NAVY
        p.add_run(f['symptome'])

        p = doc.add_paragraph()
        r = p.add_run("Diagnostic : ")
        r.bold = True
        r.font.color.rgb = C_EMERALD
        p.add_run(f['diagnostic'])

        p = doc.add_paragraph()
        r = p.add_run("Actions réalisées : ")
        r.bold = True
        r.font.color.rgb = C_NAVY
        for line in f['actions']:
            p = doc.add_paragraph(line, style='List Bullet')
            p.paragraph_format.space_after = Pt(3)

        p = doc.add_paragraph()
        r = p.add_run("Résultat : ")
        r.bold = True
        r.font.color.rgb = C_GREEN
        run = p.add_run(f['resultat'])
        run.font.color.rgb = C_GREEN
        run.bold = True
        p.paragraph_format.space_after = Pt(10)
        doc.add_paragraph()

    # Method
    doc.add_heading("Méthodologie", level=1)
    meth = [
        ("1. SÉCURISATION", "Consignation LOTO, EPI, signalisation de zone."),
        ("2. DIAGNOSTIC", "Mesures, alarmes HMI, schémas, GMAO. Pas d'hypothèse sans vérification."),
        ("3. RCA", "5 pourquoi, Ishikawa, Pareto. Cause physique + systémique."),
        ("4. INTERVENTION", "Réparation, réglage, programmation. Respect procédures."),
        ("5. VALIDATION", "Essais charge, mesures A/A, GMAO, formation."),
    ]
    for t, d in meth:
        p = doc.add_paragraph()
        r = p.add_run(t + " — ")
        r.bold = True
        r.font.color.rgb = C_NAVY
        p.add_run(d)

    doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(f"{PROFILE['name']} — {PROFILE['title']}\n{PROFILE['contact']}")
    r.italic = True
    r.font.size = Pt(9)
    r.font.color.rgb = C_MUTED

    path = os.path.join(OUTPUT_DIR, "Portfolio_Salah_Barki_v4.docx")
    doc.save(path)
    print(f"[OK] Word v4 : {path}")


# ============================================================
# 2) PPT V4 — Futuristic Dark Navy Header + Emerald Accents
# ============================================================
def create_ppt():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    C_BG = RGBColor(0xFF, 0xFF, 0xFF)
    C_NAVY = RGBColor(0x0F, 0x17, 0x2A)
    C_EMERALD = RGBColor(0x10, 0xB9, 0x81)
    C_GRAY = RGBColor(0xF8, 0xFA, 0xFC)
    C_TEXT = RGBColor(0x37, 0x41, 0x51)
    C_MUTED = RGBColor(0x94, 0xA3, 0xB8)
    C_GREEN = RGBColor(0x05, 0x96, 0x69)
    C_CARD = RGBColor(0xF1, 0xF5, 0xF9)

    def add_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = C_BG

    def add_text(slide, left, top, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Segoe UI'
        p.alignment = align
        return box

    def add_bullet_text(slide, left, top, width, height, lines, font_size, color):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = 'Segoe UI'
            p.space_after = Pt(5)
        return box

    blank = prs.slide_layouts[6]

    # ===== SLIDE 1: HERO =====
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    # Top navy bar
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(3.5))
    s.fill.solid(); s.fill.fore_color.rgb = C_NAVY; s.line.fill.background()
    # Emerald accent line
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.48), Inches(13.333), Inches(0.04))
    s.fill.solid(); s.fill.fore_color.rgb = C_EMERALD; s.line.fill.background()
    add_text(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(1.2), "PORTFOLIO", 54, C_BG, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.8), "d'Interventions Maintenance Industrielle", 24, C_EMERALD, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.6), PROFILE["name"], 28, C_NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.5), PROFILE["contact"], 11, C_MUTED, align=PP_ALIGN.CENTER)
    # Summary card
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.6), Inches(10.333), Inches(1.5))
    s.fill.solid(); s.fill.fore_color.rgb = C_CARD; s.line.fill.background()
    add_text(slide, Inches(1.7), Inches(5.8), Inches(9.9), Inches(1.2), PROFILE["summary"], 12, C_TEXT, align=PP_ALIGN.CENTER)

    # ===== SLIDE 2: STATS + SKILLS =====
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
    s.fill.solid(); s.fill.fore_color.rgb = C_NAVY; s.line.fill.background()
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), "PROFIL & COMPÉTENCES", 22, C_BG, bold=True)
    # Stats cards
    stats = [("10+", "Ans"), ("7", "Entreprises"), ("4", "Pays"), ("6", "Fiches")]
    x = 0.5
    for num, lab in stats:
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.2), Inches(2.8), Inches(1.3))
        s.fill.solid(); s.fill.fore_color.rgb = C_CARD; s.line.fill.background()
        add_text(slide, Inches(x+0.1), Inches(1.35), Inches(2.6), Inches(0.6), num, 28, C_EMERALD, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x+0.1), Inches(1.95), Inches(2.6), Inches(0.4), lab, 12, C_TEXT, align=PP_ALIGN.CENTER)
        x += 3.1
    # Skills
    comps = [
        ("⚡ ÉLECTROTECHNIQUE", "Siemens S7\nTIA Portal / Step 7\nLadder, câblage BT\nVariateurs, capteurs"),
        ("🔧 MÉCANIQUE", "Presses injection\nLignes PET/Canettes\nHydraulique/Pneumatique\nTransmissions, guidages"),
        ("📊 MÉTHODES", "GMAO SAP PM\nRCA, 5 pourquoi\nPareto, SPC\nKPI MTBF, MTTR, OEE"),
    ]
    x = 0.5
    for title, body in comps:
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.8), Inches(4.0), Inches(4.2))
        s.fill.solid(); s.fill.fore_color.rgb = C_CARD; s.line.fill.background()
        add_text(slide, Inches(x+0.15), Inches(2.95), Inches(3.7), Inches(0.5), title, 13, C_EMERALD, bold=True)
        add_text(slide, Inches(x+0.15), Inches(3.5), Inches(3.7), Inches(3.4), body, 12, C_TEXT)
        x += 4.3

    # ===== FICHES =====
    for f in FICHES:
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        # Header navy
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
        s.fill.solid(); s.fill.fore_color.rgb = C_NAVY; s.line.fill.background()
        # Emerald accent strip
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.96), Inches(13.333), Inches(0.06))
        s.fill.solid(); s.fill.fore_color.rgb = C_EMERALD; s.line.fill.background()
        add_text(slide, Inches(0.4), Inches(0.12), Inches(12), Inches(0.5), f"{f['num']}  {f['titre']}", 20, C_BG, bold=True)
        add_text(slide, Inches(0.4), Inches(0.55), Inches(12), Inches(0.4), f"{f['entreprise']} — {f['lieu']}  |  {f['type']}", 10, C_MUTED)

        # Left card: Diagnostic
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.25), Inches(6.2), Inches(5.8))
        s.fill.solid(); s.fill.fore_color.rgb = C_CARD; s.line.fill.background()
        add_text(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.4), "DIAGNOSTIC", 12, C_EMERALD, bold=True)
        diag_text = f"Équipement : {f['equip']}\n\nContexte : {f['contexte']}\n\nSymptôme : {f['symptome']}\n\n{f['diagnostic']}"
        add_text(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5.2), diag_text, 10, C_TEXT)

        # Right: Actions
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.25), Inches(6.3), Inches(4.8))
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF); s.line.fill.background()
        # Emerald left border
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(1.25), Inches(0.05), Inches(4.8))
        s.fill.solid(); s.fill.fore_color.rgb = C_EMERALD; s.line.fill.background()
        add_text(slide, Inches(6.9), Inches(1.4), Inches(6), Inches(0.4), "ACTIONS", 12, C_NAVY, bold=True)
        add_bullet_text(slide, Inches(6.9), Inches(1.8), Inches(6), Inches(4.2), f['actions'], 10, C_TEXT)

        # Result bar
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.9))
        s.fill.solid(); s.fill.fore_color.rgb = C_GREEN; s.line.fill.background()
        add_text(slide, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.7), f"RÉSULTAT : {f['resultat']}", 11, C_BG, bold=True)

    # ===== FINAL: METHODOLOGY =====
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
    s.fill.solid(); s.fill.fore_color.rgb = C_NAVY; s.line.fill.background()
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), "MÉTHODOLOGIE D'INTERVENTION", 22, C_BG, bold=True)
    meths = [
        ("1", "SÉCURISATION", "Consignation LOTO, EPI, signalisation"),
        ("2", "DIAGNOSTIC", "Mesures, alarmes HMI, schémas, GMAO"),
        ("3", "RCA", "5 pourquoi, Ishikawa, Pareto"),
        ("4", "INTERVENTION", "Réparation, réglage, programmation"),
        ("5", "VALIDATION", "Essais charge, mesures A/A, GMAO, formation"),
    ]
    x = 0.4
    for num, title, desc in meths:
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.4), Inches(2.4), Inches(3.8))
        s.fill.solid(); s.fill.fore_color.rgb = C_CARD; s.line.fill.background()
        add_text(slide, Inches(x+0.1), Inches(1.55), Inches(2.2), Inches(0.5), num, 28, C_EMERALD, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x+0.1), Inches(2.2), Inches(2.2), Inches(0.4), title, 12, C_NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x+0.1), Inches(2.65), Inches(2.2), Inches(2.4), desc, 11, C_TEXT, align=PP_ALIGN.CENTER)
        x += 2.55
    add_text(slide, Inches(1), Inches(6.2), Inches(11.333), Inches(0.5), f"{PROFILE['name']} — {PROFILE['contact']}", 11, C_MUTED, align=PP_ALIGN.CENTER)

    path = os.path.join(OUTPUT_DIR, "Portfolio_Salah_Barki_v4.pptx")
    prs.save(path)
    print(f"[OK] PPT v4 : {path}")


# ============================================================
# 3) PDF V4
# ============================================================
def create_pdf():
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY

    path = os.path.join(OUTPUT_DIR, "Portfolio_Salah_Barki_v4.pdf")
    doc = SimpleDocTemplate(path, pagesize=A4, rightMargin=1.8*cm, leftMargin=1.8*cm, topMargin=1.8*cm, bottomMargin=1.8*cm)
    styles = getSampleStyleSheet()

    C_NAVY = colors.HexColor('#0f172a')
    C_EMERALD = colors.HexColor('#10b981')
    C_TEXT = colors.HexColor('#374151')
    C_MUTED = colors.HexColor('#64748b')
    C_GREEN = colors.HexColor('#059669')
    C_CARD = colors.HexColor('#f1f5f9')
    C_BG = colors.HexColor('#ffffff')

    s_title = ParagraphStyle(name='s_title_v4', parent=styles['Title'], fontSize=26, textColor=C_NAVY, alignment=TA_CENTER, spaceAfter=4, fontName='Helvetica-Bold')
    s_sub = ParagraphStyle(name='s_sub_v4', fontSize=14, textColor=C_EMERALD, alignment=TA_CENTER, spaceAfter=10, fontName='Helvetica-Bold')
    s_name = ParagraphStyle(name='s_name_v4', fontSize=16, textColor=C_NAVY, alignment=TA_CENTER, spaceAfter=4, fontName='Helvetica-Bold')
    s_contact = ParagraphStyle(name='s_contact_v4', fontSize=9, textColor=C_MUTED, alignment=TA_CENTER, spaceAfter=14)
    s_h1 = ParagraphStyle(name='s_h1_v4', parent=styles['Heading1'], fontSize=16, textColor=C_NAVY, spaceAfter=8, spaceBefore=14, fontName='Helvetica-Bold')
    s_h2 = ParagraphStyle(name='s_h2_v4', parent=styles['Heading2'], fontSize=12, textColor=C_NAVY, spaceAfter=6, spaceBefore=10, fontName='Helvetica-Bold')
    s_body = ParagraphStyle(name='s_body_v4', parent=styles['Normal'], fontSize=9.5, leading=13, alignment=TA_JUSTIFY, spaceAfter=5, textColor=C_TEXT)
    s_bullet = ParagraphStyle(name='s_bullet_v4', parent=s_body, leftIndent=14, spaceAfter=3, bulletIndent=6, bulletFontName='Helvetica-Bold', bulletFontSize=9, bulletColor=C_EMERALD)
    s_result = ParagraphStyle(name='s_result_v4', parent=s_body, textColor=C_GREEN, fontSize=9.5, leading=13, spaceAfter=8, leftIndent=6, backColor=colors.HexColor('#f0fdf4'), borderColor=C_GREEN, borderWidth=2, borderPadding=6, leftPadding=8, rightPadding=8, topPadding=6, bottomPadding=6)
    s_diag = ParagraphStyle(name='s_diag_v4', parent=s_body, backColor=C_CARD, leftIndent=6, borderColor=colors.HexColor('#e2e8f0'), borderWidth=1, borderPadding=6, leftPadding=8, rightPadding=8, topPadding=6, bottomPadding=6)

    story = []
    story.append(Paragraph("PORTFOLIO", s_title))
    story.append(Paragraph("d'Interventions", s_sub))
    story.append(Paragraph("Maintenance Industrielle", s_sub))
    story.append(Paragraph(PROFILE["name"], s_name))
    story.append(Paragraph(PROFILE["contact"].replace('📍','').replace('📱','').replace('✉️',''), s_contact))
    story.append(Paragraph(PROFILE["summary"], ParagraphStyle(name='s_sum_v4', parent=s_body, backColor=C_CARD, borderColor=colors.HexColor('#e2e8f0'), borderWidth=1, borderPadding=8, leftPadding=10, rightPadding=10, topPadding=8, bottomPadding=8)))
    story.append(Spacer(1, 0.3*cm))

    data = [
        [Paragraph("<b>10+</b>", ParagraphStyle(name='sk1_v4', alignment=TA_CENTER, fontSize=18, textColor=C_EMERALD)),
         Paragraph("<b>7</b>", ParagraphStyle(name='sk2_v4', alignment=TA_CENTER, fontSize=18, textColor=C_EMERALD)),
         Paragraph("<b>4</b>", ParagraphStyle(name='sk3_v4', alignment=TA_CENTER, fontSize=18, textColor=C_EMERALD)),
         Paragraph("<b>6</b>", ParagraphStyle(name='sk4_v4', alignment=TA_CENTER, fontSize=18, textColor=C_EMERALD))],
        [Paragraph("Années d'exp.", ParagraphStyle(name='sk5_v4', alignment=TA_CENTER, fontSize=9, textColor=C_MUTED)),
         Paragraph("Entreprises", ParagraphStyle(name='sk6_v4', alignment=TA_CENTER, fontSize=9, textColor=C_MUTED)),
         Paragraph("Pays", ParagraphStyle(name='sk7_v4', alignment=TA_CENTER, fontSize=9, textColor=C_MUTED)),
         Paragraph("Fiches phares", ParagraphStyle(name='sk8_v4', alignment=TA_CENTER, fontSize=9, textColor=C_MUTED))],
    ]
    t = Table(data, colWidths=[4*cm, 4*cm, 4*cm, 4*cm])
    t.setStyle(TableStyle([
        ('ALIGN', (0,0), (-1,-1), 'CENTER'), ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('LEFTPADDING', (0,0), (-1,-1), 6), ('RIGHTPADDING', (0,0), (-1,-1), 6), ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
        ('BACKGROUND', (0,0), (-1,0), C_CARD),
    ]))
    story.append(t)
    story.append(Spacer(1, 0.3*cm))

    story.append(Paragraph("COMPÉTENCES & TECHNOLOGIES", s_h1))
    comps = [
        ("ÉLECTROTECHNIQUE & AUTOMATISME", "Siemens S7, TIA Portal, Step 7, programmation Ladder, câblage BT, lecture schémas électriques, variateurs, servo-moteurs, capteurs inductifs/photoélectriques/encodeurs."),
        ("MÉCANIQUE & FLUIDIQUE", "Presses injection, lignes PET/canettes, hydraulique (pompes, vérins, distributeurs, accumulateurs), pneumatique (distributeurs 5/2, vérins, préparation air), transmissions, guidages, roulements."),
        ("MÉTHODES & GESTION", "GMAO (SAP PM), planification préventive/corrective, analyse RCA (5 pourquoi, Ishikawa, Pareto), suivi KPI (MTBF, MTTR, OEE), consignation LOTO, normes HSE."),
    ]
    for cat, detail in comps:
        story.append(Paragraph(f"<b>{cat}</b> — {detail}", s_body))
    story.append(PageBreak())

    story.append(Paragraph("FICHES D'INTERVENTION", s_h1))
    story.append(Paragraph("Format : Contexte → Diagnostic structuré → Actions concrètes → Résultats chiffrés", s_body))
    story.append(Spacer(1, 0.2*cm))

    for f in FICHES:
        story.append(Paragraph(f"{f['num']}. {f['titre']}", s_h2))
        story.append(Paragraph(f"<i>{f['entreprise']} — {f['lieu']}  |  {f['type']}</i>", ParagraphStyle(name='s_meta_v4', parent=s_body, textColor=C_MUTED, fontSize=9)))
        story.append(Paragraph(f"<b>Équipement :</b> {f['equip']}", s_body))
        story.append(Paragraph(f"<b>Contexte :</b> {f['contexte']}", s_body))
        story.append(Paragraph(f"<b>Symptôme :</b> {f['symptome']}", s_body))
        story.append(Paragraph(f"<b>Diagnostic :</b><br/>{f['diagnostic']}", s_diag))
        story.append(Paragraph("<b>Actions réalisées :</b>", s_body))
        for line in f['actions']:
            story.append(Paragraph(line, s_bullet))
        story.append(Paragraph(f"<b>RÉSULTAT :</b> {f['resultat']}", s_result))
        story.append(Spacer(1, 0.2*cm))

    story.append(PageBreak())
    story.append(Paragraph("MÉTHODOLOGIE D'INTERVENTION", s_h1))
    meths = [
        ("1. SÉCURISATION", "Consignation LOTO, EPI, signalisation de zone."),
        ("2. DIAGNOSTIC", "Mesures, alarmes HMI, schémas, GMAO. Pas d'hypothèse sans vérification."),
        ("3. RCA", "5 pourquoi, Ishikawa, Pareto. Cause physique + systémique."),
        ("4. INTERVENTION", "Réparation, réglage, programmation. Respect procédures."),
        ("5. VALIDATION", "Essais charge, mesures A/A, GMAO, formation."),
    ]
    for t, d in meths:
        story.append(Paragraph(f"<b>{t}</b> — {d}", s_body))
    story.append(Spacer(1, 0.4*cm))
    story.append(Paragraph(f"<b>{PROFILE['name']}</b> — {PROFILE['title']}<br/>{PROFILE['contact']}", ParagraphStyle(name='s_footer_v4', alignment=TA_CENTER, fontSize=9, textColor=C_MUTED)))

    doc.build(story)
    print(f"[OK] PDF v4 : {path}")


# ============================================================
# 4) HTML V4 — Futuristic SaaS Dashboard Style
# ============================================================
def create_html():
    html_content = """<!DOCTYPE html>
<html lang="fr">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Portfolio — Salah Eddine Barki</title>
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&display=swap');
:root {
  --bg: #ffffff;
  --navy: #0f172a;
  --emerald: #10b981;
  --emerald-dark: #059669;
  --card: #f1f5f9;
  --card-hover: #e2e8f0;
  --text: #334155;
  --muted: #64748b;
  --border: #e2e8f0;
  --shadow: 0 4px 6px -1px rgba(0,0,0,0.05), 0 2px 4px -2px rgba(0,0,0,0.05);
  --shadow-lg: 0 10px 15px -3px rgba(0,0,0,0.08), 0 4px 6px -4px rgba(0,0,0,0.08);
}
* { margin:0; padding:0; box-sizing:border-box; }
body {
  font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
  background: var(--bg);
  color: var(--text);
  line-height: 1.6;
  font-size: 15px;
  -webkit-font-smoothing: antialiased;
}
.container { max-width: 1020px; margin: 0 auto; padding: 0 24px; }

/* HERO */
.hero {
  background: var(--navy);
  padding: 72px 0 48px;
  text-align: center;
  position: relative;
  overflow: hidden;
}
.hero::after {
  content: "";
  position: absolute;
  bottom: 0;
  left: 0;
  right: 0;
  height: 4px;
  background: linear-gradient(90deg, var(--emerald), #34d399, var(--emerald));
}
.hero h1 {
  font-size: 3.2rem;
  font-weight: 800;
  color: #fff;
  letter-spacing: -0.04em;
  line-height: 1.1;
}
.hero .subtitle {
  font-size: 1.15rem;
  color: var(--emerald);
  font-weight: 600;
  margin-top: 12px;
  letter-spacing: 0.02em;
}
.hero .name {
  font-size: 1.6rem;
  color: #94a3b8;
  font-weight: 500;
  margin-top: 20px;
}
.hero .contact {
  font-size: 0.9rem;
  color: #64748b;
  margin-top: 8px;
}
.hero-summary {
  max-width: 760px;
  margin: 28px auto 0;
  background: rgba(255,255,255,0.03);
  border: 1px solid rgba(255,255,255,0.08);
  border-radius: 16px;
  padding: 22px 26px;
  font-size: 0.95rem;
  color: #cbd5e1;
  text-align: justify;
  line-height: 1.65;
  backdrop-filter: blur(8px);
}

/* STATS */
.stats-bar {
  display: grid;
  grid-template-columns: repeat(auto-fit, minmax(160px, 1fr));
  gap: 16px;
  margin: -28px auto 40px;
  max-width: 900px;
  padding: 0 24px;
  position: relative;
  z-index: 2;
}
.stat-pill {
  background: var(--bg);
  border: 1px solid var(--border);
  border-radius: 14px;
  padding: 20px 14px;
  text-align: center;
  box-shadow: var(--shadow-lg);
  transition: transform 0.2s, box-shadow 0.2s;
}
.stat-pill:hover {
  transform: translateY(-3px);
  box-shadow: 0 20px 25px -5px rgba(0,0,0,0.1), 0 8px 10px -6px rgba(0,0,0,0.1);
}
.stat-pill .num {
  font-size: 2.2rem;
  font-weight: 800;
  color: var(--emerald);
  line-height: 1;
}
.stat-pill .lab {
  font-size: 0.78rem;
  color: var(--muted);
  text-transform: uppercase;
  letter-spacing: 0.08em;
  margin-top: 8px;
  font-weight: 600;
}

/* SECTIONS */
.section { padding: 48px 0; }
.section-title {
  font-size: 1.35rem;
  font-weight: 700;
  color: var(--navy);
  margin-bottom: 24px;
  display: flex;
  align-items: center;
  gap: 12px;
}
.section-title .line {
  flex: 1;
  height: 1px;
  background: linear-gradient(90deg, var(--border), transparent);
}

/* COMP CARDS */
.comp-grid {
  display: grid;
  grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
  gap: 18px;
}
.comp-card {
  background: var(--card);
  border-radius: 14px;
  padding: 22px 20px;
  border: 1px solid var(--border);
  transition: all 0.2s;
  position: relative;
}
.comp-card::before {
  content: "";
  position: absolute;
  top: 0;
  left: 0;
  right: 0;
  height: 3px;
  background: linear-gradient(90deg, var(--emerald), #34d399);
  border-radius: 14px 14px 0 0;
}
.comp-card:hover {
  background: var(--card-hover);
  transform: translateY(-2px);
  box-shadow: var(--shadow);
}
.comp-card h3 {
  font-size: 0.88rem;
  color: var(--navy);
  margin-bottom: 10px;
  text-transform: uppercase;
  letter-spacing: 0.04em;
  font-weight: 700;
}
.comp-card p {
  font-size: 0.9rem;
  color: var(--text);
  line-height: 1.55;
}

/* FICHE */
.fiche {
  margin-bottom: 36px;
  border-radius: 16px;
  overflow: hidden;
  box-shadow: var(--shadow-lg);
  border: 1px solid var(--border);
}
.fiche-header {
  background: var(--navy);
  color: #fff;
  padding: 18px 22px;
  display: flex;
  justify-content: space-between;
  align-items: center;
  flex-wrap: wrap;
  gap: 10px;
}
.fiche-header h3 {
  font-size: 1.05rem;
  font-weight: 600;
  letter-spacing: -0.01em;
}
.fiche-header .meta {
  font-size: 0.78rem;
  color: #94a3b8;
}
.fiche-header .badge {
  background: linear-gradient(135deg, var(--emerald), #34d399);
  color: #fff;
  padding: 5px 12px;
  border-radius: 20px;
  font-size: 0.7rem;
  font-weight: 700;
  text-transform: uppercase;
  letter-spacing: 0.04em;
}
.fiche-body {
  background: var(--bg);
  padding: 22px;
}
.fiche-grid {
  display: grid;
  grid-template-columns: 1fr 1fr;
  gap: 20px;
}
@media (max-width: 768px) { .fiche-grid { grid-template-columns: 1fr; } }

.fiche-col {
  background: var(--card);
  border-radius: 12px;
  padding: 16px;
  border: 1px solid var(--border);
}
.fiche-col h4 {
  font-size: 0.72rem;
  text-transform: uppercase;
  letter-spacing: 0.12em;
  color: var(--emerald-dark);
  margin-bottom: 12px;
  font-weight: 700;
}
.fiche-col p, .fiche-col li {
  font-size: 0.86rem;
  color: var(--text);
  margin-bottom: 7px;
  line-height: 1.55;
}
.fiche-col ul {
  list-style: none;
  padding-left: 0;
}
.fiche-col ul li {
  position: relative;
  padding-left: 16px;
}
.fiche-col ul li::before {
  content: "›";
  position: absolute;
  left: 0;
  color: var(--emerald);
  font-weight: 800;
  font-size: 1.1rem;
  line-height: 1;
}
.result-box {
  background: linear-gradient(135deg, #f0fdf4, #dcfce7);
  border-left: 4px solid var(--emerald-dark);
  padding: 14px 18px;
  border-radius: 0 12px 12px 0;
  margin-top: 14px;
  font-size: 0.9rem;
  color: var(--emerald-dark);
  font-weight: 700;
  line-height: 1.5;
}

/* METHOD */
.method-grid {
  display: grid;
  grid-template-columns: repeat(auto-fit, minmax(190px, 1fr));
  gap: 16px;
}
.method-card {
  background: var(--card);
  border-radius: 14px;
  padding: 20px 14px;
  text-align: center;
  border: 1px solid var(--border);
  transition: all 0.2s;
}
.method-card:hover {
  background: var(--card-hover);
  transform: translateY(-2px);
  box-shadow: var(--shadow);
}
.method-card .step {
  font-size: 2rem;
  font-weight: 800;
  color: var(--emerald);
  line-height: 1;
  margin-bottom: 10px;
}
.method-card h4 {
  font-size: 0.85rem;
  color: var(--navy);
  margin-bottom: 6px;
  font-weight: 700;
}
.method-card p {
  font-size: 0.82rem;
  color: var(--text);
  line-height: 1.45;
}

/* FOOTER */
footer {
  text-align: center;
  padding: 36px 0;
  margin-top: 24px;
  border-top: 1px solid var(--border);
  font-size: 0.85rem;
  color: var(--muted);
}

/* PRINT */
@media print {
  body { background: #fff; }
  .hero { padding: 30px 0; }
  .fiche { break-inside: avoid; }
}
</style>
</head>
<body>

<div class="hero">
  <div class="container">
    <h1>Portfolio</h1>
    <div class="subtitle">d'Interventions Maintenance Industrielle</div>
    <div class="name">Salah Eddine Barki</div>
    <div class="contact">📍 Amiens, France &nbsp;|&nbsp; 📱 +33 6 88 69 07 04 &nbsp;|&nbsp; ✉️ salahbarki.seb@gmail.com</div>
    <div class="hero-summary">""" + PROFILE["summary"] + """</div>
  </div>
</div>

<div class="stats-bar">
  <div class="stat-pill"><div class="num">10+</div><div class="lab">Années</div></div>
  <div class="stat-pill"><div class="num">7</div><div class="lab">Entreprises</div></div>
  <div class="stat-pill"><div class="num">4</div><div class="lab">Pays</div></div>
  <div class="stat-pill"><div class="num">6</div><div class="lab">Fiches phares</div></div>
</div>

<section class="section">
  <div class="container">
    <div class="section-title">Compétences & Technologies<span class="line"></span></div>
    <div class="comp-grid">
      <div class="comp-card"><h3>⚡ Électrotechnique & Automatisme</h3><p>Siemens S7, TIA Portal, Step 7, programmation Ladder, câblage BT, lecture schémas électriques, variateurs, servo-moteurs, capteurs inductifs/photoélectriques/encodeurs.</p></div>
      <div class="comp-card"><h3>🔧 Mécanique & Fluidique</h3><p>Presses injection, lignes PET/canettes, hydraulique (pompes, vérins, distributeurs, accumulateurs), pneumatique (distributeurs 5/2, vérins, préparation air), transmissions, guidages, roulements.</p></div>
      <div class="comp-card"><h3>📊 Méthodes & Gestion</h3><p>GMAO (SAP PM), planification préventive/corrective, analyse RCA (5 pourquoi, Ishikawa, Pareto), suivi KPI (MTBF, MTTR, OEE), consignation LOTO, normes HSE.</p></div>
    </div>
  </div>
</section>

<section class="section">
  <div class="container">
    <div class="section-title">Fiches d'Intervention<span class="line"></span></div>
    <p style="margin-bottom:24px; color:var(--muted); font-size:0.9rem;">Format structuré : Contexte → Diagnostic → Actions concrètes → Résultats chiffrés</p>
"""
    for f in FICHES:
        actions_html = "\n".join([f"<li>{l}</li>" for l in f['actions']])
        html_content += f"""
    <div class="fiche">
      <div class="fiche-header">
        <div><h3>{f['num']}. {f['titre']}</h3><div class="meta">{f['entreprise']} — {f['lieu']}</div></div>
        <span class="badge">{f['type']}</span>
      </div>
      <div class="fiche-body">
        <p style="font-size:0.84rem; color:var(--muted); margin-bottom:14px;"><b>Équipement :</b> {f['equip']}</p>
        <div class="fiche-grid">
          <div class="fiche-col">
            <h4>Diagnostic</h4>
            <p><b>Contexte :</b> {f['contexte']}</p>
            <p><b>Symptôme :</b> {f['symptome']}</p>
            <p style="white-space:pre-line;">{f['diagnostic']}</p>
          </div>
          <div class="fiche-col">
            <h4>Actions</h4>
            <ul>{actions_html}</ul>
          </div>
        </div>
        <div class="result-box">RÉSULTAT : {f['resultat']}</div>
      </div>
    </div>
"""
    html_content += """
  </div>
</section>

<section class="section">
  <div class="container">
    <div class="section-title">Méthodologie<span class="line"></span></div>
    <div class="method-grid">
      <div class="method-card"><div class="step">1</div><h4>SÉCURISATION</h4><p>Consignation LOTO, EPI, signalisation de zone.</p></div>
      <div class="method-card"><div class="step">2</div><h4>DIAGNOSTIC</h4><p>Mesures, alarmes HMI, schémas, GMAO.</p></div>
      <div class="method-card"><div class="step">3</div><h4>RCA</h4><p>5 pourquoi, Ishikawa, Pareto.</p></div>
      <div class="method-card"><div class="step">4</div><h4>INTERVENTION</h4><p>Réparation, réglage, programmation.</p></div>
      <div class="method-card"><div class="step">5</div><h4>VALIDATION</h4><p>Essais charge, mesures A/A, GMAO.</p></div>
    </div>
  </div>
</section>

<footer>
  <div class="container">
    <strong>Salah Eddine Barki</strong> — Technicien de Maintenance Industrielle<br>
    📱 +33 6 88 69 07 04 &nbsp;|&nbsp; ✉️ salahbarki.seb@gmail.com &nbsp;|&nbsp; 📍 Amiens, France
  </div>
</footer>

</body>
</html>
"""
    path = os.path.join(OUTPUT_DIR, "Portfolio_Salah_Barki_v4.html")
    with open(path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    print(f"[OK] HTML v4 : {path}")


if __name__ == "__main__":
    create_word()
    create_ppt()
    create_pdf()
    create_html()
    print("\n=== PORTFOLIO V4 COMPLET ===")
