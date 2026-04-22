# -*- coding: utf-8 -*-
"""
Add 6 images to PowerPoint v5 intervention slides
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

PPT_PATH = r"C:\Users\ba2rb\Downloads\salpre\Portfolio_Salah_Barki_v5.pptx"
OUTPUT_PATH = r"C:\Users\ba2rb\Downloads\salpre\Portfolio_Salah_Barki_v5_images.pptx"
IMG_DIR = r"C:\Users\ba2rb\Downloads\salpre"

# Fiche slides are slides 4-9 (0-indexed: 3,4,5,6,7,8)
# Mapping: slide_index -> image_file
IMAGE_MAP = {
    3: "img_01_encartonneuse.png",
    4: "img_02_presse_injection.png",
    5: "img_03_automate_s7.png",
    6: "img_04_banc_test.png",
    7: "img_05_hydraulique.png",
    8: "img_06_gmao.png",
}

prs = Presentation(PPT_PATH)

for slide_idx, img_name in IMAGE_MAP.items():
    slide = prs.slides[slide_idx]
    img_path = os.path.join(IMG_DIR, img_name)
    
    if not os.path.exists(img_path):
        print(f"[SKIP] {img_name} not found")
        continue
    
    # Add image in the top-right area of the slide (next to actions column)
    # Position: right column area, below the "ACTIONS" label
    # x=6.7", y=2.1", width=5.8", height=3.2" (cropped to fit)
    # Actually better: smaller image in the right column, below title
    pic = slide.shapes.add_picture(
        img_path,
        Inches(6.9),
        Inches(1.85),
        width=Inches(6.0)
    )
    
    # Crop height to max 3.0 inches (to not overlap with result bar)
    max_height = Inches(3.0)
    if pic.height > max_height:
        # Scale width proportionally
        ratio = max_height / pic.height
        pic.height = max_height
        pic.width = int(pic.width * ratio)
    
    print(f"[OK] Slide {slide_idx+1} — {img_name} ({pic.width.pt:.0f}x{pic.height.pt:.0f} pt)")

prs.save(OUTPUT_PATH)
print(f"\n[SAVED] {OUTPUT_PATH}")
