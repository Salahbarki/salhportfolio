# -*- coding: utf-8 -*-
"""
Génère Portfolio Salah Barki : Word, PowerPoint, PDF
"""
import os

OUTPUT_DIR = r"C:\Users\ba2rb\Downloads\salpre"

# ============================================================
# 1) WORD (.docx)
# ============================================================
def create_word():
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_PARAGRAPH_ALIGNMENT
    from docx.oxml.ns import qn

    doc = Document()

    # Style par défaut
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)
    style.element.rPr.rFonts.set(qn('w:eastAsia'), 'Calibri')

    # --- TITRE ---
    title = doc.add_heading("Portfolio d'Interventions Techniques", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = sub.add_run("Salah Eddine Barki — Technicien de Maintenance Industrielle")
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0x6F, 0x00)

    info = doc.add_paragraph()
    info.alignment = WD_ALIGN_PARAGRAPH.CENTER
    info.add_run("📍 Amiens, France  |  📱 +33 6 88 69 07 04  |  ✉️ salahbarki.seb@gmail.com\n\n")

    # --- PROFIL ---
    doc.add_heading("Profil & Expertise", level=1)
    doc.add_paragraph(
        "Plus de 10 ans d’expérience dans la maintenance industrielle multitechnique "
        "(mécanique, électrique, automatisme, hydraulique, pneumatique). Spécialisé dans le "
        "diagnostic avancé, la programmation d’automates Siemens et l’amélioration continue "
        "sur des lignes de production automatisées."
    )
    p = doc.add_paragraph()
    p.add_run("10+ années  |  7 entreprises  |  4 pays  |  Expertise multitechnique").bold = True

    # --- MACHINES ---
    doc.add_heading("Machines & Équipements Maîtrisés", level=1)
    machines = [
        ("Lignes Agroalimentaire", "Lignes remplissage PET & canettes, souffleuses, encolleuses, fardeleuses, convoyeurs"),
        ("Injection & Plasturgie", "Presses injection, sécheurs, trémies, détecteurs métal, tubes extrusion"),
        ("Automatisme & Électricité", "Automates Siemens S7 (TIA/Step7), variateurs, capteurs, armoires BT"),
        ("Mécanique & Fluidique", "Transmissions, guidages, vérins, pneumatique, hydraulique, compresseurs"),
        ("Automobile & Test", "Lignes assemblage câblage, bancs test électrique, inspection qualité"),
        ("GMAO & Supervision", "GMAO, plannings, collecte données, supervision process, MPSM"),
    ]
    for m, d in machines:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(f"{m} : ").bold = True
        p.add_run(d)

    # --- FICHES ---
    doc.add_page_break()
    doc.add_heading("Fiches d’Intervention Phares", level=1)

    fiches = [
        {
            "titre": "1. Panne d’encartonneuse sur ligne PET — Coca-Cola (Curatif urgent)",
            "equip": "Encartonneuse / Casseuse de colis sur ligne bouteilles PET 2L",
            "sympt": "Arrêt intempestif répétitif avec défaut 'produit coincé'. Rejets fréquents de colis mal formés.",
            "tech": "Automate Siemens, variateurs Servo, capteurs photoélectriques, pneumatique",
            "enjeu": "Ligne critique : 60 000 bouteilles/h. Risque rupture de stock client.",
            "diag": "Analyse programme Step 7 : capteur photo instable (condensation buée) + servo désynchronisé de 15ms sur convoyeur.",
            "actions": [
                "Remplacement capteur par modèle anti-buée IP67 + purge d’air",
                "Recalage mécanique guide d’entrée produit",
                "Optimisation programme : filtre logiciel 50ms sur signal capteur",
                "Resynchronisation axe servo via TIA Portal",
                "Validation run 30 min à cadence nominale"
            ],
            "result": "Temps d’arrêt réduit de 45 min → 12 min. Taux rejet ÷4. Mise à jour gamme maintenance préventive."
        },
        {
            "titre": "2. Mise en service presse d’injection — TE Connectivity (Installation)",
            "equip": "Presse injection 120T + sécheur + trémie + détecteur métal",
            "sympt": "Déchargement, installation et démarrage d’une nouvelle cellule injection connecteurs.",
            "tech": "Hydraulique, électrique BT, thermorégulation, automatisation cycles",
            "enjeu": "Respect délai mise en production. Qualité premier jet exigée.",
            "diag": "Ancrage mécanique, raccordement circuits, paramétrage sécheur (-40°C point rosée), intégration détecteur métal.",
            "actions": [
                "Ancrage presse sur plots anti-vibration",
                "Raccordement hydraulique (centraux), électrique (armoire + IHM)",
                "Mise en service sécheur et détecteur métal, réglage sensibilité",
                "Essais à vide puis production : ajustement paramètres cycle",
                "Rédaction PV recette et formation opérateurs"
            ],
            "result": "Mise en production en 48h. Zéro défaut qualité sur 500 premières pièces. Procédure standardisée partagée."
        },
        {
            "titre": "3. Débogage programme automate Siemens — TE Connectivity (Automatisme)",
            "equip": "Cellule injection automatisée Siemens S7-1200 / TIA Portal",
            "sympt": "Temps cycle anormalement long (+12s vs nominal). Arrêts aléatoires 'attente bras robot'.",
            "tech": "Siemens TIA Portal, Step 7, Ladder, PROFINET",
            "enjeu": "Perte cadence = sous-capacité. Impact OEE mensuel visible.",
            "diag": "Temporisation sécurité T37 fixée à 8s (bras met 5.5s). Bit validation capteur mal rafraîchi dans OB1. Programme séquentiel bloquant au lieu de parallèle.",
            "actions": [
                "Diagnostic E/S : vérification signaux capteurs/actionneurs (table visualisation)",
                "Réécriture FC12 pour paralléliser mouvements compatibles",
                "Réduction tempo sécurité à 6.5s avec validation encoder",
                "Ajout compteur cycle pour détection précoce dérive",
                "Sauvegarde, test pas-à-pas, suivi 1 semaine"
            ],
            "result": "Temps cycle ramené à nominal (-12s). OEE cellule +8%. Documentation mise à jour avec commentaires."
        },
        {
            "titre": "4. RCA banc de test électrique — Kromberg & Schubert (Qualité)",
            "equip": "Banc test électrique (continuité, isolation, résistance) ligne faisceaux",
            "sympt": "Taux rejet anormal 12% sur test continuité pour une même référence.",
            "tech": "Banc test paramétrable, base données tests, Excel",
            "enjeu": "Risque blocage ligne client automobile (just-in-time).",
            "diag": "Analyse données 3 semaines : dérive résistance contact pins (oxydation). Maintenance ne prévoyait pas nettoyage pins test. Extraction fumées soudure altérait contacts.",
            "actions": [
                "Collecte et analyse données historiques (tendance dérive)",
                "Démonstration oxydation par mesure pin à pin",
                "Remise à niveau banc : remplacement pins usés, réglage pression contact",
                "Amélioration extraction fumées à proximité",
                "Modification gamme préventive : contrôle/nettoyage pins toutes les 2 semaines"
            ],
            "result": "Taux rejet ramené à 1.2% (standard). Économie ~15 000€/an de rebut évitée. Procédure RCA partagée usine."
        },
        {
            "titre": "5. Dépannage hydraulique presse — Sovireso (Curatif mécanique)",
            "equip": "Presse découpe/pliage hydraulique",
            "sympt": "Perte pression fin de course descente. Bruit 'coups de bélier' circuit.",
            "tech": "Hydraulique 80 bar, distributeur 4/3, vérin double effet, schéma hydraulique",
            "enjeu": "Machine critique sans remplaçant sur site.",
            "diag": "Chute pression 80→35 bar côté tige fin course. Clapet anti-retour non refermé (corps étranger/joint). Accumulateur hydropneumatique déchargé (azote perdu).",
            "actions": [
                "Isolation circuit et vidange sécurisée",
                "Démontage clapet anti-retour, nettoyage, remplacement joint torique",
                "Vérification et recharge accumulateur à 65 bar (précharge azote)",
                "Remplacement filtre retour très encrassé",
                "Mise en pression progressive et essais sous charge"
            ],
            "result": "Intervention en 3h30 (prévision 1 jour). Pression stable rétablie. Disparition à-coups hydrauliques. Plan remplacement huile/filtres avancé de 2 mois."
        },
        {
            "titre": "6. Planification maintenance & GMAO — TE Connectivity (Planification)",
            "equip": "Site >120 équipements, lignes injection & assemblage",
            "sympt": "40% pannes liées à préventifs manqués ou mal planifiés. Pas de priorisation ABC.",
            "tech": "GMAO, supervision, Excel, plannings constructeurs",
            "enjeu": "Réduire arrêts imprévus de 25% sur l’année.",
            "diag": "Conflits avec production, pièces non disponibles, absence classification criticité.",
            "actions": [
                "Collecte analyse data GMAO + supervision (MTBF, MTTR, top 10 pannes)",
                "Élaboration planning annuel basé plans constructeurs + historique fiabilité",
                "Classification ABC équipements (criticité production/sécurité/coût)",
                "Coordination fenêtres maintenance avec production",
                "Approvisionnement anticipé pièces critiques + élaboration devis améliorations"
            ],
            "result": "Taux réalisation préventif 68% → 94%. Arrêts imprévus -28%. Coûts maintenance optimisés (-15% sous-traitants)."
        },
    ]

    for f in fiches:
        doc.add_heading(f["titre"], level=2)
        p = doc.add_paragraph()
        p.add_run("Équipement : ").bold = True
        p.add_run(f["equip"])
        p = doc.add_paragraph()
        p.add_run("Symptôme : ").bold = True
        p.add_run(f["sympt"])
        p = doc.add_paragraph()
        p.add_run("Technologies : ").bold = True
        p.add_run(f["tech"])
        p = doc.add_paragraph()
        p.add_run("Enjeu : ").bold = True
        p.add_run(f["enjeu"])
        p = doc.add_paragraph()
        p.add_run("Diagnostic : ").bold = True
        p.add_run(f["diag"])
        p = doc.add_paragraph()
        p.add_run("Actions réalisées : ").bold = True
        for a in f["actions"]:
            doc.add_paragraph(a, style='List Bullet')
        p = doc.add_paragraph()
        p.add_run("Résultat : ").bold = True
        run = p.add_run(f["result"])
        run.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)
        run.bold = True
        doc.add_paragraph()

    # --- COMPETENCES ---
    doc.add_page_break()
    doc.add_heading("Compétences Techniques", level=1)
    comp = [
        ("Électrotechnique & Automatisme", ["Siemens TIA Portal / Step 7", "Programmation Ladder / E/S", "Câblage électrique BT", "Lecture schémas électriques", "Variateurs & Servomoteurs"]),
        ("Mécanique & Fluidique", ["Maintenance mécanique industrielle", "Hydraulique", "Pneumatique", "Transmission / Guidage", "Presse injection / Lignes PET"]),
        ("Méthodes & Gestion", ["GMAO & Planification", "Analyse données / Reporting", "RCA (Root Cause Analysis)", "Amélioration continue", "Formation & transfert compétences"]),
    ]
    for cat, items in comp:
        doc.add_heading(cat, level=2)
        for i in items:
            doc.add_paragraph(i, style='List Bullet')

    # --- METHODOLOGIE ---
    doc.add_heading("Méthodologie d’Intervention", level=1)
    method = [
        ("1. Écoute & Sécurisation", "Recueil symptôme conducteur. Consignation LOTO. Protection personnelle/collective."),
        ("2. Diagnostic structuré", "Analyse alarmes, schémas/GMAO, mesures électriques/mécaniques, tests fonctionnels."),
        ("3. Intervention ciblée", "Réparation/remplacement. Optimisation si dérive récurrente. Respect standards."),
        ("4. Validation & Reporting", "Essais fonctionnels, validation qualité production. Saisie GMAO, mise à jour doc, transfert équipe."),
    ]
    for t, d in method:
        p = doc.add_paragraph()
        p.add_run(t + " : ").bold = True
        p.add_run(d)

    # Footer
    doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run("Salah Eddine Barki — Technicien de Maintenance Industrielle\n📱 +33 6 88 69 07 04  |  ✉️ salahbarki.seb@gmail.com").italic = True

    path = os.path.join(OUTPUT_DIR, "Portfolio_Maintenance_Salah_Barki.docx")
    doc.save(path)
    print(f"[OK] Word créé : {path}")


# ============================================================
# 2) POWERPOINT (.pptx)
# ============================================================
def create_ppt():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    PRIMARY = RGBColor(0x1A, 0x23, 0x7E)
    ACCENT = RGBColor(0xFF, 0x6F, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x26, 0x32, 0x38)
    GREEN = RGBColor(0x2E, 0x7D, 0x32)

    def add_title_slide(title, subtitle):
        blank = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(blank)
        # fond
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY
        # titre
        box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        # sous-titre
        box2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = ACCENT
        p2.alignment = PP_ALIGN.CENTER
        # contact
        box3 = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(0.8))
        tf3 = box3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = "📍 Amiens, France  |  📱 +33 6 88 69 07 04  |  ✉️ salahbarki.seb@gmail.com"
        p3.font.size = Pt(14)
        p3.font.color.rgb = WHITE
        p3.alignment = PP_ALIGN.CENTER
        return slide

    def add_content_slide(heading, lines, color=PRIMARY):
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)
        # barre haut
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.6))  # rectangle
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        # titre
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.05), Inches(12), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = heading
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        # contenu
        cbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12.333), Inches(6))
        ctf = cbox.text_frame
        ctf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                cp = ctf.paragraphs[0]
            else:
                cp = ctf.add_paragraph()
            cp.text = line
            cp.font.size = Pt(18)
            cp.font.color.rgb = DARK
            cp.space_after = Pt(10)
        return slide

    def add_fiche_slide(num, titre, equip, sympt, diag, actions, resultat):
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)
        # barre accent
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        # titre
        box = slide.shapes.add_textbox(Inches(0.4), Inches(0.05), Inches(12.5), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{num}. {titre}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE

        y = 0.7
        def block(label, text, y_offset, h=0.8, bold_label=True):
            b = slide.shapes.add_textbox(Inches(0.4), Inches(y_offset), Inches(12.5), Inches(h))
            t = b.text_frame
            t.word_wrap = True
            pr = t.paragraphs[0]
            if bold_label:
                pr.add_run().text = f"{label} "
                pr.runs[0].font.bold = True
                pr.runs[0].font.size = Pt(14)
                pr.runs[0].font.color.rgb = PRIMARY
                pr.add_run().text = text
                pr.runs[1].font.size = Pt(14)
                pr.runs[1].font.color.rgb = DARK
            else:
                pr.text = text
                pr.font.size = Pt(14)
                pr.font.color.rgb = DARK
            return y_offset + h + 0.05

        y = block("Équipement :", equip, y)
        y = block("Symptôme :", sympt, y)
        y = block("Diagnostic :", diag, y)
        # actions
        act_box = slide.shapes.add_textbox(Inches(0.4), Inches(y), Inches(12.5), Inches(1.4))
        act_tf = act_box.text_frame
        act_tf.word_wrap = True
        ap = act_tf.paragraphs[0]
        ap.add_run().text = "Actions : "
        ap.runs[0].font.bold = True
        ap.runs[0].font.size = Pt(14)
        ap.runs[0].font.color.rgb = PRIMARY
        for a in actions:
            ap = act_tf.add_paragraph()
            ap.text = f"• {a}"
            ap.font.size = Pt(13)
            ap.font.color.rgb = DARK
            ap.space_after = Pt(4)
        y += 1.5
        # résultat
        res_box = slide.shapes.add_textbox(Inches(0.4), Inches(y), Inches(12.5), Inches(0.7))
        res_tf = res_box.text_frame
        res_tf.word_wrap = True
        rp = res_tf.paragraphs[0]
        rp.add_run().text = "Résultat : "
        rp.runs[0].font.bold = True
        rp.runs[0].font.size = Pt(14)
        rp.runs[0].font.color.rgb = GREEN
        rp.add_run().text = resultat
        rp.runs[1].font.size = Pt(14)
        rp.runs[1].font.color.rgb = GREEN
        rp.runs[1].font.bold = True
        return slide

    # === SLIDES ===
    add_title_slide("Portfolio d’Interventions Techniques", "Salah Eddine Barki — Technicien de Maintenance Industrielle")

    add_content_slide("Profil & Chiffres Clés", [
        "• 10+ années d’expérience en maintenance industrielle multitechnique",
        "• 7 entreprises : agroalimentaire, plasturgie, automobile, connectique",
        "• Compétences couvrant : mécanique, électrique, automatisme, hydraulique, pneumatique",
        "• Spécialisations : diagnostic avancé, programmation Siemens S7, GMAO, amélioration continue",
        "",
        "📍 Amiens, France  |  📱 +33 6 88 69 07 04  |  ✉️ salahbarki.seb@gmail.com"
    ])

    add_content_slide("Machines & Équipements Maîtrisés", [
        "🥤 Lignes Agroalimentaire : remplissage PET & canettes, souffleuses, encolleuses, convoyeurs",
        "🔧 Injection & Plasturgie : presses injection, sécheurs, trémies, détecteurs métal",
        "⚡ Automatisme & Électricité : Siemens S7, TIA Portal, variateurs, capteurs, armoires BT",
        "🛠️ Mécanique & Fluidique : transmissions, vérins, hydraulique, pneumatique, compresseurs",
        "🚗 Automobile & Test : lignes assemblage câblage, bancs test électrique, contrôle qualité",
        "💻 GMAO & Supervision : planification, collecte données, supervision process, MPSM"
    ])

    fiches_ppt = [
        ("Panne encartonneuse ligne PET", "Coca-Cola", "Curatif urgent",
         "Encartonneuse/Casseuse colis ligne PET 2L",
         "Arrêts répétitifs 'produit coincé', rejets colis mal formés.",
         "Capteur photo instable (condensation) + servo désynchronisé 15ms.",
         ["Remplacement capteur anti-buée IP67 + purge d’air", "Recalage guide entrée", "Filtre logiciel 50ms programme", "Resynchronisation servo TIA Portal", "Validation run 30 min"],
         "Arrêt 45 min → 12 min. Rejet ÷4. Gamme préventive mise à jour."),
        ("Mise en service presse injection", "TE Connectivity", "Installation",
         "Presse injection 120T + sécheur + trémie + détecteur métal",
         "Installation et démarrage nouvelle cellule injection connecteurs.",
         "Ancrage, raccordements, paramétrage sécheur (-40°C rosée), détecteur métal.",
         ["Ancrage plots anti-vibration", "Raccordement hydraulique/électrique/IHM", "Mise en service sécheur + réglage détecteur", "Essais vide et production", "PV recette + formation"],
         "Mise en prod 48h. 0 défaut sur 500 pièces. Procédure standardisée."),
        ("Débogage programme Siemens S7", "TE Connectivity", "Automatisme",
         "Cellule injection Siemens S7-1200 / TIA Portal",
         "Cycle +12s, arrêts aléatoires 'attente bras robot'.",
         "Tempo T37 trop longue (8s vs 5.5s). Bit capteur mal rafraîchi. Programme bloquant.",
         ["Diagnostic E/S table visualisation", "Réécriture FC12 mouvements parallèles", "Réduction tempo 6.5s + validation encoder", "Ajout compteur cycle dérive", "Test pas-à-pas + suivi 1 semaine"],
         "Cycle nominal retrouvé (-12s). OEE +8%. Doc commentée mise à jour."),
        ("RCA banc test électrique", "Kromberg & Schubert", "Qualité / RCA",
         "Banc test électrique (continuité/isolation/résistance)",
         "Taux rejet 12% sur test continuité référence stable.",
         "Oxydation pins test + extraction fumées insuffisante.",
         ["Analyse tendance données 3 semaines", "Mesure comparative pin à pin", "Remplacement pins usés + réglage pression", "Amélioration extraction fumées", "Nouvelle gamme nettoyage pins 2 semaines"],
         "Rejet 12% → 1.2%. Économie ~15 000€/an. RCA partagée usine."),
        ("Dépannage hydraulique presse", "Sovireso", "Curatif mécanique",
         "Presse découpe/pliage hydraulique",
         "Perte pression fin descente. Coups de bélier circuit.",
         "Clapet anti-retour bloqué (joint). Accumulateur déchargé. Filtre encrassé.",
         ["Isolation + vidange sécurisée", "Démontage/nettoyage clapet + joint neuf", "Recharge accumulateur 65 bar azote", "Remplacement filtre retour", "Essais sous charge"],
         "3h30 au lieu de 1 jour. Pression stable. À-coups supprimés."),
        ("Planification GMAO annuelle", "TE Connectivity", "Planification & Fiabilité",
         "Site >120 équipements (injection + assemblage)",
         "40% pannes = préventifs manqués. Pas de classification criticité.",
         "Conflits production, pièces indisponibles, absence priorisation ABC.",
         ["Analyse data GMAO/supervision (MTBF/MTTR)", "Planning annuel constructeurs + historique", "Classification ABC criticité", "Coordination fenêtres maintenance", "Approvisionnement pièces + devis améliorations"],
         "Préventif 68% → 94%. Arrêts imprévus -28%. Coûts -15%."),
    ]

    for i, (titre, entreprise, badge, equip, sympt, diag, actions, result) in enumerate(fiches_ppt, 1):
        add_fiche_slide(i, f"{titre} — {entreprise} ({badge})", equip, sympt, diag, actions, result)

    add_content_slide("Compétences Clés", [
        "⚡ Électrotechnique & Automatisme : Siemens TIA Portal / Step 7, Ladder, câblage BT, schémas électriques, variateurs",
        "🔧 Mécanique & Fluidique : maintenance industrielle, hydraulique, pneumatique, transmissions, presses injection",
        "📊 Méthodes & Gestion : GMAO, planification, RCA, analyse données, amélioration continue, formation équipes",
        "",
        "Langues : Arabe (Expert), Français (Courant), Anglais (Avancé), Allemand (Débutant)"
    ], color=PRIMARY)

    add_content_slide("Méthodologie d’Intervention", [
        "1️⃣ Écoute & Sécurisation : recueil symptôme, consignation LOTO, protections",
        "2️⃣ Diagnostic structuré : alarmes, schémas, mesures, tests fonctionnels",
        "3️⃣ Intervention ciblée : réparation, optimisation, respect standards",
        "4️⃣ Validation & Reporting : essais, validation qualité, saisie GMAO, transfert équipe",
        "",
        "Merci de votre attention !",
        "📱 +33 6 88 69 07 04  |  ✉️ salahbarki.seb@gmail.com"
    ], color=PRIMARY)

    path = os.path.join(OUTPUT_DIR, "Portfolio_Maintenance_Salah_Barki.pptx")
    prs.save(path)
    print(f"[OK] PowerPoint créé : {path}")


# ============================================================
# 3) PDF avec ReportLab
# ============================================================
def create_pdf():
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY

    path = os.path.join(OUTPUT_DIR, "Portfolio_Maintenance_Salah_Barki.pdf")
    doc = SimpleDocTemplate(path, pagesize=A4,
                            rightMargin=2*cm, leftMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm)

    styles = getSampleStyleSheet()
    styleN = ParagraphStyle(name='NormalCustom', parent=styles['Normal'], fontSize=10, leading=14, alignment=TA_JUSTIFY, spaceAfter=6)
    styleH1 = ParagraphStyle(name='H1Custom', parent=styles['Heading1'], fontSize=18, textColor=colors.HexColor('#1a237e'), spaceAfter=12, spaceBefore=12)
    styleH2 = ParagraphStyle(name='H2Custom', parent=styles['Heading2'], fontSize=14, textColor=colors.HexColor('#1a237e'), spaceAfter=8, spaceBefore=10)
    styleH3 = ParagraphStyle(name='H3Custom', parent=styles['Heading3'], fontSize=12, textColor=colors.HexColor('#ff6f00'), spaceAfter=6, spaceBefore=8)
    styleTitle = ParagraphStyle(name='TitleCustom', parent=styles['Title'], fontSize=22, textColor=colors.HexColor('#1a237e'), alignment=TA_CENTER, spaceAfter=12)
    styleSub = ParagraphStyle(name='SubTitle', fontSize=14, textColor=colors.HexColor('#ff6f00'), alignment=TA_CENTER, spaceAfter=10)
    styleResult = ParagraphStyle(name='Result', parent=styleN, textColor=colors.HexColor('#2e7d32'), fontSize=10, spaceAfter=8)

    story = []

    # Titre
    story.append(Paragraph("Portfolio d’Interventions Techniques", styleTitle))
    story.append(Paragraph("Salah Eddine Barki — Technicien de Maintenance Industrielle", styleSub))
    story.append(Paragraph("<b>📍</b> Amiens, France  |  <b>📱</b> +33 6 88 69 07 04  |  <b>✉️</b> salahbarki.seb@gmail.com", ParagraphStyle(name='contact', alignment=TA_CENTER, fontSize=10, spaceAfter=20)))
    story.append(Spacer(1, 0.3*cm))

    # Profil
    story.append(Paragraph("Profil & Expertise", styleH1))
    story.append(Paragraph("""Plus de <b>10 ans</b> d’expérience dans la maintenance industrielle multitechnique 
    (mécanique, électrique, automatisme, hydraulique, pneumatique). Spécialisé dans le <b>diagnostic avancé</b>, 
    la <b>programmation d’automates Siemens</b> et l’<b>amélioration continue</b> sur des lignes de production automatisées.""", styleN))
    story.append(Spacer(1, 0.2*cm))

    # Table chiffres clés
    data = [
        [Paragraph("<b>10+</b>", ParagraphStyle(name='big', alignment=TA_CENTER, fontSize=16, textColor=colors.HexColor('#ff6f00'))),
         Paragraph("<b>7</b>", ParagraphStyle(name='big', alignment=TA_CENTER, fontSize=16, textColor=colors.HexColor('#ff6f00'))),
         Paragraph("<b>4</b>", ParagraphStyle(name='big', alignment=TA_CENTER, fontSize=16, textColor=colors.HexColor('#ff6f00'))),
         Paragraph("<b>∞</b>", ParagraphStyle(name='big', alignment=TA_CENTER, fontSize=16, textColor=colors.HexColor('#ff6f00')))],
        [Paragraph("Années d’expérience", ParagraphStyle(name='a', alignment=TA_CENTER, fontSize=9)),
         Paragraph("Entreprises", ParagraphStyle(name='b', alignment=TA_CENTER, fontSize=9)),
         Paragraph("Pays", ParagraphStyle(name='c', alignment=TA_CENTER, fontSize=9)),
         Paragraph("Pannes résolues", ParagraphStyle(name='d', alignment=TA_CENTER, fontSize=9))]
    ]
    t = Table(data, colWidths=[4*cm, 4*cm, 4*cm, 4*cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#f5f5f5')),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('LEFTPADDING', (0,0), (-1,-1), 6),
        ('RIGHTPADDING', (0,0), (-1,-1), 6),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ('GRID', (0,0), (-1,-1), 0.5, colors.lightgrey),
    ]))
    story.append(t)
    story.append(Spacer(1, 0.4*cm))

    # Machines
    story.append(Paragraph("Machines & Équipements Maîtrisés", styleH1))
    machines = [
        ("🥤 Lignes Agroalimentaire", "Remplissage PET & canettes, souffleuses, encolleuses, fardeleuses, convoyeurs"),
        ("🔧 Injection & Plasturgie", "Presses injection, sécheurs, trémies, détecteurs métal, tubes extrusion"),
        ("⚡ Automatisme & Électricité", "Siemens S7, TIA Portal, variateurs, capteurs, armoires BT"),
        ("🛠️ Mécanique & Fluidique", "Transmissions, guidages, vérins, pneumatique, hydraulique, compresseurs"),
        ("🚗 Automobile & Test", "Lignes assemblage câblage, bancs test électrique, inspection qualité"),
        ("💻 GMAO & Supervision", "GMAO, plannings, collecte données, supervision, MPSM"),
    ]
    for m, d in machines:
        story.append(Paragraph(f"<b>{m}</b> : {d}", styleN))
    story.append(PageBreak())

    # Fiches
    story.append(Paragraph("Fiches d’Intervention Phares", styleH1))
    story.append(Paragraph("Présentation selon la méthode Diagnostic → Action → Résultat.", styleN))
    story.append(Spacer(1, 0.3*cm))

    fiches = [
        {
            "titre": "1. Panne d’encartonneuse ligne PET — Coca-Cola (Curatif urgent)",
            "equip": "Encartonneuse/Casseuse de colis sur ligne PET 2L",
            "sympt": "Arrêt répétitif 'produit coincé'. Rejets colis mal formés.",
            "diag": "Capteur photo instable (condensation) + servo désynchronisé 15ms sur convoyeur.",
            "actions": [
                "Remplacement capteur anti-buée IP67 + purge d’air",
                "Recalage mécanique guide entrée produit",
                "Filtre logiciel 50ms sur signal capteur (programme)",
                "Resynchronisation axe servo via TIA Portal",
                "Validation run 30 min à cadence nominale"
            ],
            "result": "Arrêt réduit de 45 min → 12 min. Rejet ÷4. Gamme préventive mise à jour."
        },
        {
            "titre": "2. Mise en service presse injection — TE Connectivity (Installation)",
            "equip": "Presse injection 120T + sécheur + trémie + détecteur métal",
            "sympt": "Installation et démarrage nouvelle cellule injection connecteurs.",
            "diag": "Ancrage, raccordements, paramétrage sécheur (-40°C point rosée), détecteur métal.",
            "actions": [
                "Ancrage presse sur plots anti-vibration",
                "Raccordement hydraulique/électrique/IHM",
                "Mise en service sécheur + réglage détecteur métal",
                "Essais à vide puis en production",
                "PV recette et formation opérateurs"
            ],
            "result": "Mise en production 48h. 0 défaut sur 500 premières pièces. Procédure standardisée."
        },
        {
            "titre": "3. Débogage programme Siemens S7 — TE Connectivity (Automatisme)",
            "equip": "Cellule injection Siemens S7-1200 / TIA Portal",
            "sympt": "Cycle +12s vs nominal. Arrêts aléatoires 'attente bras robot'.",
            "diag": "Tempo T37 trop longue (8s vs 5.5s). Bit capteur mal rafraîchi dans OB1. Programme bloquant.",
            "actions": [
                "Diagnostic E/S table visualisation",
                "Réécriture FC12 pour paralléliser mouvements",
                "Réduction tempo 6.5s + validation encoder",
                "Ajout compteur cycle détection dérive",
                "Test pas-à-pas + suivi 1 semaine"
            ],
            "result": "Cycle nominal retrouvé (-12s). OEE cellule +8%. Documentation commentée."
        },
        {
            "titre": "4. RCA banc test électrique — Kromberg & Schubert (Qualité)",
            "equip": "Banc test électrique (continuité/isolation/résistance)",
            "sympt": "Taux rejet 12% sur test continuité référence stable.",
            "diag": "Oxydation pins test + extraction fumées insuffisante. Non prévu dans gamme.",
            "actions": [
                "Analyse tendance données 3 semaines",
                "Mesure comparative pin à pin",
                "Remplacement pins usés + réglage pression contact",
                "Amélioration extraction fumées",
                "Nouvelle gamme nettoyage pins toutes les 2 semaines"
            ],
            "result": "Rejet 12% → 1.2%. Économie ~15 000€/an. Procédure RCA partagée usine."
        },
        {
            "titre": "5. Dépannage hydraulique presse — Sovireso (Curatif)",
            "equip": "Presse découpe/pliage hydraulique",
            "sympt": "Perte pression fin descente. Coups de bélier circuit.",
            "diag": "Clapet anti-retour bloqué (corps étranger/joint). Accumulateur déchargé. Filtre encrassé.",
            "actions": [
                "Isolation circuit + vidange sécurisée",
                "Démontage/nettoyage clapet + joint torique neuf",
                "Recharge accumulateur 65 bar (précharge azote)",
                "Remplacement filtre retour encrassé",
                "Essais sous charge"
            ],
            "result": "3h30 au lieu de 1 jour. Pression stable. À-coups supprimés. Plan avancé 2 mois."
        },
        {
            "titre": "6. Planification GMAO annuelle — TE Connectivity (Planification)",
            "equip": "Site >120 équipements (injection + assemblage)",
            "sympt": "40% pannes = préventifs manqués. Pas de classification criticité.",
            "diag": "Conflits production, pièces indisponibles, absence priorisation ABC.",
            "actions": [
                "Analyse data GMAO/supervision (MTBF/MTTR/top 10)",
                "Planning annuel constructeurs + historique fiabilité",
                "Classification ABC criticité production/sécurité/coût",
                "Coordination fenêtres maintenance avec production",
                "Approvisionnement pièces critiques + devis améliorations"
            ],
            "result": "Préventif 68% → 94%. Arrêts imprévus -28%. Coûts sous-traitants -15%."
        },
    ]

    for f in fiches:
        story.append(Paragraph(f["titre"], styleH2))
        story.append(Paragraph(f"<b>Équipement :</b> {f['equip']}", styleN))
        story.append(Paragraph(f"<b>Symptôme :</b> {f['sympt']}", styleN))
        story.append(Paragraph(f"<b>Diagnostic :</b> {f['diag']}", styleN))
        story.append(Paragraph("<b>Actions :</b>", styleN))
        for a in f["actions"]:
            story.append(Paragraph(f"• {a}", ParagraphStyle(name='bullet', parent=styleN, leftIndent=12)))
        story.append(Paragraph(f"<b>Résultat :</b> {f['result']}", styleResult))
        story.append(Spacer(1, 0.2*cm))

    story.append(PageBreak())

    # Compétences
    story.append(Paragraph("Compétences Techniques", styleH1))
    comps = [
        ("⚡ Électrotechnique & Automatisme", "Siemens TIA Portal / Step 7, Programmation Ladder, câblage BT, schémas électriques, variateurs/servos"),
        ("🔧 Mécanique & Fluidique", "Maintenance industrielle, hydraulique, pneumatique, transmissions/guidages, presses injection/lignes PET"),
        ("📊 Méthodes & Gestion", "GMAO, planification, RCA, analyse données, amélioration continue, formation équipes"),
    ]
    for c, d in comps:
        story.append(Paragraph(f"<b>{c}</b>", styleH3))
        story.append(Paragraph(d, styleN))

    # Méthodologie
    story.append(Paragraph("Méthodologie d’Intervention", styleH1))
    method = [
        ("1️⃣ Écoute & Sécurisation", "Recueil symptôme, consignation LOTO, protections"),
        ("2️⃣ Diagnostic structuré", "Alarmes, schémas, mesures, tests fonctionnels"),
        ("3️⃣ Intervention ciblée", "Réparation, optimisation, respect standards"),
        ("4️⃣ Validation & Reporting", "Essais, validation qualité, saisie GMAO, transfert équipe"),
    ]
    for t, d in method:
        story.append(Paragraph(f"<b>{t}</b> — {d}", styleN))

    story.append(Spacer(1, 0.5*cm))
    story.append(Paragraph("<b>Salah Eddine Barki</b> — Technicien de Maintenance Industrielle<br/>📱 +33 6 88 69 07 04  |  ✉️ salahbarki.seb@gmail.com  |  📍 Amiens, France", ParagraphStyle(name='footer', alignment=TA_CENTER, fontSize=10, textColor=colors.grey)))

    doc.build(story)
    print(f"[OK] PDF créé : {path}")


if __name__ == "__main__":
    create_word()
    create_ppt()
    create_pdf()
    print("\n=== TOUS LES FICHIERS SONT PRÊTS ===")
